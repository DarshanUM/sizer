# coding: utf-8
""" This script generates PPT """
# import logging
import json
import datetime
import os
import xlsxwriter
import math
import copy
import urllib
import gettext

# Commented out since no longer drawing the UI doughnut chart
# import matplotlib
# Fix for no DISPLAY-Generating PNG with matplotlib when DISPLAY is undefined
# matplotlib.use('Agg')
# import matplotlib.pyplot as plt

from hyperconverged.models import Part, Thresholds
from base_sizer.solver.attrib import BaseConstants
from hyperconverged.solver.attrib import HyperConstants
from hyperconverged.models import SpecIntData
from base_sizer.solver.wl import WL
from hyperconverged.views.generate_layout import fetch_layout
from django.core.exceptions import ObjectDoesNotExist

from math import ceil
from collections import OrderedDict
from collections import defaultdict
from collections import Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_TICK_MARK
# from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.chart import XL_TICK_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
# from pptx.enum.dml import MSO_THEME_COLOR
# from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
# from pptx.enum.action import PP_ACTION

from sizer.local_settings import BASE_DIR
from . sizing_calc_report import SizingCalculatorReport

appname = "base"
dir = os.path.join(BASE_DIR, "sizer/locales")
# dir = "./locales"
gettext.install(appname, dir, codeset='utf8')

languages = {
    "english": gettext.translation(appname, dir, languages=["en"], codeset="utf8"),
    "japanese": gettext.translation(appname, dir, languages=["ja"], codeset="utf8")
}

_ = languages["english"].gettext

ENGLISH = "english"
JAPANESE = "japanese"

'''
import django
from django.conf import settings
os.environ.setdefault("DJANGO_SETTINGS_MODULE", "sizer.local_settings")
django.setup()
'''

VDI_KEY_HOLDER = ["User Type", "Provisioning Type", "Number Of Desktops", "Concurrency (%)", "Video RAM Size (MiB)",
                  "vCPUs per Desktop", "Clock per vCPU(MHZ)", "RAM per Desktop", "Average IOPS per Desktop",
                  "User / Application Data Size", "Gold Image per Desktop", "Number of Snapshots",
                  "Compression Savings(%)", "Dedupe Savings(%)", "Working Set(%)", "Home directories", "User Data IOPS"]
VDI_KEY_LIST = ["profile_type", "provisioning_type", "num_desktops", "concurrency", "video_RAM", "vcpus_per_desktop",
                "clock_per_desktop", "ram_per_desktop", "avg_iops_per_desktop", "disk_per_desktop", "gold_image_size",
                "snapshots", "compression_factor", "dedupe_factor", "working_set", "vdi_directory", "user_iops"]

VSI_KEY_HOLDER = ["Profile Type", "Number Of VMs", "vCPUs per VM", "RAM per VM", "Average IOPS per VM",
                  "User Data Capacity per VM", "Base Image Size", "Number of Snapshots", "Compression Savings (%)",
                  "Dedupe Savings (%)", "Working Set(%)", "vCPU Overprovisioning Ratio", "Cluster Type",
                  "Enable Remote Replication"]
VSI_KEY_LIST = ["profile_type", "num_vms", "vcpus_per_vm", "ram_per_vm", "avg_iops_per_vm", "disk_per_vm",
                "base_image_size", "snapshots", "compression_factor", "dedupe_factor", "working_set", "vcpus_per_core",
                "cluster_type", "remote_replication_enabled"]

RAW_KEY_HOLDER = ["CPU Unit", "CPU Model", "CPU Overprovisioning Ratio", "RAM Size", "RAM Overprovisioning Ratio",
                  "HDD Size", "SSD Size (GB)", "IOPS", "Compression Savings (%)", "Dedupe Savings(%)",
                  "Working Set(%)", "Future growth(%)", "Cluster Type"]
RAW_KEY_LIST = ["cpu_attribute", "cpu_model", "vcpus_per_core", "ram_size", "ram_opratio", "hdd_size", "ssd_size",
                "iops_value", "compression_factor", "dedupe_factor", "working_set",
                "overhead_percentage", "cluster_type"]

DB_KEY_HOLDER = ["Profile Type", "DB Type", "DB Size", "vCPUs per Core", "vCPU per DB", "No of DB instances",
                 "RAM per DB", "Average IOPS per DB", "Compression Savings (%)", "Dedupe Savings(%)",
                 "Working Set(%)", "Cluster Type", "Enable Remote Replication"]
DB_KEY_LIST = ["profile_type", "db_type", "db_size", "vcpus_per_core", "vcpus_per_db", "num_db_instances",
               "ram_per_db", "avg_iops_per_db", "compression_factor", "dedupe_factor", "working_set", "cluster_type",
               "remote_replication_enabled"]

ROBO_KEY_HOLDER = ["Profile Type", "Number Of VMs", "vCPUs per VM", "RAM per VM", "Average IOPS per VM",
                   "User Data Capacity per VM", "Base Image Size", "Number of Snapshots", "Compression Savings (%)",
                   "Dedupe Savings (%)", "Working Set(%)", "vCPU Overprovisioning Ratio"]
ROBO_KEY_LIST = ["profile_type", "num_vms", "vcpus_per_vm", "ram_per_vm", "avg_iops_per_vm", "disk_per_vm",
                 "base_image_size", "snapshots", "compression_factor", "dedupe_factor", "working_set", "vcpus_per_core"]

ORACLE_KEY_HOLDER = ["Profile Type", "DB Type", "DB Size", "vCPUs per Core", "vCPU per DB", "No of DB instances",
                     "RAM per DB", "Average IOPS per DB", "Compression Savings (%)", "Dedupe Savings(%)",
                     "Working Set(%)", "Cluster Type", "Enable Remote Replication"]
ORACLE_KEY_LIST = ["profile_type", "db_type", "db_size", "vcpus_per_core", "vcpus_per_db", "num_db_instances",
                   "ram_per_db", "avg_iops_per_db", "compression_factor", "dedupe_factor", "working_set",
                   "cluster_type", "remote_replication_enabled"]

EXCHANGE_KEY_HOLDER = ["vCPUs ", "vCPU Overprovisioning Ratio ", "Total RAM Size(GiB)", "Effective User Capacity(GB)",
                       "DB IOPS", "Log IOPS", "Maintenance IOPS", "Compression Savings (%)", "Dedupe Savings(%)",
                       "Working Set(%)", "Future growth(%)", "Cluster Type"]
EXCHANGE_KEY_LIST = ["vcpus", "vcpus_per_core", "ram_size", "hdd_size", "EXCHANGE_16KB", "EXCHANGE_32KB",
                     "EXCHANGE_64KB", "compression_factor", "dedupe_factor", "working_set", "overhead_percentage",
                     "cluster_type"]

VDIINFRA_KEY_HOLDER = ["Broker Type", "CPU Overprovisioning Ratio", "RAM Overprovisioning Ratio",
                       "Compression Savings (%)", "Dedupe Savings(%)", "Cluster Type"]
VDIINFRA_KEY_LIST = ["infra_type", "vcpus_per_core", "ram_opratio", "compression_factor", "dedupe_factor",
                     "cluster_type"]

VMDETAILS_KEY_HOLDER = ['-vCPUs', '-RAM', '-Storage (GB)', '-Count']
VMDETAILS_KEY_LIST = ["vcpus_per_vm", "ram_per_vm", "disk_per_vm", "num_vms"]

EPIC_KEY_HOLDER = ["Total Users", "CPU SKU", "Users per Host", "RAM per Guest", "VM Guests per Host",
                   "Expected Number of Hosts", "Compression Savings (%)", "Dedupe Savings(%)"]
EPIC_KEY_LIST = ["total_users", "cpu", "users_per_host", "ram_per_guest", "guests_per_host", "expected_hosts",
                 "compression_factor", "dedupe_factor"]

VEEAM_KEY_HOLDER = ["HDD Size", "Performance Headroom", "Compression Savings (%)", "Dedupe Savings(%)"]
VEEAM_KEY_LIST = ["hdd_size", "fault_tolerance", "compression_factor", "dedupe_factor"]

SPLUNK_KEY_HOLDER = ["Profile Type", "Daily data ingest", "Max Volume per Indexer", "Storage Accumulation-",
                     "VM Details-",
                     "Compression Savings (%)", "Dedupe Savings(%)", "Cluster Type"]
SPLUNK_KEY_LIST = ["profile_type", "daily_data_ingest", "max_vol_ind", "storage_acc", "vm_details",
                   "compression_factor", "dedupe_factor",
                   "cluster_type"]

RDSH_KEY_HOLDER = ["User Type", "Broker Type", "Total Users", "GPU Users", "Video RAM Size (MiB)",
                   "RDSH directories", "User IOPS", "User Capacity", "vCPUs", "Users", "Clock per Session",
                   "Max vCPU Overprovisioning Ratio", "RAM", "OS Image Size",
                   "Compression Savings(%)", "Dedupe Savings(%)", ]
RDSH_KEY_LIST = ["profile_type", "broker_type", "total_users", "gpu_users", "video_RAM", "rdsh_directory",
                 "user_iops", "hdd_per_user", "vcpus_per_vm", "sessions_per_vm", "clock_per_session",
                 "max_vcpus_per_core", "ram_per_vm", "os_per_vm", "compression_factor", "dedupe_factor"]

CONTAINER_KEY_HOLDER = ["Container Type", "Number Of Containers", "vCPUs per Container", "vCPU Overprovisioning Ratio",
                        "RAM per Container", "Average Storage IOPS", "User / Application Data Size",
                        "OS Image Size", "Working Set(%)", "Compression Savings (%)", "Dedupe Savings (%)"]
CONTAINER_KEY_LIST = ["container_type", "num_containers", "vcpus_per_container", "vcpus_per_core",
                      "ram_per_container", "iops_per_container", "disk_per_container", "base_image_size",
                      "working_set", "compression_factor", "dedupe_factor"]

AIML_KEY_HOLDER = ["Input Source", "Expected Utilization", "Number Of Data Scientist", "vCPUs per Data Scientist",
                   "vCPU Overprovisioning Ratio", "RAM per Data Scientist(GiB)", "GPUs per Data Scientist",
                   "Storage on HX cluster", "Capacity per Data Scientist", "Compression Savings (%)",
                   "Dedupe Savings (%)"]
AIML_KEY_LIST = ["input_type", "expected_util", "num_data_scientists", "vcpus_per_ds", "vcpus_per_core",
                 "ram_per_ds", "gpu_per_ds", "enablestorage", "disk_per_ds", "compression_factor", "dedupe_factor"]

AWR_KEY_HOLDER = ["Input Type", "DB Type", "Size for", "DB Size", "vCPUs per Core", "Total vCPUs",
                  "RAM per DB (GiB)", "Compression Savings (%)", "Dedupe Savings(%)", "Database Overhead (%)",
                  "Cluster Type", "Enable Remote Replication"]
AWR_KEY_LIST = ["input_type", "db_type", "provisioned", "db_size", "vcpus_per_core", "vcpus_per_db",
                "ram_per_db", "compression_factor", "dedupe_factor", "metadata_size",
                "cluster_type", "remote_replication_enabled"]

ANTHOS_POD_KEY_HOLDER = ["Pod Name", "Pod CPU", "Pod Memory", "Pod Storage", "Pod Qty", "Worker Node CPU",
                         "Worker Node RAM", "Worker Node Storage", "Prometheous", "User Cluster HA"]
ANTHOS_POD_KEY_LIST = ["pod_name", "pod_cpu", "pod_ram", "pod_storage", "pod_quantity", "worker_node_cpu",
                     "worker_node_ram", "worker_node_storage", "prometheous_on", "pod_ha"]

ANTHOS_KEY_HOLDER = ["Load balancer CPU", "Load balancer RAM", "Controller VM CPU", "Controller VM RAM", "Compute CPU",
                     "Compute RAM", "User Cluster VM Storage", "Audit Logs /User Cluster", "ETCD Object Data/User Cluster(GB)",
                     "Prometheous/Grafana Add-on Data/User cluster", "GC Ops Suite overhead Addition(GB)",
                     "Anthos Master VM CPU", "Anthos Master VM RAM", "Admin Cluster VM Storage",
                     "GC Ops Suite for Admin cluster(GB)", "ETCD data for Admin cluster", "Dedupe Savings(%)"]

ANTHOS_KEY_LIST = ["load_balancer_cpu", "load_balancer_ram", "controller_vm_cpu", "controller_vm_ram", "user_vm_cpu",
                   "user_vm_ram", "user_vm_storage", "audit_log", "etcd_event", "prometheous_storage",
                   "gc_ops_overhead_user_vm", "vm_cpu", "vm_ram", "vm_storage", "gc_ops_overhead", "etcd_anthos_master",
                   "dedupe_factor"]

CLUSTER_KEY_HOLDER = ["Replication Factor", "Fault Tolerance Nodes"]
CLUSTER_KEY_LIST = ["replication_factor", "fault_tolerance"]

IOPS_DESC = defaultdict(str)
IOPS_DESC['VDI'] = "Workload Model: 20K IO Block Size, 100% Write"
IOPS_DESC['VDI_INFRA'] = "Workload Model: 20K IO Block Size, 100% Write"
IOPS_DESC['VSI'] = "Workload Model: 8K IO Block Size, 70/30 Read/Write"
IOPS_DESC['DB'] = "Workload Model: 8K IO Block Size, 70/30 Read/Write"
IOPS_DESC['ORACLE'] = "Workload Model: 8K IO Block Size, 70/30 Read/Write"
IOPS_DESC['ROBO'] = "Workload Model: 8K IO Block Size, 70/30 Read/Write"
IOPS_DESC['EXCHANGE'] = "Workload Model: Mixture of 32K IO Block Size 60/40 Read/Write, " \
                        "\n16K IO Block Size 100% Write and 8K IO Block Size 70/30 Read/Write"

IOPS_DESC['EPIC'] = "Workload Model: 20K IO Block Size, 100% Write"
IOPS_DESC['VEEAM'] = "Performance was not modelled for this workload"
IOPS_DESC['SPLUNK'] = "Workload Model: 8K IO Block Size, 70/30 Read/Write"
IOPS_DESC['RDSH'] = "Workload Model: 20K IO Block Size, 100% Write"
IOPS_DESC['CONTAINER'] = "Workload Model: 8K IO Block Size, 70/30 Read/Write"
IOPS_DESC['AIML'] = "Performance was not modelled for this workload"
IOPS_DESC['RAW'] = "Workload Model: 8K IO Block Size, 70/30 Read/Write"
IOPS_DESC['RAW_FILE'] = "Workload Model: 8K IO Block Size, 70/30 Read/Write"
IOPS_DESC['AWR_FILE'] = "Workload Model: 8K IO Block Size, 70/30 Read/Write"
IOPS_DESC['OLAP'] = "Workload Model: 64K IO Block Size, 100% Writes"
IOPS_DESC['ANTHOS'] = "Performance was not modelled for this workload"
IOPS_DESC['ROBO_BACKUP'] = "Workload Model: 8K IO Block Size, 70/30 Read/Write"

total_consumed_slots = 0

therm_cpulist = ["6242R [Cascade]", "6246R [Cascade]", "6248R [Cascade]", "6258R [Cascade]"]
cputhermnote = False

CORES = 1
CLOCK = 0

FULL_CAPACITY = 3

class Report(object):

    def __init__(self, s_data):
        self.scenario_data = json.loads(s_data)


class PPTReport(Report):
    SLIDEMASTER_TITLE_LAYOUT_INDEX = 0
    SLIDEMASTER_SECTION_LAYOUT_INDEX = 1
    SLIDEMASTER_BULLET_LAYOUT_INDEX = 2
    SLIDEMASTER_BOX_LAYOUT_INDEX = 5
    SLIDEMASTER_BLANK_LAYOUT_INDEX = 6
    SLIDEMASTER_BIG_STMT_LAYOUT_INDEX = 8
    SLIDEMASTER_BACKUP_LAYOUT_INDEX = 10
    SLIDEMASTER_LAST_LAYOUT_INDEX = 11
    SLIDEMASTER_VDI_CALC_LAYOUT_INDEX = 17
    SLIDEMASTER_CALC_LAYOUT_INDEX = 18
    SLIDEMASTER_NODECALC_LAYOUT_INDEX1 = 19
    SLIDEMASTER_NODECALC_LAYOUT_INDEX2 = 20
    SLIDEMASTER_RAW_CALC_LAYOUT_INDEX = 21

    SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_1 = 12
    SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_2 = 13
    SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_3 = 14
    SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_4 = 15
    SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_5 = 16
    SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_6 = 17

    SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_1 = 24
    SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_2 = 25
    SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_3 = 26
    SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_4 = 27
    SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_5 = 28
    SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_6 = 29

    SLIDEMASTER_JA_GLOSSARY_INDEX = 30
    SLIDEMASTER_LAST_LAYOUT_INDEX_JA = 31

    SLIDEMASTER_ANTHOS_LAYOUT_INDEX = 32

    def __init__(self, scenario_data, req_data):

        Report.__init__(self, scenario_data)
        # current_dir = os.path.dirname(os.path.realpath(__file__))

        template_path = os.path.join(BASE_DIR, "sizer/hflexsizer_bom_template.pptx")
        self.prs = Presentation(template_path)
        self.report_name = None
        self.workload_note = False

        # Create an utilization image name to be unique
        scenario_id = self.scenario_data['id']
        date_time = datetime.datetime.now().strftime('%Y-%m-%d_%H-%M')
        image_name = str(scenario_id) + "_" + str(date_time) + ".png"

        image_name = os.path.join(BASE_DIR, image_name)

        self.util_image_name = image_name

        self.thresholds = dict()

        self.wl_capsum = dict()

        # Load data from database
        self.load_thresholds()

        self.fi_count = req_data.get('fipackage_count', 0)
        self.fi_option = req_data.get('fipackage_name', "HX-FI-48P")

        self.language = req_data.get('language', ENGLISH)

        global _
        if self.language == JAPANESE:
            _ = languages[JAPANESE].gettext
            languages[JAPANESE].install()
        else:
            _ = languages[ENGLISH].gettext
            languages[ENGLISH].install()

        if 'slides' not in req_data:
            req_data['slides'] = ["Agenda", "Sizing Report", "Lowest_Cost Sizing Report", "All-Flash Sizing Report",
                                  "All NVMe Sizing Report", "Workload Summary"]

        self.slides = req_data['slides']

        self.fixed_cluster_name = None

    def load_thresholds(self):

        threshold_table = Thresholds.objects.all()

        for workload_type in HyperConstants.WORKLOAD_TYPES:
            self.thresholds[workload_type] = dict()
            for threshold_factor in range(0, 4):
                self.thresholds[workload_type][threshold_factor] = dict()

        for threshold in threshold_table:
            self.thresholds[threshold.workload_type][threshold.threshold_category][threshold.threshold_key] = \
                threshold.threshold_value

    def get_threshold_value(self, workload_type, threshold_factor, threshold_key):

        try:
            workload_type = HyperConstants.ROBO if (workload_type == HyperConstants.ROBO_BACKUP_SECONDARY) else workload_type
            workload_type = HyperConstants.VSI if (workload_type == HyperConstants.ROBO_BACKUP) else workload_type
            return 100 - self.thresholds[workload_type][threshold_factor][threshold_key]
        except KeyError:
            # raise HXException("Missing_Threshold_Value" + self.logger_header)
            pass

    def threshold_value(self, workload_type, threshold_factor, threshold_key):

        try:
            return self.thresholds[workload_type][threshold_factor][threshold_key] / 100.0
        except KeyError:
            # raise HXException("Missing_Threshold_Value" + self.logger_header)
            pass

    def get_header_data(self):

        hdata = dict()

        hdata['title'] = _("Cisco HyperFlex Sizing Report")
        hdata['subhead'] = _("Scenario") + ": " + self.scenario_data['name']

        if 'account' in self.scenario_data['settings_json'][0] and self.scenario_data['settings_json'][0]['account']:
            hdata['customer'] = _("Customer") + ": " + self.scenario_data['settings_json'][0]['account']
        else:
            hdata['customer'] = " "

        if 'username' in self.scenario_data['workload_json']:
            hdata['generated'] = _("Generated by") + ": " + self.scenario_data['workload_json']['username']
        else:
            hdata['generated'] = _("Generated by") + ": N/A"

        current_date = datetime.datetime.now().strftime('%Y-%m-%d')
        hdata['created_date'] = _("Date") + ": " + current_date

        return hdata

    def create_header_page(self):

        prs = self.prs
        header_data = self.get_header_data()

        slide1 = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_TITLE_LAYOUT_INDEX])
        title1 = slide1.shapes.title
        subhead = slide1.placeholders[13]
        customertext = slide1.placeholders[1]
        genbytext = slide1.placeholders[10]
        datetext = slide1.placeholders[11]

        title1.text = header_data["title"]
        subhead.text = header_data["subhead"]
        customertext.text = header_data["customer"]
        genbytext.text = header_data["generated"]
        datetext.text = header_data["created_date"]

    def create_disclaimer_page(self):
        # SLIDE 2 - Disclaimer
        prs = self.prs

        disclaimer_path = os.path.join(BASE_DIR, "sizer/disclaimer.txt")

        if self.language == JAPANESE:
            disclaimer_path = os.path.join(BASE_DIR, "sizer/disclaimer_ja.txt")

        with open(disclaimer_path, 'r', encoding='utf-8') as string_file:
            disclaimer_data = string_file.read()

        slide2 = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
        shapes2 = slide2.shapes
        shapes2.title.text = _("Disclaimer")
        tbleft = Inches(0.26)
        tbtop = Inches(0.5)
        tbwidth = Inches(9.4)
        tbheight = Inches(4)
        tx_box = shapes2.add_textbox(tbleft, tbtop, tbwidth, tbheight)
        tx_box.text = disclaimer_data
        tx_box.text_frame.paragraphs[0].font.size = Pt(14)
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        tx_box.text_frame.word_wrap = True
        tx_box.text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        tx_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    def create_agenda_page(self):

        # SLIDE 3 - Agenda

        prs = self.prs
        slide2 = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BULLET_LAYOUT_INDEX])
        shapes2 = slide2.shapes

        title_shape = shapes2.title
        body_shape = shapes2.placeholders[10]
        title_shape.text = _("Content")

        tf = body_shape.text_frame
        p = tf.add_paragraph()
        p.text = _("HyperFlex Product Overview")
        p.level = 1

        if "Sizing Report" in self.slides:
            p = tf.add_paragraph()
            p.text = _("User Inputs")
            p.level = 1

        if len(self.scenario_data['workload_result']) > 2 and "Lowest_Cost/All-Flash/All NVMe Comparison" in self.slides:
            p = tf.add_paragraph()
            p.text = _("Comparison of Lowest Cost Option v/s All Flash Option v/s All NVMe Option")
            p.level = 1

        if "Lowest_Cost Sizing Report" in self.slides or \
                "All-Flash Sizing Report" in self.slides or \
                "Fixed Config Sizing Report" in self.slides or \
                "All NVMe Sizing Report" in self.slides:
            p = tf.add_paragraph()
            p.text = _("Overall Sizing Results")
            p.level = 1

        if "Lowest_Cost Sizing Report" in self.slides:
            p = tf.add_paragraph()
            p.text = 'Lowest_Cost Sizing Report'
            p.level = 2

        if "All-Flash Sizing Report" in self.slides:
            p = tf.add_paragraph()
            p.text = 'All-Flash Sizing Report'
            p.level = 2

        if "All NVMe Sizing Report" in self.slides:
            p = tf.add_paragraph()
            p.text = 'All NVMe Sizing Report'
            p.level = 2

        if "Fixed Config Sizing Report" in self.slides:
            p = tf.add_paragraph()
            p.text = 'Fixed Config Sizing Report'
            p.level = 2

        if "Workload Calculation" in self.slides:
            p = tf.add_paragraph()
            p.text = _('Workload Calculation')
            p.level = 1

        if "Node Calculation" in self.slides:
            p = tf.add_paragraph()
            p.text = _('Node Calculation')
            p.level = 1

        if "Workload Summary" in self.slides:
            p = tf.add_paragraph()
            p.text = _("Workload Summary")
            p.level = 1

        if "Glossary" in self.slides:
            p = tf.add_paragraph()
            p.text = _("Glossary")
            p.level = 1

        # p.level = 1
        # p = tf.add_paragraph()
        # p.text = 'Comparison of HyperFlex Only and HyperFlex & Compute'
        # p = tf.add_paragraph()
        # p.text = 'Comparison of HyperFlex(Hybrid) and HyperFlex(All Flash)'

    def create_overview_page(self):

        prs = self.prs

        if self.language == JAPANESE:
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_1])
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_2])
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_3])
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_4])
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_5])
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_JA_INDEX_6])
        else:
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_1])
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_2])

            slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_3])
            note_filepath = os.path.join(BASE_DIR, "sizer/slide_note3.txt")

            with open(note_filepath, 'r', encoding='utf-8') as string_file:
                notes = string_file.read()

            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes

            slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_4])
            note_filepath = os.path.join(BASE_DIR, "sizer/slide_note4.txt")

            with open(note_filepath, 'r', encoding='utf-8') as string_file:
                notes = string_file.read()

            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes

            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_5])

            slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_OVERVIEW_LAYOUT_INDEX_6])
            note_filepath = os.path.join(BASE_DIR, "sizer/slide_note6.txt")

            with open(note_filepath, 'r', encoding='utf-8') as string_file:
                notes = string_file.read()

            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes

    @staticmethod
    def get_wl_summary(wl_list):

        wl_summary_table = OrderedDict()
        row_index = 0
        rvtools_flag = False

        # HEADER FORMATION
        wl_summary_table['header'] = [_('Workload'), _('Type'), _('Description'), _('Count'),
                                      'RF', _('Fault Tolerance')]

        for wdata in wl_list:
            wl_summary_per_row = list()

            wl_count = ''

            if wdata['wl_type'] == "VDI":

                wl_name = wdata['wl_name']
                wl_description = wdata['profile_type'] + ", " + wdata['provisioning_type']
                wl_count = "{:,}".format(wdata['num_desktops'])

            elif wdata['wl_type'] == "VSI":

                wl_name = wdata['wl_name']
                dr_input = \
                    ", " + str(wdata['replication_amt']) + "% Replicated" if wdata['remote_replication_enabled'] else ""
                wl_description = wdata['profile_type'] + dr_input
                wl_count = "{:,}".format(wdata['num_vms'])

            elif wdata['wl_type'] == "RAW":

                wl_name = wdata['wl_name']
                wl_description = ""
                cpu_unit = wdata['cpu_attribute']
                cunit = 'Cores' if cpu_unit == 'vcpus' else 'Clock (GHz),'

                if wdata['cpu_attribute'] == 'vcpus':
                    cores_desc = \
                        ", vCPUs -" + "{:,}".format(wdata[cpu_unit]) + ", CPU Overprovisioning Ratio-" + \
                        str(wdata['vcpus_per_core'])
                else:
                    cores_desc = \
                        ", Clock (GHz) -" + "{:,}".format(wdata[cpu_unit]) + ", CPU Overprovisioning Ratio-" + \
                        str(wdata['vcpus_per_core'])

                inputcpu_model = wdata.get('cpu_model', "Intel Platinum 8164")
                ram_opratio = wdata.get('ram_opratio', 1)

                wl_desc = \
                    "CPU Unit- " + cunit + ", CPU Model-" + inputcpu_model + cores_desc + ", RAM Size (" + \
                    wdata['ram_size_unit'] + ")-" + "{:,}".format(wdata['ram_size']) + ", RAM Overprovisioning Ratio-" \
                    + str(ram_opratio) + ", HDD Size (" + wdata['hdd_size_unit'] + ")- " + \
                    "{:,}".format(wdata['hdd_size'])

                if 'overhead_percentage' in wdata:
                    wl_desc += ", Future growth (%)- " + str(wdata['overhead_percentage'])

                if 'input_type' in wdata:
                    wl_description = "Input- " + wdata['input_type'] + ", "

                wl_description += wl_desc

            elif wdata['wl_type'] == "RAW_FILE":

                wl_name = wdata['wl_name']
                wl_description = ""
                profiler_warning = ""
                cpu_unit = wdata['cpu_attribute']
                cunit = 'Cores' if cpu_unit == 'vcpus' else 'Clock (GHz),'

                if wdata['cpu_attribute'] == 'vcpus':
                    cores_desc = \
                        ", vCPUs -" + "{:,}".format(wdata[cpu_unit]) + ", CPU Overprovisioning Ratio-" + \
                        str(wdata['vcpus_per_core'])
                else:
                    cores_desc = \
                        ", Clock (GHz) -" + "{:,}".format(wdata[cpu_unit]) + ", CPU Overprovisioning Ratio-" + \
                        str(wdata['vcpus_per_core'])

                inputcpu_model = wdata.get('cpu_model', "Intel Platinum 8164")
                ram_opratio = wdata.get('ram_opratio', 1)

                wl_desc = \
                    "CPU Unit- " + cunit + ", CPU Model-" + inputcpu_model + cores_desc + ", RAM Size (" + \
                    wdata['ram_size_unit'] + ")-" + "{:,}".format(wdata['ram_size']) + ", RAM Overprovisioning Ratio-" \
                    + str(ram_opratio) + ", HDD Size (" + wdata['hdd_size_unit'] + ")- " + \
                    "{:,}".format(wdata['hdd_size'])

                if 'overhead_percentage' in wdata:
                    wl_desc += ", Future growth (%)- " + str(wdata['overhead_percentage'])

                if 'input_type' in wdata:

                    if wdata['input_type'] == "RV Tools XLSX":
                        rvtools_flag = True

                    wl_description = "Input- " + wdata['input_type'] + ", "

                    if 'hasWarning' in wdata and wdata['hasWarning']:
                        profiler_warning = " \n(Warning: Insufficient historical data for some of the servers for " \
                                           "effective sizing)"

                wl_description += wl_desc + profiler_warning

            elif wdata['wl_type'] == "DB":

                wl_name = wdata['wl_name']
                dr_input = \
                    ", " + str(wdata['replication_amt']) + "% Replicated" if wdata['remote_replication_enabled'] else ""
                wl_description = wdata['profile_type'] + ", " + wdata['db_type'] + dr_input
                wl_count = "{:,}".format(wdata['num_db_instances'])

            elif wdata['wl_type'] == "ROBO":

                wl_name = wdata['wl_name']
                wl_description = wdata['profile_type']
                wl_count = "{:,}".format(wdata['num_vms'])

            elif wdata['wl_type'] == "ROBO_BACKUP":
    
                wl_name = wdata['wl_name']
                wl_description = wdata['profile_type']
                wl_count = "{:,}".format(wdata['num_vms'])
                
            elif wdata['wl_type'] == "ORACLE":

                wl_name = wdata['wl_name']
                dr_input = \
                    ", " + str(wdata['replication_amt']) + "% Replicated" if wdata['remote_replication_enabled'] else ""
                wl_description = wdata['profile_type'] + ", " + wdata['db_type'] + dr_input
                wl_count = "{:,}".format(wdata['num_db_instances'])

            elif wdata['wl_type'] == "EXCHANGE":

                wl_name = wdata['wl_name']
                wl_description = \
                    "vCPUs- " + str(wdata['vcpus']) + ", CPU Overprovisioning Ratio-" + str(wdata['vcpus_per_core']) + \
                    ", RAM Size (GiB)- " + "{:,}".format(wdata['ram_size']) + ", HDD Size (GB)- " + \
                    "{:,}".format(wdata['hdd_size'])

                if 'overhead_percentage' in wdata:
                    wl_description += ", Future growth (%)- " + str(wdata['overhead_percentage'])

            elif wdata['wl_type'] == "VDI_INFRA":

                wl_name = wdata['wl_name']
                wl_description = "Broker Type-" + wdata['infra_type']

            elif wdata['wl_type'] == "EPIC":

                wl_name = wdata['wl_name']
                dc1 = wdata['datacentres'][0]
                dc2 = wdata['datacentres'][1]
                wl_description = "Datacenter 1: Total Users Supported (%)-" + str(dc1['concurrent_user_pcnt']) + \
                                 ", Number of Clusters-" + str(
                    dc1['num_clusters']) + ". Datacenter 2: Total Users Supported (%)-" + \
                                 str(dc2['concurrent_user_pcnt']) + ", Number of Clusters-" + str(dc2['num_clusters'])

            elif wdata['wl_type'] == "VEEAM":
                wl_name = wdata['wl_name']
                wl_description = "Total Storage Capacity Requirement-" + str(wdata['hdd_size']) + wdata['hdd_size_unit']

            elif wdata['wl_type'] == "SPLUNK":

                wl_name = wdata['wl_name']
                wl_description = "Profile Type-" + wdata['profile_type'] + ", Daily data ingest-" + \
                                 str(wdata['daily_data_ingest']) + wdata['daily_data_ingest_unit'] + \
                                 ", Max Volume per Indexer-" + str(wdata['max_vol_ind']) + wdata['max_vol_ind_unit']

            elif wdata['wl_type'] == "RDSH":

                wl_name = wdata['wl_name']
                wl_description = "Profile Type-" + wdata['profile_type'] + ", Broker Type-" + wdata['broker_type']
                wl_count = "{:,}".format(wdata['total_users'])

            elif wdata['wl_type'] == "CONTAINER":

                wl_name = wdata['wl_name']
                wl_description = "Container Type-" + wdata['container_type']
                wl_count = "{:,}".format(wdata['num_containers'])

            elif wdata['wl_type'] == "AIML":

                wl_name = wdata['wl_name']
                wl_description = "Input source-" + wdata['input_type'] + \
                                 ", Expected Utilization-" + wdata['expected_util']
                wl_count = "{:,}".format(wdata['num_data_scientists'])

            elif wdata['wl_type'] == "AWR_FILE":

                wl_name = wdata['wl_name']
                wl_description = "AWR File input, " + "Database Type-" + wdata['db_type']

            elif wdata['wl_type'] == "ANTHOS":

                wl_name = wdata['wl_name']
                wl_description = "Number of Pods configured: " + str(len(wdata['pod_detail']))

            else:
                continue

            wl_summary_per_row.append(wl_name)

            if wdata['wl_type'] == "ROBO":
                wl_summary_per_row.append("Edge")
            else:
                wl_summary_per_row.append(wdata['wl_type'])

            wl_summary_per_row.append(wl_description)
            wl_summary_per_row.append(wl_count)
            wl_summary_per_row.append(str(wdata['replication_factor']))
            wl_fault_tolerance = "N+" + str(wdata['fault_tolerance'])
            wl_summary_per_row.append(wl_fault_tolerance)

            row_name = "row" + str(row_index)
            wl_summary_table[row_name] = wl_summary_per_row
            row_index += 1

        return wl_summary_table, rvtools_flag

    def get_sizing_configuration(self):

        sizing_config_table = OrderedDict()

        sizing_config_table['header'] = [_('Configuration'), _('Description')]

        hx_version = self.scenario_data['settings_json'][0]['hx_version']
        sizing_config_per_row = list()
        sizing_config_per_row.append(_('HXDP Version'))
        sizing_config_per_row.append(str(hx_version))

        row_index = 0
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        sizer_version = self.scenario_data['settings_json'][0]['sizer_version']
        sizing_config_per_row = list()
        sizing_config_per_row.append(_('Sizer Version'))
        sizing_config_per_row.append(str(sizer_version))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        settings_data = self.scenario_data['settings_json'][0]

        heterogeneous = settings_data['heterogenous']
        if not heterogeneous:
            hetero_value = _('Hyperflex only')
        else:
            hetero_value = _('Hyperflex and Compute')

        sizing_config_per_row = list()
        cluster_conf = _("Cluster Type")
        sizing_config_per_row.append(cluster_conf)
        sizing_config_per_row.append(hetero_value)

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        threshold = settings_data['threshold']
        if threshold == 0:
            thr_value = 'Conservative'
        elif threshold == 1:
            thr_value = 'Standard'
        elif threshold == 2:
            thr_value = 'Aggressive'
        else:
            thr_value = 'Full Capacity'

        sizing_config_per_row = list()
        threshold_conf = _("Threshold")
        sizing_config_per_row.append(threshold_conf)
        sizing_config_per_row.append(_(thr_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        hypervisor = settings_data['hypervisor']
        if hypervisor == 'esxi':
            hypervisor_value = 'ESXi'
        else:
            hypervisor_value = 'Hyper-V'

        sizing_config_per_row = list()
        hypervisor_conf = _("Hypervisor")
        sizing_config_per_row.append(hypervisor_conf)
        sizing_config_per_row.append(hypervisor_value)

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        sizing_option = settings_data['bundle_only']
        if sizing_option == HyperConstants.BUNDLE_ONLY:
            sizing_value = 'Bundle only'
        elif sizing_option == HyperConstants.CTO_ONLY:
            sizing_value = 'CTO only'
        else:
            sizing_value = 'Bundle and CTO'

        sizing_config_per_row = list()
        sizing_conf = _("Sizing Option")
        sizing_config_per_row.append(sizing_conf)
        sizing_config_per_row.append(_(sizing_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        # server_option = settings_data['server_type']
        # if server_option == "ALL":
        #     server_value = 'All Servers'
        # else:
        #     server_value = server_option + " Server"

        # sizing_config_per_row = list()
        # server_conf = "Server Type"
        # sizing_config_per_row.append(server_conf)
        # sizing_config_per_row.append(server_value)
        #
        # row_index += 1
        # row_name = "row" + str(row_index)
        # sizing_config_table[row_name] = sizing_config_per_row

        disk_option = settings_data['disk_option']
        if disk_option == "SED":
            disk_value = 'Storage Encrypted Disks only'
        elif disk_option == "NVME":
            disk_value = 'NVMe only'
        elif disk_option == "COLDSTREAM":
            disk_value = 'Optane only'
        elif disk_option == "LFF":
            disk_value = 'Large Form Factor Disks'
        else:
            disk_value = 'All Disks'

        sizing_config_per_row = list()
        disk_conf = _("Disk Option")
        sizing_config_per_row.append(disk_conf)
        sizing_config_per_row.append(_(disk_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        disk_option = settings_data.get('cache_option', 'ALL')
        if disk_option == "SED":
            disk_value = 'Storage Encrypted Disks only'
        elif disk_option == "NVMe":
            disk_value = 'NVMe only'
        elif disk_option == "Optane":
            disk_value = 'Optane only'
        else:
            disk_value = 'All Cache Disks'

        sizing_config_per_row = list()
        disk_conf = _("Cache Option")
        sizing_config_per_row.append(disk_conf)
        sizing_config_per_row.append(_(disk_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        # modular_lan = settings_data['modular_lan']
        # if modular_lan == "40G_VIC":
        #     modular_lan_value = '40G VIC'
        # else:
        #     modular_lan_value = 'All'
        #
        # sizing_config_per_row = list()
        # modular_lan_conf = "Modular LAN"
        # sizing_config_per_row.append(modular_lan_conf)
        # sizing_config_per_row.append(modular_lan_value)
        #
        # row_index += 1
        # row_name = "row" + str(row_index)
        # sizing_config_table[row_name] = sizing_config_per_row

        sw_cost = settings_data['license_yrs']
        sw_cost_value = str(sw_cost) + _(" Year(s)")

        sizing_config_per_row = list()
        sw_cost_conf = _("Software Cost")
        sizing_config_per_row.append(sw_cost_conf)
        sizing_config_per_row.append(sw_cost_value)

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        sizing_config_per_row = list()
        sizing_config_per_row.append(_("Hardware Acceleration"))
        hercules_conf = settings_data.get('hercules_conf', 'enabled')
        if hercules_conf == "enabled":
            hercules_conf_value = 'Auto'
        elif hercules_conf == "forced":
            hercules_conf_value = 'On'
        else:
            hercules_conf_value = 'Off'
        sizing_config_per_row.append(_(hercules_conf_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        sizing_config_per_row = list()
        sizing_config_per_row.append(_("HyperFlex Boost"))
        hx_boost_conf = settings_data.get('hx_boost_conf', 'enabled')
        if hx_boost_conf == "enabled":
            hx_boost_conf_value = 'Auto'
        elif hx_boost_conf == "forced":
            hx_boost_conf_value = 'On'
        else:
            hx_boost_conf_value = 'Off'
        sizing_config_per_row.append(_(hx_boost_conf_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        sizing_config_per_row = list()
        sizing_config_per_row.append(_("Single Cluster"))
        single_cluster = settings_data.get('single_cluster', False)
        single_cluster_value = "Yes" if single_cluster == True else "No"
        sizing_config_per_row.append(_(single_cluster_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        sizing_config_per_row = list()
        sizing_config_per_row.append(_("CPU Generation"))
        cpu_generation = settings_data.get('cpu_generation', "recommended")
        if cpu_generation == "recommended":
            cpu_generation_value = 'Recommended'
        elif cpu_generation == "sky":
            cpu_generation_value = 'Skylake'
        elif cpu_generation == "cascade":
            cpu_generation_value = "Cascade Lake"
        else:
            cpu_generation_value = 'All'
        sizing_config_per_row.append(_(cpu_generation_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        return sizing_config_table

    def get_sizing_configuration_fixed_sizing(self):

        sizing_config_table = OrderedDict()

        sizing_config_table['header'] = [_('Configuration'), _('Description')]

        hx_version = self.scenario_data['settings_json'][0]['hx_version']
        sizing_config_per_row = list()
        sizing_config_per_row.append(_('HXDP Version'))
        sizing_config_per_row.append(str(hx_version))

        row_index = 0
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        sizer_version = self.scenario_data['settings_json'][0]['sizer_version']
        sizing_config_per_row = list()
        sizing_config_per_row.append(_('Sizer Version'))
        sizing_config_per_row.append(str(sizer_version))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        settings_data = self.scenario_data['settings_json'][0]

        threshold = settings_data['threshold']
        if threshold == 0:
            thr_value = 'Conservative'
        elif threshold == 1:
            thr_value = 'Standard'
        elif threshold == 2:
            thr_value = 'Aggressive'
        else:
            thr_value = 'Full Capacity'

        sizing_config_per_row = list()
        threshold_conf = "Threshold"
        sizing_config_per_row.append(_(threshold_conf))
        sizing_config_per_row.append(_(thr_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        hypervisor = settings_data['hypervisor']
        if hypervisor == 'esxi':
            hypervisor_value = 'ESXi'
        else:
            hypervisor_value = 'Hyper-V'

        sizing_config_per_row = list()
        hypervisor_conf = "Hypervisor"
        sizing_config_per_row.append(_(hypervisor_conf))
        sizing_config_per_row.append(hypervisor_value)

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        hercules_conf = settings_data['hercules_conf']
        if hercules_conf == 'disabled':
            hercules_conf_value = 'Off'
        else:
            hercules_conf_value = 'On'

        sizing_config_per_row = list()
        sizing_config_per_row.append(_("Hardware Acceleration"))
        sizing_config_per_row.append(_(hercules_conf_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        hx_boost_conf = settings_data['hx_boost_conf']
        if hx_boost_conf == 'disabled':
            hx_boost_conf_value = 'Off'
        else:
            hx_boost_conf_value = 'On'

        sizing_config_per_row = list()
        sizing_config_per_row.append(_("HyperFlex Boost"))
        sizing_config_per_row.append(_(hx_boost_conf_value))

        row_index += 1
        row_name = "row" + str(row_index)
        sizing_config_table[row_name] = sizing_config_per_row

        return sizing_config_table

    def get_node_configuration_fixed_sizing(self, fixed_config_json):

        node_config_table = OrderedDict()

        # node_properties = self.scenario_data['settings_json'][0]['node_properties']
        node_properties = fixed_config_json['node_properties']
        node_type = node_properties['nodeType']
        nodetype = 'CTO' if node_type == 'cto' else 'Bundle'
        node_config_per_row = list()
        node_config_per_row.append(_('Node Type'))
        node_config_per_row.append(nodetype)

        row_index = 0
        row_name = "row" + str(row_index)
        node_config_table[row_name] = node_config_per_row

        node_name = node_properties['node']
        node_config_per_row = list()
        node_config_per_row.append(_('HyperFlex Node Type'))
        node_config_per_row.append(node_name)

        row_index += 1
        row_name = "row" + str(row_index)
        node_config_table[row_name] = node_config_per_row

        no_of_nodes = node_properties['no_of_nodes']
        node_config_per_row = list()
        node_config_per_row.append(_('No. of HyperFlex Nodes'))
        node_config_per_row.append(str(no_of_nodes))

        row_index += 1
        row_name = "row" + str(row_index)
        node_config_table[row_name] = node_config_per_row

        compute_node = node_properties['compute_node']
        node_config_per_row = list()
        node_config_per_row.append(_('Compute Node Type'))
        node_config_per_row.append(compute_node)

        row_index += 1
        row_name = "row" + str(row_index)
        node_config_table[row_name] = node_config_per_row

        no_of_computes = node_properties['no_of_computes']
        node_config_per_row = list()
        node_config_per_row.append(_('No. of Compute Nodes'))
        node_config_per_row.append(str(no_of_computes))

        row_index += 1
        row_name = "row" + str(row_index)
        node_config_table[row_name] = node_config_per_row

        cpu = "%s (%s, %s)" %(node_properties['cpu'][0].split(' ')[0], node_properties['cpu'][1], node_properties['cpu'][2])
        node_config_per_row = list()
        node_config_per_row.append(_('CPU Type (Cores, GHz)'))
        node_config_per_row.append(cpu)

        row_index += 1
        row_name = "row" + str(row_index)
        node_config_table[row_name] = node_config_per_row

        ram_total = sum(node_properties['ram'][1] * ram_size for ram_size in node_properties['ram'][2])
        ram_str = ["%s * %sGiB" %(node_properties['ram'][1], ram_size) for ram_size in node_properties['ram'][2]]
        ram_calc_str = ' + '.join([str(v) for v in ram_str])
        ram = "%s [%s]" %(ram_total, ram_calc_str)

        node_config_per_row = list()
        node_config_per_row.append(_('RAM per Node (GiB)'))
        node_config_per_row.append(str(ram))

        row_index += 1
        row_name = "row" + str(row_index)
        node_config_table[row_name] = node_config_per_row

        disks_per_node = node_properties['disks_per_node']
        node_config_per_row = list()
        node_config_per_row.append(_('No. of Disk Drives per Node'))
        node_config_per_row.append(str(disks_per_node))

        row_index += 1
        row_name = "row" + str(row_index)
        node_config_table[row_name] = node_config_per_row

        disk_capacity = node_properties['disk_capacity'][0]
        node_config_per_row = list()
        node_config_per_row.append(_('Disk Drive Size (GB)'))
        node_config_per_row.append(str(disk_capacity))

        row_index += 1
        row_name = "row" + str(row_index)
        node_config_table[row_name] = node_config_per_row

        cache_size = node_properties['cache_size'][0]
        node_config_per_row = list()
        node_config_per_row.append(_('Cache Size'))
        node_config_per_row.append(cache_size)

        row_index += 1
        row_name = "row" + str(row_index)
        node_config_table[row_name] = node_config_per_row

        return node_config_table

    def get_aggwl_requirement(self):

        wl_require_table = OrderedDict()

        wl_require_table['header'] = ['CPU', 'RAM', 'Total Workload Disk Capacity', 'IOPS']
        wl_require_per_row = list()

        util_list = self.scenario_data['workload_result'][0]['summary_info']['Utilization']

        for util_data in util_list:

            tag_name = util_data['tag_name']

            if (tag_name == 'Cache') or (tag_name == 'GPU Users'):
                continue

            node_val = round(util_data['workload_val'], 2)
            units = util_data['units']

            if tag_name == 'Physical Cores' or tag_name == 'CPU':
                units = "Cores"

            row_value = str(node_val) + " " + units
            wl_require_per_row.append(row_value)

        row_index = 0
        row_name = "row" + str(row_index)
        wl_require_table[row_name] = wl_require_per_row

        return wl_require_table

    @staticmethod
    def render_single_table(slide_obj, x, y, wid, heig, tabletitle, col_width, tabledata, first_row_merge = False):

        lentable = len(tabledata) - 1

        slide = slide_obj
        shapes = slide.shapes

        rows = len(tabledata)
        cols = len(tabledata['header'])
        left = Inches(x)
        top = Inches(y)
        width = Inches(wid)
        height = Inches(heig)

        tbtop = top - Inches(0.64)
        tx_box = slide.shapes.add_textbox(left, tbtop, Inches(2), Inches(1))
        tf = tx_box.text_frame
        p = tf.add_paragraph()
        p.text = tabletitle
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x92, 0xC8, 0xE9)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        for colindex in range(0, cols):
            table.columns[colindex].width = Inches(col_width[colindex])

        for colindex in range(0, cols):
            table.cell(0, colindex).text = tabledata['header'][colindex]
            table.cell(0, colindex).text_frame.paragraphs[0].font.size = Pt(9)
            table.cell(0, colindex).text_frame.paragraphs[0].font.bold = True
            table.cell(0, colindex).fill.solid()
            table.cell(0, colindex).fill.fore_color.rgb = RGBColor(0xB2, 0xB2, 0xB2)
            table.cell(0, colindex).vertical_anchor = MSO_ANCHOR.MIDDLE

        for rowindex in range(0, lentable):
            for colindex in range(0, cols):
                key = 'row' + str(rowindex)
                table.cell(rowindex + 1, colindex).text = tabledata[key][colindex]
                table.cell(rowindex + 1, colindex).text_frame.paragraphs[0].font.size = Pt(8)
                table.cell(rowindex + 1, colindex).vertical_anchor = MSO_ANCHOR.MIDDLE

            if first_row_merge and rowindex == 1:
                PPTReport.merge_cells_horizontally(table=table, row_idx=1, start_col_idx=1, end_col_idx=3)
                table.cell(rowindex, 1).text_frame.paragraphs[0].font.size = Pt(8)

    @staticmethod
    def render_util_table(slide_obj, x, y, wid, heig, col_width, tabledata):

        lentable = len(tabledata) - 1

        slide = slide_obj
        shapes = slide.shapes

        rows = len(tabledata)
        cols = len(tabledata['header'])
        left = Inches(x)
        top = Inches(y)
        width = Inches(wid)
        height = Inches(heig)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        for colindex in range(0, cols):
            table.columns[colindex].width = Inches(col_width[colindex])

        for colindex in range(0, cols):
            table.cell(0, colindex).text = tabledata['header'][colindex]
            table.cell(0, colindex).text_frame.paragraphs[0].font.size = Pt(9)
            table.cell(0, colindex).text_frame.paragraphs[0].font.bold = True
            table.cell(0, colindex).fill.solid()
            table.cell(0, colindex).fill.fore_color.rgb = RGBColor(0xB2, 0xB2, 0xB2)
            table.cell(0, colindex).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for rowindex in range(0, lentable):
            for colindex in range(0, cols):
                key = 'row' + str(rowindex)
                table.cell(rowindex + 1, colindex).text = tabledata[key][colindex]
                table.cell(rowindex + 1, colindex).text_frame.paragraphs[0].font.size = Pt(8)
                table.cell(rowindex + 1, colindex).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                if colindex == 0:
                    table.cell(rowindex + 1, colindex).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                    table.cell(rowindex + 1, colindex).text_frame.paragraphs[0].font.bold = True

    def create_user_in_page(self):
        prs = self.prs
        slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
        shapes = slide.shapes
        shapes.title.text = _('User Inputs')

        wl_summary_table, rvtools_flag = self.get_wl_summary(self.scenario_data['workload_json']['wl_list'])

        x = 0.29  # LEFT POSITION
        y = 0.73  # TOP POSITION
        wid = 6.9  # table Width
        heig = 2.7  # table height
        col_width = [1.3, 0.9, 2.9, 0.6, 0.4, 0.8]

        if len(wl_summary_table) <= 2:
            heig = 0.7
        elif 2 < len(wl_summary_table) <= 4:
            heig = 1.2
        elif 4 < len(wl_summary_table) <= 6:
            heig = 2

        wl_summary_truncate_table = OrderedDict()
        if len(wl_summary_table) <= 6:
            wl_summary_truncate_table = wl_summary_table
            table_title = _("Workload Summary")
        else:
            wl_summary_truncate_table['header'] = wl_summary_table['header']
            table_title = _("Workload Summary (Top 6 workloads*)")
            for n in range(0, 6):
                key = 'row' + str(n)
                wl_summary_truncate_table[key] = wl_summary_table[key]

        self.render_single_table(slide, x, y, wid, heig, table_title, col_width, wl_summary_truncate_table)

        if len(wl_summary_table) > 6:
            tx_box = slide.shapes.add_textbox(Inches(6.5), Inches(5.2), Inches(4), Inches(0.5))
            tx_box.text = _("*Please go to Workload Summary Slide for complete details")
            tx_box.text_frame.paragraphs[0].font.size = Pt(8)
            tx_box.text_frame.paragraphs[0].font.bold = True
            tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        sizing_conf_table = self.get_sizing_configuration()
        x = 7.29  # LEFT POSITION
        y = 0.73  # TOP POSITION
        wid = 2.6  # table Width
        heig = 3.5  # table height
        col_width = [1.0, 1.6]
        self.render_single_table(slide, x, y, wid, heig, _('Sizing Configurations'), col_width, sizing_conf_table)

        # Commenting Image as its not placed properly in case of Language support
        # hypervisor = self.scenario_data['settings_json'][0]['hypervisor']
        #
        # if hypervisor == 'hyperv':
        #     image_path = os.path.join(BASE_DIR, "sizer/hyperv_icon.png")
        #     slide.shapes.add_picture(image_path, Inches(8.37), Inches(2.39))
        # else:
        #     image_path = os.path.join(BASE_DIR, "sizer/esxi_icon.png")
        #     slide.shapes.add_picture(image_path, Inches(8.77), Inches(2.36), width=Inches(0.28), height=Inches(0.28))

        if rvtools_flag:
            tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.9), Inches(6.9), Inches(0.5))
            tx_box.text = _("NOTE: RVTool gives a point in time estimate of the usage which may lead to large errors " 
                          "in the estimates. We recommend using the option 30 Day Summary from HX Workload Profiler "
                          "to estimate the resources used by existing environments.")
            tx_box.text_frame.paragraphs[0].font.size = Pt(8)
            tx_box.text_frame.word_wrap = True
            tx_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.6), Inches(9.5), Inches(0.5))
        '''
        tx_box.text = "Portal on best practices needed for the workloads: " \
                      "\"https://www.cisco.com/c/en/us/products/hyperconverged-infrastructure/hyperflex-solutions.html#~Coredate\""
        tx_box.text_frame.word_wrap = True
        tx_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        '''
        tx_box.text_frame.paragraphs[0].font.size = Pt(9)
        tx_box.text_frame.paragraphs[0].font.bold = True
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        run = tx_box.text_frame.paragraphs[0].add_run()
        run.text = _('Click here to know more on the best practices needed for the workloads')
        run.hyperlink.address = 'https://www.cisco.com/c/en/us/products/hyperconverged-infrastructure/hyperflex-solutions.html#~Coredate'

        '''
        #Commented out bcoz aggregate Utilization is removed from 5.2 Version
        wl_require_table = self.get_aggwl_requirement()
        x = 2.9            # LEFT POSITION
        y = 4.5     # TOP POSITION
        wid = 4.5    # table Width
        heig = 0.6      # table height
        col_width = [1.0, 1.0, 1.5, 1.0]
        self.render_single_table(slide, x, y, wid, heig, 'Aggregate Workload Requirements', col_width, wl_require_table)


        left = Inches(3.4)
        top = Inches(3.9)
        width = Inches(3.2)
        height = Inches(0.3)

        shape = shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x00, 0xB0, 0xF0)

        # rotate 45 degrees counter-clockwise
        shape.rotation = 180.0
        line = shape.line
        line.color.rgb = RGBColor(0x00, 0xB0, 0xF0)
        '''

    def create_fixed_user_in_page(self, cluster_name,wl_list, fixed_setting_json):

        prs = self.prs
        slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
        shapes = slide.shapes
        shapes.title.text = _('User Inputs') +  ' (1/2)'

        wl_summary_table, rvtools_flag = self.get_wl_summary(self.scenario_data['workload_json']['wl_list'])

        x = 0.29  # LEFT POSITION
        y = 0.73  # TOP POSITION
        wid = 6.9  # table Width
        heig = 2.7  # table height
        col_width = [1.3, 0.9, 2.9, 0.6, 0.4, 0.8]

        if len(wl_summary_table) <= 2:
            heig = 0.7
        elif 2 < len(wl_summary_table) <= 4:
            heig = 1.2
        elif 4 < len(wl_summary_table) <= 6:
            heig = 2

        wl_summary_truncate_table = OrderedDict()
        if len(wl_summary_table) <= 6:
            wl_summary_truncate_table = wl_summary_table
            table_title = _("Workload Summary")
        else:
            wl_summary_truncate_table['header'] = wl_summary_table['header']
            table_title = _("Workload Summary (Top 6 workloads*)")
            for n in range(0, 6):
                key = 'row' + str(n)
                wl_summary_truncate_table[key] = wl_summary_table[key]

        self.render_single_table(slide, x, y, wid, heig, table_title, col_width, wl_summary_truncate_table)

        if len(wl_summary_table) > 6:
            tx_box = slide.shapes.add_textbox(Inches(x), Inches(4.57), Inches(4), Inches(0.5))
            tx_box.text = _("*Please go to Workload Summary Slide for complete details")
            tx_box.text_frame.paragraphs[0].font.size = Pt(8)
            tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        sizing_conf_table = self.get_sizing_configuration_fixed_sizing()
        x = 7.29  # LEFT POSITION
        y = 0.73  # TOP POSITION
        wid = 2.6  # table Width
        heig = 1.6  # table height
        col_width = [1.0, 1.6]
        self.render_single_table(slide, x, y, wid, heig, _('Sizing Configurations'), col_width, sizing_conf_table)

        # Commenting Image as its not placed properly in case of Language support
        # hypervisor = self.scenario_data['settings_json'][0]['hypervisor']
        #
        # if hypervisor == 'hyperv':
        #     image_path = os.path.join(BASE_DIR, "sizer/hyperv_icon.png")
        #     slide.shapes.add_picture(image_path, Inches(8.37), Inches(2.06))
        # else:
        #     image_path = os.path.join(BASE_DIR, "sizer/esxi_icon.png")
        #     slide.shapes.add_picture(image_path, Inches(8.77), Inches(2.03), width=Inches(0.28), height=Inches(0.28))

        if rvtools_flag:
            tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.8), Inches(9.5), Inches(0.5))
            tx_box.text = _("NOTE: RVTool gives a point in time estimate of the usage which may lead to large errors "
                          "in the estimates. We recommend using the option 30 Day Summary from HX Workload Profiler "
                          "to estimate the resources used by existing environments.")
            tx_box.text_frame.paragraphs[0].font.size = Pt(8)
            tx_box.text_frame.word_wrap = True
            tx_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.6), Inches(9.5), Inches(0.5))
        tx_box.text_frame.paragraphs[0].font.size = Pt(9)
        tx_box.text_frame.paragraphs[0].font.bold = True
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        run = tx_box.text_frame.paragraphs[0].add_run()
        run.text = _('Click here to know more on the best practices needed for the workloads')
        run.hyperlink.address = 'https://www.cisco.com/c/en/us/products/hyperconverged-infrastructure/hyperflex-solutions.html#~Coredate'


        prs = self.prs
        slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
        shapes = slide.shapes
        shapes.title.text = _('User Inputs') + ' (2/2)'

        node_conf_table = self.get_node_configuration_fixed_sizing(fixed_setting_json)
        self.render_rectangle_fixed(slide, node_conf_table)

    @staticmethod
    def get_filter_settings(settings_list):

        setting_filters_difference = [_('Sizing configuration/Filter Option used')]
        setting_filters = ''

        node_filter = settings_list[0]['filters']['Node_Type']
        if len(node_filter):
            node_filter_value = ', '.join(map(str, node_filter))
            setting_filters = "Hyperflex Nodes: " + node_filter_value + "\n"

        compute_filter = settings_list[0]['filters']['Compute_Type']
        if len(compute_filter):
            compute_filter_value = ', '.join(map(str, compute_filter))

            setting_filters += "Compute Nodes: " + compute_filter_value + "\n"

        cpu_filter = settings_list[0]['filters']['CPU_Type']
        if len(cpu_filter):
            cpu_filter_value = ', '.join(map(str, cpu_filter))
            setting_filters += "CPU (Cores, GHz): " + cpu_filter_value + "\n"

        clock_filter = settings_list[0]['filters']['Clock']
        if len(clock_filter):
            clock_filter_value = ', '.join(map(str, clock_filter))
            setting_filters += "Clock (GHz): " + clock_filter_value + "\n"

        rslots_filter = settings_list[0]['filters']['RAM_Slots']
        if len(rslots_filter):
            rslots_filter_value = ', '.join(map(str, rslots_filter))
            setting_filters += "RAM Slots: " + rslots_filter_value + "\n"

        ram_filter = settings_list[0]['filters']['RAM_Options']
        if len(ram_filter):
            ram_filter_value = ', '.join(map(str, ram_filter))
            setting_filters += "RAM Slots: " + ram_filter_value + "\n"

        disk_filter = settings_list[0]['filters']['Disk_Options']
        if len(disk_filter):
            disk_filter_value = ', '.join(map(str, disk_filter))
            setting_filters += "Disk Capacity: " + disk_filter_value + "\n"

        cache_filter = settings_list[0]['filters']['Cache_Options']
        if len(cache_filter):
            cache_filter_value = ', '.join(map(str, cache_filter))
            setting_filters += "Cache Capacity: " + cache_filter_value + "\n"

        gpu_filter = settings_list[0]['filters']['GPU_Type']
        if len(gpu_filter):
            gpu_filter_value = ', '.join(map(str, gpu_filter))
            setting_filters += "GPU: " + gpu_filter_value

        if setting_filters:
            setting_filters_difference.append(setting_filters)
            setting_filters_difference.append("")       # Keep empty for All-Flash Column
            setting_filters_difference.append("")       # Keep empty for All-NVMe Column

        return setting_filters_difference

    def get_compare_summary(self, wl_result_list, settings_list):

        compare_summary_table = OrderedDict()

        # HEADER FORMATION
        compare_summary_table['header'] = ['', _('Lowest_Cost Option'), _('All-Flash Option'), _('All NVMe Option')]
        self.compare_row_merge = False

        row_index = 0

        setting_filters_difference = self.get_filter_settings(settings_list)

        # Add entry only if difference in filter/settings is there.
        if len(setting_filters_difference) > 1:
            self.compare_row_merge = True
            row_name = "row" + str(row_index)
            compare_summary_table[row_name] = setting_filters_difference
            row_index += 1

        clusters_details = [_('Number of Clusters')]
        nodes_details = [_('Number of Nodes')]
        ru_details = [_('Number of RU')]
        cost_differential = [_('List Price Comparison')]
        ru_differential = [_('RU Comparison')]
        recommended_nodes_details = [_('Recommended Nodes')]

        cost_values = list()
        ru_values = list()

        remain_node_table = OrderedDict()

        for wl_result_data in wl_result_list:
            key = wl_result_data['result_name']

            no_of_clusters = sum([len(cluster) for cluster in wl_result_data['clusters']])
            clusters_details.append(str(no_of_clusters))

            sdata = wl_result_data['summary_info']
            num_nodes = int(sdata['num_nodes'])
            num_ft_nodes = int(sdata['num_ft_nodes'])
            nodes = num_nodes - num_ft_nodes
            node_value = str(nodes) + " + " + str(num_ft_nodes) + " (FT)"
            nodes_details.append(node_value)

            num_ru = wl_result_data['summary_info']['rack_units']
            ru_values.append(num_ru)
            ru_details.append(str(num_ru))
            cost_values.append(sdata['capex'])

            cluster_list = wl_result_data['clusters']
            flat_list = [item for sublist in cluster_list for item in sublist]
            node_details_table = self.get_nodes_details(flat_list)

            if len(node_details_table) <= 2:
                node_details_truncate_table = node_details_table
                remain_node_table[key] = list()
            else:
                node_details_truncate_table = node_details_table[0:2]
                del node_details_table[:2]
                remain_node_table[key] = node_details_table

            node_str = ""
            for node_data in node_details_truncate_table:
                node_str += node_data['header'] + ": " + node_data['count'] + " " + node_data['type'] + "\n" + \
                            node_data['description'] + "\n\n"

            recommended_nodes_details.append(node_str)

        row_name = "row" + str(row_index)
        compare_summary_table[row_name] = clusters_details
        row_index += 1

        row_name = "row" + str(row_index)
        compare_summary_table[row_name] = nodes_details
        row_index += 1

        row_name = "row" + str(row_index)
        compare_summary_table[row_name] = ru_details
        row_index += 1

        if ru_values[1] > ru_values[0]:
            rudiff = round(((ru_values[1] - ru_values[0]) * 100) / ru_values[1], 2)
            rumsg = _("Lowest_Cost Option is ") + str(rudiff) + "% " + _("lower RU than All-Flash Option")
        else:
            rudiff = round(((ru_values[0] - ru_values[1]) * 100) / ru_values[0], 2)
            rumsg = _("All-Flash Option is ") + str(rudiff) + "% " + _("lower RU than Lowest_Cost Option")

        if ru_values[2] > ru_values[0]:
            rudiff = round(((ru_values[2] - ru_values[0]) * 100) / ru_values[2], 2)
            rumsgcont = _("Lowest_Cost Option is ") + str(rudiff) + "% " + _("lower RU than All NVMe Option")
        else:
            rudiff = round(((ru_values[0] - ru_values[2]) * 100) / ru_values[0], 2)
            rumsgcont = _("All NVMe Option is ") + str(rudiff) + "% " + _("lower RU than Lowest_Cost Option")

        ru_differential.append(rumsg)
        ru_differential.append(" ")
        ru_differential.append(rumsgcont)
        row_name = "row" + str(row_index)
        compare_summary_table[row_name] = ru_differential
        row_index += 1

        if cost_values[1] > cost_values[0]:
            costdiff = round(((cost_values[1] - cost_values[0]) * 100) / cost_values[1], 2)
            costmsg = _("Lowest_Cost Option is ") + str(costdiff) + "% " + _("cheaper than All-Flash Option")
        else:
            costdiff = round(((cost_values[0] - cost_values[1]) * 100) / cost_values[0], 2)
            costmsg = _("All-Flash Option is ") + str(costdiff) + "% " + _("cheaper than Lowest_Cost Option")

        if cost_values[2] > cost_values[0]:
            costdiff = round(((cost_values[2] - cost_values[0]) * 100) / cost_values[2], 2)
            costmsgcont = _("Lowest_Cost Option is ") + str(costdiff) + "% " + _("cheaper than All NVMe Option")
        else:
            costdiff = round(((cost_values[0] - cost_values[2]) * 100) / cost_values[0], 2)
            costmsgcont = _("All NVMe Option is ") + str(costdiff) + "% " + _("cheaper than Lowest_Cost Option")

        cost_differential.append(costmsg)
        cost_differential.append(" ")
        cost_differential.append(costmsgcont)
        row_name = "row" + str(row_index)
        compare_summary_table[row_name] = cost_differential
        row_index += 1

        row_name = "row" + str(row_index)
        compare_summary_table[row_name] = recommended_nodes_details

        # Recommended nodes list is more than 3
        compare_summary_table_continue = OrderedDict()
        row_index = 0

        while len(remain_node_table) != 0:
            lowest_cost_node_table = remain_node_table['Lowest_Cost']
            allflash_node_table = remain_node_table['All-Flash']
            allnvme_node_table = remain_node_table['All NVMe']

            if len(lowest_cost_node_table) != 0 or len(allflash_node_table) != 0 or len(allnvme_node_table) != 0:
                remain_node_table = OrderedDict()

                recommended_nodes_details_cont = [_('Recommended Nodes')]

                if len(lowest_cost_node_table) <= 7:
                    node_details_truncate_table = lowest_cost_node_table
                    remain_node_table["Lowest_Cost"] = list()
                else:
                    node_details_truncate_table = lowest_cost_node_table[0:7]
                    del lowest_cost_node_table[:7]
                    remain_node_table["Lowest_Cost"] = lowest_cost_node_table

                node_str = ""
                for node_data in node_details_truncate_table:
                    node_str += node_data['header'] + ": " + node_data['count'] + " " + node_data['type'] + "\n" + \
                                node_data['description'] + "\n\n"

                recommended_nodes_details_cont.append(node_str[:-2])

                if len(allflash_node_table) <= 7:
                    node_details_truncate_table = allflash_node_table
                    remain_node_table["All-Flash"] = list()
                else:
                    node_details_truncate_table = allflash_node_table[0:7]
                    del allflash_node_table[:7]
                    remain_node_table["All-Flash"] = allflash_node_table

                node_str = ""
                for node_data in node_details_truncate_table:
                    node_str += node_data['header'] + ": " + node_data['count'] + " " + node_data['type'] + "\n" + \
                                node_data['description'] + "\n\n"

                recommended_nodes_details_cont.append(node_str[:-2])

                if len(allnvme_node_table) <= 7:
                    node_details_truncate_table = allnvme_node_table
                    remain_node_table["All NVMe"] = list()
                else:
                    node_details_truncate_table = allnvme_node_table[0:7]
                    del allnvme_node_table[:7]
                    remain_node_table["All NVMe"] = allnvme_node_table

                node_str = ""
                for node_data in node_details_truncate_table:
                    node_str += node_data['header'] + ": " + node_data['count'] + " " + node_data['type'] + "\n" + \
                                node_data['description'] + "\n\n"

                recommended_nodes_details_cont.append(node_str[:-2])

                row_name = "row" + str(row_index)
                compare_summary_table_continue[row_name] = recommended_nodes_details_cont
                row_index += 1
            else:
                break

        return compare_summary_table, compare_summary_table_continue

    def create_comparison_page(self):

        prs = self.prs
        slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
        shapes = slide.shapes
        shapes.title.text = _('Comparison of Lowest Cost Option v/s All Flash Option v/s All NVMe Option')

        wl_result_list = self.scenario_data['workload_result']
        settings_list = self.scenario_data['settings_json']
        compare_summary_table, compare_summary_table_continue = self.get_compare_summary(wl_result_list, settings_list)

        x = 0.29  # LEFT POSITION
        y = 0.50  # TOP POSITION
        wid = 6.5  # table Width
        heig = 2.5
        tabletitle = ""
        col_width = [1.1, 2.8, 2.8, 2.8]

        self.render_single_table(slide, x, y, wid, heig, tabletitle, col_width,
                                 compare_summary_table, self.compare_row_merge)

        if len(compare_summary_table_continue) != 0:
            heig = 0.6
            y = 0.45  # TOP POSITION
            # HEADER FORMATION
            compare_summary_table_continue['header'] = ['', _('Lowest_Cost Option'), _('All-Flash Option'),
                                                        _('All NVMe Option')]

            slide_title = _('Comparison of Lowest Cost Option v/s All Flash Option v/s All NVMe Option')
            self.render_table_rec_nodes(slide_title, x, y, wid, heig, col_width, compare_summary_table_continue)

    def get_wl_cluster_mapping(self, cl_list):

        wl_cluster_table = OrderedDict()
        row_index = 0
        cluster_no = 0
        drcluster_no = 0

        # HEADER FORMATION
        wl_cluster_table['header'] = [_('Workload'), _('Type'), _('Description'), _('Count'), 'RF',
                                      _('Fault Tolerance'), _('Cluster Assignment')]

        for cluster_data_index in cl_list:
            for cluster_data in cluster_data_index:

                # cluster_name = cluster_data['node_info'][0]['display_name']
                cluster_no += 1
                if len(cluster_data_index) > 1:
                    drcluster_no = cluster_no + 1

                cluster_name = self.fixed_cluster_name if self.fixed_cluster_name else "Cluster " + str(cluster_no)
                wl_list = cluster_data['wl_list']

                stretchcluster = any(
                    wl_list_tmp['cluster_type'] == 'stretch' for wl_list_tmp in cluster_data['wl_list'])

                epiccluster = any(wl_list_tmp['wl_type'] == 'EPIC' for wl_list_tmp in cluster_data['wl_list'])

                dcname = list()
                dc_count_map = defaultdict()
                if epiccluster:
                    for cluster_tmp_data in cluster_data_index:
                        epicdc = [tmp_wl_list['dc_name'] for tmp_wl_list in cluster_tmp_data['wl_list']]
                        dcname.extend(epicdc)

                    dc_count_map = dict(Counter(dcname))

                for wdata in wl_list:
                    wl_cluster_per_row = list()
                    if wdata['wl_type'] == "ROBO":
                        wl_type = "Edge"
                    else:
                        wl_type = wdata['wl_type']

                    wl_count = ''

                    if wdata['wl_type'] == "VDI":
                        if 'primary_wl_name' in wdata:
                            wl_description = "Primary workload:" + wdata['primary_wl_name'] + ", Profile:" + \
                                             wdata['profile'] + ", Note: Windows Server Licenses will be required"
                            wl_count = "{:,}".format(wdata['number_of_vms'])
                        else:
                            wl_description = wdata['profile_type'] + ", " + wdata['provisioning_type']
                            wl_count = "{:,}".format(wdata['num_desktops'])
                            if 'desktops_per_node' in wdata:
                                wl_count += " (Desktops per node for this workload = " + "{:,}".format(
                                    wdata['desktops_per_node']) + ")"
                            if 'maxdesktop' in cluster_data:
                                wl_count += " (Max possible = " + "{:,}".format(cluster_data['maxdesktop']) + ")"

                    elif wdata['wl_type'] == "VSI":
                        wl_description = wdata['profile_type']
                        wl_count = "{:,}".format(wdata['num_vms'])

                    elif wdata['wl_type'] == "RAW":
                        wl_description = ""
                        if 'input_type' in wdata:
                            wl_description = wdata['input_type']

                    elif wdata['wl_type'] == "RAW_FILE":
                        wl_description = ""
                        if 'input_type' in wdata:
                            wl_description = wdata['input_type']

                    elif wdata['wl_type'] == "DB":
                        wl_description = wdata['profile_type'] + ", " + wdata['db_type']
                        wl_count = "{:,}".format(wdata['num_db_instances'])

                    elif wdata['wl_type'] == "ROBO":
                        wl_description = wdata['profile_type']
                        wl_count = "{:,}".format(wdata['num_vms'])

                    elif wdata['wl_type'] == "ROBO_BACKUP":
                        wl_description = wdata['profile_type']
                        wl_count = "{:,}".format(wdata['num_vms'])
                        
                    elif wdata['wl_type'] == "ORACLE":
                        wl_description = wdata['profile_type'] + ", " + wdata['db_type']
                        wl_count = "{:,}".format(wdata['num_db_instances'])

                    elif wdata['wl_type'] == "EXCHANGE":
                        wl_description = ""
                        if 'input_type' in wdata:
                            wl_description = wdata['input_type']

                    elif wdata['wl_type'] == "VDI_INFRA":
                        wl_description = "Broker Type-" + wdata['infra_type']

                    elif wdata['wl_type'] == "EPIC":
                        wl_description = wdata['cpu']

                    elif wdata['wl_type'] == "VEEAM":
                        wl_description = "Total Storage Capacity Requirement-" + str(wdata['hdd_size']) + \
                                         wdata['hdd_size_unit']

                    elif wdata['wl_type'] == "SPLUNK":
                        wl_description = "Profile Type-" + wdata['profile_type']

                    elif wdata['wl_type'] == "RDSH":
                        if 'primary_wl_name' in wdata:
                            wl_description = "Primary workload:" + wdata['primary_wl_name'] + ", Profile:" + \
                                             wdata['profile']
                            wl_count = "{:,}".format(wdata['number_of_vms'])
                        else:
                            wl_description = "Profile Type-" + wdata['profile_type']
                            wl_count = "{:,}".format(wdata['total_users'])

                    elif wdata['wl_type'] == "CONTAINER":
                        wl_description = "Container Type-" + wdata['container_type']
                        wl_count = "{:,}".format(wdata['num_containers'])

                    elif wdata['wl_type'] == "AIML":

                        wl_description = "Input source-" + wdata['input_type'] + \
                                         ", Expected Utilization-" + wdata['expected_util']
                        wl_count = "{:,}".format(wdata['num_data_scientists'])

                    elif wdata['wl_type'] == "AWR_FILE":

                        wl_name = wdata['wl_name']
                        wl_description = "AWR File input, " + "Database Type-" + wdata['db_type']

                    elif wdata['wl_type'] == "ANTHOS":

                        wl_description = "Number of Pods configured: " + str(len(wdata['pod_detail']))

                    else:
                        continue
                    cluster_placement = ""
                    if len(cluster_data_index) == 1:
                        cluster_placement = cluster_name

                    elif len(cluster_data_index) > 1 and wdata['cluster_type'] == 'stretch':
                        stretchcluster_no = cluster_no + 1
                        stretchcluster_name = "Cluster " + str(stretchcluster_no)
                        cluster_placement = cluster_name + " (Stretch Pair - " + stretchcluster_name + ")"

                    elif len(cluster_data_index) > 1 and 'remote_replication_enabled' in wdata and \
                            not wdata['remote_replication_enabled'] and wdata['cluster_type'] != 'stretch':
                        cluster_placement = cluster_name + " (Site A)"

                    elif len(cluster_data_index) > 1 and 'remote_replication_enabled' in wdata and \
                            wdata['remote_replication_enabled'] and 'remote' in wdata:

                        drcluster_name = "Cluster " + str(drcluster_no)
                        wl_type_prefix = "VMs" if wdata['wl_type'] == "VSI" else "DBs"

                        if not wdata['remote']:
                            cluster_placement = cluster_name + " (Site A), "
                            cluster_placement += "DR mirror for " + str(wdata['replication_amt']) + "% " + \
                                                 wl_type_prefix + ": " + drcluster_name + " (Site B) "
                        else:
                            cluster_placement = drcluster_name + " (Site B), "
                            cluster_placement += "DR mirror for " + str(wdata['replication_amt']) + "% " + \
                                                 wl_type_prefix + ": " + cluster_name + " (Site A) "

                        cluster_placement = cluster_placement.replace(', ', '\n')

                    elif len(cluster_data_index) > 1 and epiccluster:
                        if 'DC1' in dc_count_map:
                            dc1_end = dc_count_map['DC1'] + cluster_no - 1
                            dc1_cluster = str(cluster_no) + " - " + str(dc1_end)

                        if 'DC2' in dc_count_map:
                            if 'DC1' in dc_count_map:
                                dc2_begin = dc_count_map['DC1'] + cluster_no
                            else:
                                dc2_begin = str(cluster_no)
                            dc2_end = dc2_begin + dc_count_map['DC2'] - 1
                            dc2_cluster = str(dc2_begin) + " - " + str(dc2_end)

                        if 'DC1' in dc_count_map and 'DC2' in dc_count_map:
                            cluster_placement = "Cluster [" + dc1_cluster + ", " + dc2_cluster + "]"
                        elif 'DC1' in dc_count_map:
                            cluster_placement = "Cluster [" + dc1_cluster + "]"
                        else:
                            cluster_placement = "Cluster [" + dc2_cluster + "]"

                    wl_cluster_per_row.append(wdata['wl_name'])
                    wl_cluster_per_row.append(wl_type)
                    wl_cluster_per_row.append(wl_description)
                    wl_cluster_per_row.append(wl_count)
                    wl_cluster_per_row.append(str(wdata['replication_factor']))
                    wl_fault_tolerance = "N+" + str(wdata['fault_tolerance'])
                    wl_cluster_per_row.append(wl_fault_tolerance)
                    wl_cluster_per_row.append(cluster_placement)

                    row_name = "row" + str(row_index)
                    wl_cluster_table[row_name] = wl_cluster_per_row
                    row_index += 1

                if stretchcluster:
                    cluster_no += 1

                if epiccluster:
                    cluster_no += len(cluster_data_index) - 1  # -1 since will be increamenting in loop again

                # in case of DR cluster before break increment cluster number
                if not stretchcluster and not epiccluster and len(cluster_data_index) > 1:
                    cluster_no += 1

                break
                # break because doesnt need DR cluster WL to repeat

        return wl_cluster_table

    def create_wl_cluster_map_page(self, wl_result_data):

        cluster_list = wl_result_data['clusters']

        wl_cluster_table = self.get_wl_cluster_mapping(cluster_list)

        no_of_clusters = sum([len(cluster) for cluster in wl_result_data['clusters']])
        x = 0.19  # LEFT POSITION
        y = 0.89  # TOP POSITION
        wid = 8.6  # table Width
        heig = 3.5
        tabletitle = _("Number of Clusters Created") + ": " + str(no_of_clusters)
        col_width = [1.2, 0.9, 1.9, 1.9, 0.4, 1.1, 2.4]

        if len(wl_cluster_table) <= 2:
            heig = 0.8
        elif 2 < len(wl_cluster_table) <= 4:
            heig = 1.5
        elif 4 < len(wl_cluster_table) <= 6:
            heig = 2

        slide_title = _('Assignment of Workload to Clusters')
        self.render_table_across_pages_for_mapping(slide_title, tabletitle, x, y, wid, heig, col_width,
                                                   wl_cluster_table)

        # wl_cluster_truncate_table = OrderedDict()
        # if len(wl_cluster_table) <= 11:
        #     wl_cluster_truncate_table = wl_cluster_table
        # else:
        #     wl_cluster_truncate_table['header'] = wl_cluster_table['header']
        #     tabletitle = tabletitle + " (Top 10 workloads*)"
        #     for n in range(0, 10):
        #         key = 'row' + str(n)
        #         wl_cluster_truncate_table[key] = wl_cluster_table[key]
        #
        # self.render_single_table(slide, x, y, wid, heig, tabletitle, col_width, wl_cluster_truncate_table)

    @staticmethod
    def get_sizing_summary(wl_result_data):

        sizing_summary_table = OrderedDict()
        row_index = 0

        no_of_clusters = sum([len(cluster) for cluster in wl_result_data['clusters']])
        clusters_details = [_('Number of Clusters'), str(no_of_clusters)]
        row_name = "row" + str(row_index)
        sizing_summary_table[row_name] = clusters_details
        row_index += 1

        sdata = wl_result_data['summary_info']
        num_nodes = int(sdata['num_nodes'])
        num_ft_nodes = int(sdata['num_ft_nodes'])
        nodes = num_nodes - num_ft_nodes
        node_value = str(nodes) + " + " + str(num_ft_nodes) + " (FT)"
        nodes_details = [_('Number of Nodes'), node_value]
        row_name = "row" + str(row_index)
        sizing_summary_table[row_name] = nodes_details
        row_index += 1

        num_ru = wl_result_data['summary_info']['rack_units']
        num_chassis = wl_result_data['num_chassis']
        no_of_slots_per_chassis = 8
        total_slots = num_chassis * no_of_slots_per_chassis
        free_slots = total_slots - int(total_consumed_slots)

        if int(total_consumed_slots) and num_chassis and free_slots:
            ru_details = [_('Number of RU'), str(num_ru) + "*"]
        else:
            ru_details = [_('Number of RU'), str(num_ru)]

        row_name = "row" + str(row_index)
        sizing_summary_table[row_name] = ru_details
        row_index += 1

        '''
        capex = wl_result_data['summary_info']['capex']
        total_capex = "US $" + str(capex)
        capex_details = ['Total Capex', total_capex]
        row_name = "row" + str(row_index)
        sizing_summary_table[row_name] = capex_details
        row_index += 1

        power_consumption = wl_result_data['power_consumption']
        power_consum = str(power_consumption/1000) + " KW" 
        power_details = ['Power Cooling', power_consum]
        row_name = "row" + str(row_index)
        sizing_summary_table[row_name] = power_details
        '''

        return sizing_summary_table

    @staticmethod
    def render_rectangle(slide_obj, tabledata, slot_msg):

        slide = slide_obj
        shapes = slide.shapes

        left = Inches(0.09)
        actual_left = left
        top = Inches(1.1)
        width = Inches(1.4)
        height = Inches(0.3)

        tbleft = Inches(0.05)
        tbtop = Inches(0.62)
        tbwidth = Inches(2.7)
        tbheight = Inches(0.26)
        tx_box = slide.shapes.add_textbox(tbleft, tbtop, tbwidth, tbheight)
        tx_box.text = _('Sizing Summary')
        tx_box.text_frame.paragraphs[0].font.size = Pt(12)
        tx_box.text_frame.paragraphs[0].font.bold = True
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x92, 0xC8, 0xE9)

        for rowindex in range(0, len(tabledata)):
            if rowindex % 2 == 0:
                red = green = blue = 0xF9

            for colindex in range(0, len(tabledata['row0'])):
                key = 'row' + str(rowindex)
                shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                shape.text_frame.paragraphs[0].font.size = Pt(9)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(red, green, blue)
                shape.text = tabledata[key][colindex]
                left = left + width + Inches(0.05)
            left = actual_left
            top = top + height + Inches(0.03)
            red = green = blue = 0xFF

        top = top + Inches(0.03)
        tx_box = slide.shapes.add_textbox(left, top, tbwidth, tbheight)
        tx_box.text = slot_msg
        tx_box.text_frame.paragraphs[0].font.size = Pt(8)
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    @staticmethod
    def render_rectangle_fixed(slide_obj, tabledata):

        slide = slide_obj
        shapes = slide.shapes

        left = Inches(0.29)
        actual_left = left
        top = Inches(1.1)
        width = Inches(2.3)
        height = Inches(0.3)

        tbleft = Inches(0.25)
        tbtop = Inches(0.62)
        tbwidth = Inches(2.7)
        tbheight = Inches(0.26)
        tx_box = slide.shapes.add_textbox(tbleft, tbtop, tbwidth, tbheight)
        tx_box.text = _('Node Properties')
        tx_box.text_frame.paragraphs[0].font.size = Pt(12)
        tx_box.text_frame.paragraphs[0].font.name = 'CiscoSansTT Light'
        tx_box.text_frame.paragraphs[0].font.bold = True
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x92, 0xC8, 0xE9)

        keys = ['row0', 'row1', 'row2', 'row3', 'row4']
        firsttabledata = {x: tabledata[x] for x in keys}

        for rowindex in range(0, len(firsttabledata)):
            if not rowindex % 2:
                red = green = blue = 0xF9

            for colindex in range(0, len(firsttabledata['row0'])):
                key = 'row' + str(rowindex)
                shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                shape.text_frame.paragraphs[0].font.size = Pt(9)
                shape.text_frame.paragraphs[0].font.name = 'CiscoSansTT Light'
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(red, green, blue)
                shape.text = firsttabledata[key][colindex]
                left = left + width + Inches(0.05)
            left = actual_left
            top = top + height + Inches(0.03)
            red = green = blue = 0xFF

        left = Inches(5.2)
        actual_left = left
        top = Inches(1.1)
        width = Inches(2.3)
        height = Inches(0.3)

        keys = ['row5', 'row6', 'row7', 'row8', 'row9']
        secondtabledata = {x: tabledata[x] for x in keys}

        for rowindex in range(5, len(tabledata)):
            if not rowindex % 2:
                red = green = blue = 0xF9

            for colindex in range(0, len(secondtabledata['row5'])):
                key = 'row' + str(rowindex)
                shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                shape.text_frame.paragraphs[0].font.size = Pt(9)
                shape.text_frame.paragraphs[0].font.name = 'CiscoSansTT Light'
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(red, green, blue)
                shape.text = secondtabledata[key][colindex]
                left = left + width + Inches(0.05)
            left = actual_left
            top = top + height + Inches(0.03)
            red = green = blue = 0xFF

    @staticmethod
    def get_nodes_details(cluster_list):

        node_list = list()
        accessory_list = list()
        global total_consumed_slots
        global num_chassis_per_cluster
        global cputhermnote
        cputhermnote = False
        total_consumed_slots = 0
        num_chassis_per_cluster = 0

        cluster_number = 0
        for cdata in cluster_list:
            cluster_number += 1
            node_info_list = cdata['node_info']
            for node_data in node_info_list:
                node_name = node_data['model_details']['name']
                node_desc_map = node_data['model_details']['node_description']
                subtype = node_data['model_details']['subtype']
                node_type = node_data['model_details']['type'] + " nodes"

                node_count = node_data['num_nodes']

                node_desc = list()
                node_desc.append(node_desc_map['CPU'])
                node_desc.append(node_desc_map['RAM'])
                if 'HDD' in node_desc_map:
                    node_desc.append(node_desc_map['HDD'])
                if 'SSD' in node_desc_map:
                    node_desc.append(node_desc_map['SSD'])
                if 'hx_boost_conf' in node_data and node_data['hx_boost_conf']:
                    node_desc.append("HyperFlex Boost")
                if 'hercules_conf' in node_data and node_data['hercules_conf']:
                    node_desc.append("HyperFlex Acceleration Engine")
                if 'GPU_USERS' in node_desc_map:
                    node_desc.append(node_desc_map['GPU_USERS'])

                if node_data['model_details']['cpu_part'] in therm_cpulist:
                    cputhermnote = True

                rus = list()
                rack_space = node_data['model_details']['rack_space']
                rval = str(rack_space) + " RU"
                rus.insert(0, rval)

                mod_lan = [node_data['mod_lan']] if 'mod_lan' in node_data else list()

                node_desc = node_desc + rus + mod_lan

                node_dict = dict()
                node_dict['header'] = node_name + " [Cluster " + str(cluster_number) + "]"
                node_dict['desc'] = node_desc
                node_dict['count'] = str(node_count)
                node_dict['type'] = node_type

                slot_per_blade = rack_space
                if subtype == 'compute' and node_name.find("B200") != -1:
                    if rack_space == 0.5:
                        slot_per_blade = rack_space * 2

                    consumed_slots_per_blade = slot_per_blade * node_count
                    total_consumed_slots = total_consumed_slots + consumed_slots_per_blade

                    num_chassis_per_cluster = int(math.ceil(node_count / 8.0))

                flag = False
                if not node_list:
                    node_list.append(node_dict)
                else:
                    for nlistdata in node_list:
                        if nlistdata['header'] == node_name and nlistdata['desc'] == node_desc:
                            count = int(nlistdata['count'])
                            total_count = count + int(node_count)
                            nlistdata['count'] = str(total_count)
                            flag = True
                            break

                    if not flag:
                        node_list.append(node_dict)

            if 'accessories' in cdata:
                accessory_info_list = cdata['accessories']
                for accessory_data in accessory_info_list:
                    accessory_name = accessory_data['part_name']
                    accessory_desc = accessory_data['part_description']
                    accessory_count = accessory_data['count']
                    accessory_type = accessory_data['type']

                    accessory_dict = dict()
                    accessory_dict['header'] = accessory_name
                    accessory_dict['description'] = accessory_desc
                    accessory_dict['count'] = str(accessory_count)
                    accessory_dict['type'] = accessory_type

                    flag = False
                    if not accessory_list:
                        accessory_list.append(accessory_dict)
                    else:
                        for nlistdata in accessory_list:
                            if nlistdata['header'] == accessory_name and nlistdata['description'] == accessory_desc:
                                count = int(nlistdata['count'])
                                total_count = count + int(accessory_count)
                                nlistdata['count'] = str(total_count)
                                flag = True
                                break

                        if not flag:
                            accessory_list.append(accessory_dict)

        for index in range(len(node_list)):
            node_list[index]['description'] = ' | '.join(node_list[index]['desc'])

        node_details_table = node_list + accessory_list
        return node_details_table

    @staticmethod
    def render_node_rectangle(slide_obj, leftpos, toppos, widthpos, table_title, tabledata):

        slide = slide_obj
        shapes = slide.shapes

        left = leftpos
        top = toppos
        width = widthpos

        actual_left = left
        height = Inches(0.25)

        sub_title = shapes.add_textbox(left, top - Inches(0.5), Inches(2.7), Inches(0.26))
        sub_title.text = table_title
        sub_title.text_frame.paragraphs[0].font.size = Pt(12)
        sub_title.text_frame.paragraphs[0].font.bold = True
        sub_title.text_frame.paragraphs[0].font.name = 'CiscoSansTT Light'
        sub_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x92, 0xC8, 0xE9)

        rowindex = 1
        red = green = blue = 0xF9

        for data in tabledata:
            if rowindex % 2 == 0:
                red = green = blue = 0xFF

            shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.text_frame.paragraphs[0].font.size = Pt(10)
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(red, green, blue)
            shape.text = data['header']

            desctop = top + height
            descheight = Inches(0.5)

            shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, desctop, width, descheight)
            shape.text_frame.paragraphs[0].font.size = Pt(8)
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(red, green, blue)
            shape.text = data['description']

            c_left = left + width
            c_height = Inches(0.75)
            c_width = Inches(0.4)

            shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, top, c_width, c_height)

            shape.text_frame.paragraphs[0].font.size = Pt(10)
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xDE, 0x79, 0x79)
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(red, green, blue)
            shape.text = data['count']

            rowindex += 1
            left = actual_left
            top = top + height + Inches(0.6)
            red = green = blue = 0xF9

    def create_sizing_results_page(self, wl_result_data):

        # RENDER NODE

        # Get node result
        cluster_list = wl_result_data['clusters']
        flat_list = [item for sublist in cluster_list for item in sublist]
        node_details_table = self.get_nodes_details(flat_list)

        # Get Sizing Summary
        sizing_summary_table = self.get_sizing_summary(wl_result_data)

        # calculate No of free chassis slots available
        num_chassis = wl_result_data['num_chassis']
        no_of_slots_per_chassis = 8
        total_slots = num_chassis * no_of_slots_per_chassis
        free_slots = total_slots - int(total_consumed_slots)

        slot_msg = ''

        if int(total_consumed_slots) and num_chassis and free_slots:
            slot_msg = '*includes ' + str(free_slots) + ' free chassis slots'

        # prs = self.prs
        # slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
        # shapes = slide.shapes
        # shapes.title.text = 'Overall Sizing Results'
        #
        # self.render_rectangle(slide, sizing_summary_table, slot_msg)

        # Below Code is commented as it is redundant information
        len_nodes = len(node_details_table)
        nodes_per_slide = 4

        if len_nodes % nodes_per_slide == 0:
            total_slides = len_nodes // nodes_per_slide
        else:
            total_slides = len_nodes // nodes_per_slide + 1

        slide_num = 1

        for index in range(0, int(total_slides)):
            prs = self.prs
            slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
            shapes = slide.shapes
            shapes.title.text = _('Overall Sizing Results')
            leftpos = Inches(3.21)
            toppos = Inches(1.1)
            widthpos = Inches(6.2)

            if len(node_details_table) <= 4:
                node_details_trucate_table = node_details_table
            else:
                node_details_trucate_table = node_details_table[0:4]
                del node_details_table[:4]

            table_title = _("Recommended Nodes") + " (" + str(slide_num) + "/" + str(total_slides) + ")"
            self.render_node_rectangle(slide, leftpos, toppos, widthpos, table_title, node_details_trucate_table)

            self.render_rectangle(slide, sizing_summary_table, slot_msg)
            slide_num += 1

            if cputhermnote:
                tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.9), Inches(9.5), Inches(0.3))
                tx_box.text = _(
                    "Note: Lower ambient temperature required 30deg Celsius for All NVMe and 32deg Celsius for HX*240 "
                    "in case of CPUs 6242R, 6246R, 6248R and 6258R")
                tx_box.text_frame.paragraphs[0].font.size = Pt(9)
                tx_box.text_frame.word_wrap = True
                tx_box.text_frame.paragraphs[0].font.bold = True
                tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)


    def create_backup_page(self):

        # SLIDE 10 - PageBreak
        prs = self.prs

        slide3 = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BACKUP_LAYOUT_INDEX])
        tbleft = Inches(3.9)
        tbtop = Inches(2.4)
        tbwidth = Inches(2)
        tbheight = Inches(0.5)
        tx_box = slide3.shapes.add_textbox(tbleft, tbtop, tbwidth, tbheight)
        tx_box.text = 'Backup'
        tx_box.text_frame.paragraphs[0].font.size = Pt(24)
        tx_box.text_frame.paragraphs[0].font.bold = True
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    def create_page_break(self, text):

        prs = self.prs
        slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BIG_STMT_LAYOUT_INDEX])
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = text

    @staticmethod
    def render_util_charts(slide_obj, table_title, tag, chart1_x, chart1_y, maximum_scale, series_dict_graychart,
                           wid, heig, units, sizing_threshold):

        chart_slide = slide_obj

        left = Inches(chart1_x)
        top = Inches(chart1_y)
        width = Inches(wid)
        height = Inches(heig)

        tx_width = Inches(2.0)
        tx_left = left
        tx_top = top + Inches(0.2)
        tx_box = chart_slide.shapes.add_textbox(tx_left, tx_top, tx_width, height)
        tx_box.text = table_title
        tx_box.text_frame.paragraphs[0].font.size = Pt(10)
        tx_box.text_frame.paragraphs[0].font.bold = False

        chart1_data = ChartData()
        chart1_data.categories = [tag]

        for k, v in series_dict_graychart.items():
            chart1_data.add_series(k, (v,))

        x = left
        y = top + Inches(0.9)
        cx = width
        cy = Inches(1.0)

        chart1_graphic_frame = chart_slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart1_data
        )

        chart1_category_axis = chart1_graphic_frame.chart.category_axis
        chart1_category_axis.has_major_gridlines = False
        chart1_category_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
        chart1_category_axis.tick_label_position = XL_TICK_LABEL_POSITION.HIGH

        chart1_value_axis = chart1_graphic_frame.chart.value_axis
        chart1_value_axis.has_major_gridlines = False
        chart1_value_axis.has_minor_gridlines = False
        chart1_value_axis.maximum_scale = maximum_scale
        chart1_value_axis.minimum_scale = 0
        chart1_value_axis.minor_tick_mark = XL_TICK_MARK.NONE
        chart1_value_axis.major_tick_mark = XL_TICK_MARK.NONE
        chart1_value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NEXT_TO_AXIS
        chart1_value_axis.visible = False
        chart1_category_axis.visible = False

        chart1 = chart1_graphic_frame.chart
        plot1 = chart1.plots[0]
        plot1.has_data_labels = True
        chart1_data_labels = plot1.data_labels
        # data_labels.number_format = '0 "Cores"'
        chart1_data_labels.number_format = '0 "' + units + '"'
        if tag == 'Storage IOPS':
            chart1_data_labels.number_format = '0 "%"'

        chart1_data_labels.font.size = Pt(8)
        chart1_data_labels.font.bold = False
        chart1_data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        chart1_data_labels.position = XL_LABEL_POSITION.CENTER
        chart1.has_legend = True
        chart1.legend.font.size = Pt(9)
        chart1.legend.include_in_layout = True
        chart1.legend.position = XL_LEGEND_POSITION.TOP

        chart1_colours = ['C0EA80', 'F8F8F8', 'F8F8F8', 'D3D3D3']
        for i in range(len(chart1_data)):
            ser = plot1.series[i]
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = RGBColor.from_string(chart1_colours[i])
            ser.format.line.color.rgb = RGBColor(0x00, 0x00, 0x00)

        tx_width = Inches(4.6)
        tx_left = Inches(chart1_x + 0.05)
        tx_top = Inches(chart1_y + 1.65)
        tx_height = Inches(0.40)
        tx_box = chart_slide.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
        tx_box.text = sizing_threshold
        tx_box.text_frame.paragraphs[0].font.size = Pt(9)
        tx_box.text_frame.word_wrap = True
        tx_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)

        '''
        #ADDING CHART2 AT POSITION

        left = Inches(chart2_x)
        top = Inches(chart2_y)
        width = Inches(wid)
        height = Inches(heig)

        chart2_data = ChartData()
        chart2_data.categories = [tag]
        for k, v in series_dict.items():
            if (float(v) <= 0.0):                   #For zero utils just to have one small line.
                v=0.4
            chart2_data.add_series(k, (v, ))

        x = left
        y = top + Inches(0.6)
        cx = width
        cy = Inches(1)

        chart2_graphic_frame = chart_slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart2_data
        )

        chart2_category_axis = chart2_graphic_frame.chart.category_axis
        chart2_category_axis.has_major_gridlines = False
        #category_axis.has_minor_gridlines = False
        chart2_category_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
        chart2_category_axis.tick_label_position = XL_TICK_LABEL_POSITION.HIGH

        chart2_value_axis = chart2_graphic_frame.chart.value_axis
        chart2_value_axis.has_major_gridlines = False
        chart2_value_axis.has_minor_gridlines = False
        chart2_value_axis.minimum_scale = 0
        chart2_value_axis.maximum_scale = 100.0
        chart2_value_axis.minor_tick_mark = XL_TICK_MARK.NONE
        chart2_value_axis.major_tick_mark = XL_TICK_MARK.NONE
        chart2_value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NEXT_TO_AXIS
        chart2_value_axis.visible=False
        chart2_category_axis.visible=False


        chart2 = chart2_graphic_frame.chart
        chart2_plot = chart2.plots[0]
        chart2_plot.has_data_labels = True
        chart2_data_labels = chart2_plot.data_labels
        #data_labels.number_format = '0.0"%"'
        chart2_data_labels.number_format = '0"%"'
        chart2_data_labels.font.size = Pt(8)
        chart2_data_labels.font.bold = False
        chart2_data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        chart2_data_labels.position = XL_LABEL_POSITION.CENTER
        #chart.chart_style = 31
        chart2.has_legend = True
        chart2.legend.font.size = Pt(9)
        chart2.legend.include_in_layout = False
        chart2.legend.position = XL_LEGEND_POSITION.BOTTOM

        chart2_listofColours = ['C0EA80', 'D6F9A1', 'F4D589', 'F2F25C']
        for i in range(len(chart2_data)):
            ser = chart2_plot.series[i]
            ser.fill.solid()
            ser.fill.fore_color.rgb = RGBColor.from_string(chart2_listofColours[i])
        '''

    def render_util(self, slide, slide_title, util_list, rep_factor, wl_type, ft, wl_list):

        shapes = slide.shapes
        shapes.title.text = slide_title

        cpu = slide.placeholders[10]
        memory = slide.placeholders[11]
        disk = slide.placeholders[12]
        iops = slide.placeholders[13]

        wid = 4.5
        heig = 0.5

        iops_conversion_factor = max(wl[BaseConstants.IOPS_CONV_FAC] for wl in wl_list)
        total_iops = int(iops_conversion_factor * self.num_nodes / 1000.0)

        wl_types_cluster = list(set([wl['wl_type'] for wl in wl_list]))
        vdi_directory_enabled = False
        if wl_type == 'VDI':
            vdi_directory_enabled = any(wl.get('vdi_directory', False) for wl in wl_list)

        if wl_type == 'RAW' or wl_type == 'RAW_FILE':
            block_size = list(set([wl['io_block_size'] for wl in wl_list if 'io_block_size' in wl]))

            if len(block_size) == 1:
                IOPS_DESC[wl_type] = IOPS_DESC[block_size[0]]
            elif len(block_size) > 1:
                IOPS_DESC[wl_type] = "Blended IOPS Mix"
            else:
                IOPS_DESC[wl_type] = "Performance was not modelled for this workload"

            raw_total_iops = sum(list([wl['iops_value'] if 'iops_value' in wl else 0 for wl in wl_list ]))

        for util_data in util_list:

            tag_name = util_data['tag_name']
            sizing_threshold = ""
            if util_data['status'] == True:
                node_val = util_data['node_val']
                sizing_thr = util_data['best_practice']

                healthy_util = int(round(util_data['wl_util']))
                threshold_util = int(ceil(util_data['threshold_util']))
                free_util = int(round(util_data['free_util']))

                tag_name = util_data['tag_name']
                series_dict_graychart = OrderedDict()

                if tag_name == BaseConstants.CPU:
                    display_node_val = int(ceil(node_val))
                    display_sizing_thr = int(ceil(sizing_thr))
                else:
                    display_node_val = round(node_val, 1)
                    display_sizing_thr = round(sizing_thr, 1)

            if util_data['status'] == True and (tag_name == 'Physical Cores' or tag_name == 'CPU'):

                cpu.text = tag_name

                table_title = _("Total Usable* CPUs") + " = " + "{:,}".format(display_node_val) + " " + util_data['units']
                title_tag = "\n(" + _("After System overhead, CPU overprovisioning ratio") + " = " + str(util_data['ratio'])

                if ft:
                    title_tag += ",\n" + _("Performance headroom") + " = " + _("Sustain ") \
                                 + str(ft) + _(" node failure") + ")"
                else:
                    title_tag += ")"

                table_title = table_title + title_tag

                workload = util_data['workload_val']
                available = node_val - workload
                series_dict_graychart[_('Workload')] = round(workload)
                series_dict_graychart[_('Unused')] = round(available)
                maximum_scale = round(workload) + round(available)
                chart1_x = 0.29
                chart1_y = 0.55

                sizing_threshold = _("Best practice* sizing threshold") + ": " + "{:,}".format(display_sizing_thr) + " " + \
                                   util_data['units']

            elif util_data['status'] == True and tag_name == 'RAM':

                memory.text = tag_name
                table_title = _("Total Usable* RAM") + " = " + "{:,}".format(display_node_val) + " " + util_data['units']
                title_tag = "\n(" + _("After System overhead, RAM overprovisioning ratio") + " = " + str(util_data['ratio'])

                if ft:
                    title_tag += ",\n" + _("Performance headroom") + " = " \
                                 + _("Sustain ") + str(ft) + _(" node failure") + ")"
                else:
                    title_tag += ")"

                table_title = table_title + title_tag

                workload = util_data['workload_val']
                available = node_val - workload
                series_dict_graychart[_('Workload')] = round(workload)
                series_dict_graychart[_('Unused')] = round(available)
                maximum_scale = round(workload) + round(available)
                chart1_x = 5.11
                chart1_y = 0.55
                sizing_threshold = _("Best practice* sizing threshold") + ": " + "{:,}".format(display_sizing_thr) + " " + \
                                   util_data['units']

            elif util_data['status'] == True and tag_name == 'Storage Capacity':

                disk.text = _(tag_name)
                usable_val = round(util_data['usable'], 1)
                usable_binarybyte = round(util_data['usable_binarybyte'], 1)
                node_val_binarybyte = round(util_data['node_val_binarybyte'], 1)

                table_title = _("Total Usable* Capacity") + " = " + "{:,}".format(usable_val) + " " + util_data['units'] + \
                              ' (' + str(usable_binarybyte) + " " + util_data['binarybyte_unit'] + ")\n"

                table_title += _("Total Effective* Capacity") + " = " + "{:,}".format(display_node_val) + " " + \
                               util_data['units'] + ' (' + str(node_val_binarybyte) + " " + \
                               util_data['binarybyte_unit'] + ")\n"

                title_tag = "(" + _("After Replication Factor") + " = " + str(rep_factor) + \
                            ", " + _("System overhead") + ",\n " + _("Effective storage savings") + "= " + \
                            str(round(util_data['ratio'], 1)) + "%)"

                table_title = table_title + title_tag

                # workload = (healthy_util/100) * node_val
                workload = util_data['workload_val']
                available = node_val - workload
                series_dict_graychart[_('Workload')] = round(workload)
                series_dict_graychart[_('Free Space Available')] = round(available)
                maximum_scale = round(workload) + round(available)
                chart1_x = 0.29
                chart1_y = 2.88
                best_practice_binarybyte = round(util_data['best_practice_binarybyte'], 1)
                sizing_threshold = _("Best practice* sizing threshold") + ": " + "{:,}".format(display_sizing_thr) + " " + \
                                   util_data['units'] + " (" + "{:,}".format(best_practice_binarybyte) + " " + \
                                   util_data['binarybyte_unit'] + ") " + _("of Effective Capacity")

                if self.hercules_conf:
                    sizing_threshold += "\n" + _("Effective capacity has been increased from Hardware Acceleration.")

            elif util_data['status'] == True and tag_name == 'Storage IOPS':

                iops.text = _(tag_name)
                # table_title = "Total Usable = " + str(node_val) + " " + util_data['units']
                table_title = _("Total available IOPS") + " = " + str(total_iops) + "K\n"
                if self.hercules_conf:
                    table_title = _("Total available IOPS") + "** = " + str(total_iops) + "K\n"

                if len(wl_types_cluster) > 1 or vdi_directory_enabled:
                    table_title += _("Blended IOPS Mix")
                else:
                    if wl_type in ["DB", "ORACLE"] and wl_list[0]['db_type'] == "OLAP":
                        table_title += _(IOPS_DESC[wl_list[0]['db_type']])
                    else:
                        table_title += _(IOPS_DESC[wl_type])

                # if wl_type == "EXCHANGE":
                #     table_title += "\n(Uses a weighted average to calculate equivalent 8k IOPS)"

                series_dict_graychart[_('Workload')] = round(healthy_util)
                series_dict_graychart[_('Unused / Performance Headroom')] = round(free_util)
                chart1_x = 5.11
                chart1_y = 2.88
                maximum_scale = round(healthy_util) + round(free_util)

                if self.hercules_conf:
                    sizing_threshold = _("**Total available IOPS include the additional IOPS provided by the Hardware Acceleration Card. ")

                if len(wl_types_cluster) > 1 or vdi_directory_enabled:
                    sizing_threshold += _("This is an aggregate workload derived from the weighted average of the different workloads")

                # RAW Workload has IOPS values
                if wl_type == "RAW" or wl_type == 'RAW_FILE':
                    if not raw_total_iops:
                        series_dict_graychart = dict()
                        maximum_scale = 0
                        table_title = "No performance data available."

            elif util_data['status'] == False:
                if tag_name == 'Physical Cores' or tag_name == 'CPU':
                    chart1_x = 0.29
                    chart1_y = 0.55
                    table_title = _("This solution does not allow running user VMs on the HX cluster")
                elif tag_name == 'RAM':
                    chart1_x = 5.11
                    chart1_y = 0.55
                    table_title = _("This solution does not allow running user VMs on the HX cluster")
                elif tag_name == 'Storage Capacity':
                    chart1_x = 0.29
                    chart1_y = 2.88
                    table_title = _(IOPS_DESC[wl_type])
                elif tag_name == 'Storage IOPS':
                    iops.text = _(tag_name)
                    chart1_x = 5.11
                    chart1_y = 2.88
                    table_title = _(IOPS_DESC[wl_type])

                series_dict_graychart = dict()
                maximum_scale = 0
                util_data['units'] = ""
            else:
                continue

            self.render_util_charts(slide, table_title, tag_name, chart1_x, chart1_y, maximum_scale,
                                    series_dict_graychart, wid, heig,
                                    util_data['units'], sizing_threshold)

        drenabled_cluster = any(each_wl['remote_replication_enabled'] == True for each_wl in wl_list \
                                if 'remote_replication_enabled' in each_wl)

        if drenabled_cluster:
            tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(5), Inches(9), Inches(0.3))
            tx_box.text = \
                _("Note: It is recommended that the bandwidth between the DR clusters be limited to 1Gbps. Higher "
                "bandwidth may result in higher latencies on the primary workload.")

            tx_box.text_frame.paragraphs[0].font.size = Pt(9)
            tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.24), Inches(9), Inches(0.3))
        tx_box.text = _("Expected cluster utilization based on total resources required by all workloads "
                      "and total available resources on this cluster")
        tx_box.text_frame.paragraphs[0].font.size = Pt(10)
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

        tx_box = slide.shapes.add_textbox(Inches(4.25), Inches(5.2), Inches(4), Inches(0.3))
        tx_box.text = _("* Definitions for Usable, Effective and Best Practice defined in the Glossary in the Appendix.")
        tx_box.text_frame.paragraphs[0].font.size = Pt(9)
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    def render_table_across_pages(self, slide_title, x, y, wid, hei, col_width, tabledata, rows_per_slide):

        prs = self.prs

        lentable = len(tabledata) - 1

        if lentable % rows_per_slide == 0:
            total_slides = lentable // rows_per_slide
        else:
            total_slides = lentable // rows_per_slide + 1

        slide_num = 1

        for rowindex in range(0, lentable):
            if rowindex % rows_per_slide == 0:
                slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
                shapes = slide.shapes

                slide_num_page = " (" + str(slide_num) + "/" + str(total_slides) + ")"
                shapes.title.text = slide_title + slide_num_page
                slide_num += 1

                # Add table
                rows = rows_per_slide + 1
                remaining_records = lentable - rowindex
                if remaining_records < rows_per_slide:
                    rows = remaining_records + 1

                cols = len(tabledata['header'])
                left = Inches(x)
                top = Inches(y)
                width = Inches(wid)
                height = Inches(hei)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                # set column widths
                for colindex in range(0, cols):
                    table.columns[colindex].width = Inches(col_width[colindex])

                for colindex in range(0, cols):
                    table.cell(0, colindex).text = tabledata['header'][colindex]
                    table.cell(0, colindex).text_frame.paragraphs[0].font.size = Pt(9)
                    table.cell(0, colindex).fill.solid()
                    table.cell(0, colindex).fill.fore_color.rgb = RGBColor(0xB2, 0xB2, 0xB2)

                if self.workload_note:
                    tx_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.8), Inches(9), Inches(0.3))
                    tx_box.text = \
                        "Large VMs with more than 128 GB of RAM may suffer performance degradation due to not " \
                        "following VMware Best Practices. Refer to " \
                        "https://blogs.vmware.com/performance/2017/03/" \
                        "virtual-machine-vcpu-and-vnuma-rightsizing-rules-of-thumb.html for more details."

                    tx_box.text_frame.paragraphs[0].font.size = Pt(10)
                    tx_box.text_frame.word_wrap = True
                    tx_box.text_frame.paragraphs[0].font.bold = True
                    tx_box.text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
                    tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

            row_position = (rowindex % rows_per_slide) + 1

            for colindex in range(0, cols):
                key = 'row' + str(rowindex)
                table.cell(row_position, colindex).text = tabledata[key][colindex]
                table.cell(row_position, colindex).text_frame.paragraphs[0].font.size = Pt(8)

    def render_table_across_pages_for_mapping(self, slide_title, tabletitle, x, y, wid, heig, col_width, tabledata):
        prs = self.prs

        lentable = len(tabledata) - 1

        rows_per_slide = 10
        if lentable % rows_per_slide == 0:
            total_slides = lentable // rows_per_slide
        else:
            total_slides = lentable // rows_per_slide + 1

        slide_num = 1

        for rowindex in range(0, lentable):
            if not rowindex % rows_per_slide:
                slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
                shapes = slide.shapes

                slide_num_page = " (" + str(slide_num) + "/" + str(total_slides) + ")"
                shapes.title.text = slide_title + slide_num_page
                slide_num += 1

                # Add table
                rows = rows_per_slide + 1
                remaining_records = lentable - rowindex
                if remaining_records < rows_per_slide:
                    rows = remaining_records + 1

                cols = len(tabledata['header'])
                left = Inches(x)
                top = Inches(y)
                width = Inches(wid)

                if remaining_records <= 2:
                    heig = 0.8
                elif 2 < remaining_records <= 4:
                    heig = 1.5
                elif 4 < remaining_records <= 6:
                    heig = 2
                else:
                    heig = 3.5

                height = Inches(heig)

                tbtop = top - Inches(0.64)
                tx_box = slide.shapes.add_textbox(left, tbtop, Inches(2), Inches(1))
                tf = tx_box.text_frame
                p = tf.add_paragraph()
                p.text = tabletitle
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0x92, 0xC8, 0xE9)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                # set column widths
                for colindex in range(0, cols):
                    table.columns[colindex].width = Inches(col_width[colindex])

                for colindex in range(0, cols):
                    table.cell(0, colindex).text = tabledata['header'][colindex]
                    table.cell(0, colindex).text_frame.paragraphs[0].font.size = Pt(9)
                    table.cell(0, colindex).fill.solid()
                    table.cell(0, colindex).fill.fore_color.rgb = RGBColor(0xB2, 0xB2, 0xB2)

            row_position = (rowindex % rows_per_slide) + 1

            for colindex in range(0, cols):
                key = 'row' + str(rowindex)
                table.cell(row_position, colindex).text = tabledata[key][colindex]
                table.cell(row_position, colindex).text_frame.paragraphs[0].font.size = Pt(8)

    def render_table_rec_nodes(self, slide_title, x, y, wid, hei, col_width, tabledata):

        prs = self.prs

        lentable = len(tabledata) - 1

        rows_per_slide = 1

        slide_num = 1

        for rowindex in range(0, lentable):
            if not rowindex % rows_per_slide:
                slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
                shapes = slide.shapes

                shapes.title.text = slide_title
                slide_num += 1

                # Add table
                rows = rows_per_slide + 1
                remaining_records = lentable - rowindex
                if remaining_records < rows_per_slide:
                    rows = remaining_records + 1

                cols = len(tabledata['header'])
                left = Inches(x)
                top = Inches(y)
                width = Inches(wid)
                height = Inches(hei)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                for colindex in range(0, cols):
                    table.columns[colindex].width = Inches(col_width[colindex])

                for colindex in range(0, cols):
                    table.cell(0, colindex).text = tabledata['header'][colindex]
                    table.cell(0, colindex).text_frame.paragraphs[0].font.size = Pt(9)
                    table.cell(0, colindex).fill.solid()
                    table.cell(0, colindex).fill.fore_color.rgb = RGBColor(0xB2, 0xB2, 0xB2)

            row_position = (rowindex % rows_per_slide) + 1

            for colindex in range(0, cols):
                key = 'row' + str(rowindex)
                table.cell(row_position, colindex).text = tabledata[key][colindex]
                table.cell(row_position, colindex).text_frame.paragraphs[0].font.size = Pt(8)

    def render_calc_table_across_pages(self, slide_title, x, y, wid, hei, col_width, tabledata, wl_total,
                                       rows_per_slide, op_ratio_note):

        prs = self.prs

        lentable = len(tabledata) - 1

        if lentable % rows_per_slide == 0:
            total_slides = lentable // rows_per_slide
        else:
            total_slides = lentable // rows_per_slide + 1

        slide_num = 1
        row_position = 0

        for rowindex in range(0, lentable):
            if rowindex % rows_per_slide == 0:
                slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
                shapes = slide.shapes

                slide_num_page = " (" + str(slide_num) + "/" + str(total_slides) + ")"
                shapes.title.text = slide_title + slide_num_page
                slide_num += 1

                # Add table
                rows = rows_per_slide + 1
                remaining_records = lentable - rowindex
                if remaining_records <= rows_per_slide:
                    rows = remaining_records + 1 + 1    # 1 for header, 1 for last Total workloads row

                cols = len(tabledata['header'])     # +3 because for 3 headers two columns are merged
                left = Inches(x)
                top = Inches(y)
                width = Inches(wid)
                height = Inches(hei)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                # set column widths
                row_0 = 0
                table.cell(row_0, 0).text = tabledata['header'][0]
                table.cell(row_0, 1).text = tabledata['header'][1]
                table.cell(row_0, 3).text = tabledata['header'][3]
                table.cell(row_0, 5).text = tabledata['header'][5]

                for colindex in range(0, cols):
                    table.columns[colindex].width = Inches(col_width[colindex])

                for colindex in range(0, cols):
                    table.cell(row_0, colindex).text_frame.paragraphs[0].font.size = Pt(9)
                    table.cell(row_0, colindex).text_frame.paragraphs[0].font.bold = True
                    table.cell(row_0, colindex).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    table.cell(row_0, colindex).fill.solid()
                    table.cell(row_0, colindex).fill.fore_color.rgb = RGBColor(0xB2, 0xB2, 0xB2)

                # merge cells horizontally
                self.merge_cells_horizontally(table=table, row_idx=row_0, start_col_idx=1, end_col_idx=2)
                self.merge_cells_horizontally(table=table, row_idx=row_0, start_col_idx=3, end_col_idx=4)
                self.merge_cells_horizontally(table=table, row_idx=row_0, start_col_idx=5, end_col_idx=6)

            row_position = (rowindex % rows_per_slide) + 1

            for colindex in range(0, cols):
                key = 'row' + str(rowindex)
                table.cell(row_position, colindex).text = tabledata[key][colindex]
                table.cell(row_position, colindex).text_frame.paragraphs[0].font.size = Pt(7)
                table.cell(row_position, colindex).text_frame.paragraphs[0].font.italic = True
                table.cell(row_position, colindex).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                table.cell(row_position, colindex).vertical_anchor = MSO_ANCHOR.MIDDLE

                if row_position == 3 or row_position == 6 or row_position == 9:
                    table.cell(row_position, colindex).fill.solid()
                    table.cell(row_position, colindex).fill.fore_color.rgb = RGBColor(0xB0, 0xE0, 0xE6)
                    table.cell(row_position, colindex).text_frame.paragraphs[0].font.italic = False

                if row_position == 3:
                    self.merge_cells_vertically(table=table, start_row_idx=1, end_row_idx=3, col_idx=0)

                if row_position == 6:
                    self.merge_cells_vertically(table=table, start_row_idx=4, end_row_idx=6, col_idx=0)

                if row_position == 9:
                    self.merge_cells_vertically(table=table, start_row_idx=7, end_row_idx=9, col_idx=0)

        row_index = row_position + 1    # Last row
        table.cell(row_index, 0).text = _("Total Workload Requirement")
        for colindex in range(0, cols):
            table.cell(row_index, colindex).fill.solid()
            table.cell(row_index, colindex).fill.fore_color.rgb = RGBColor(0xB0, 0xE0, 0xE6)
            table.cell(row_index, colindex).text_frame.paragraphs[0].font.italic = False
            table.cell(row_index, colindex).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(row_index, colindex).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(row_index, colindex).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(row_index, colindex).text_frame.paragraphs[0].font.bold = True

        if op_ratio_note:
            table.cell(row_index, 1).text = str(round(wl_total['CPU'], 2)) + " (*Ref Note)"
        else:
            table.cell(row_index, 1).text = str(round(wl_total['CPU'], 2))
        table.cell(row_index, 3).text = str(round(wl_total['RAM'], 2))
        table.cell(row_index, 5).text = str(round(wl_total['HDD'], 2))

        # merge cells horizontally
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=1, end_col_idx=2)
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=3, end_col_idx=4)
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=5, end_col_idx=6)

        tx_box = slide.shapes.add_textbox(Inches(0.1), Inches(4.8), Inches(9), Inches(0.3))
        tx_box.text = op_ratio_note
        tx_box.text_frame.paragraphs[0].font.size = Pt(9)
        tx_box.text_frame.paragraphs[0].font.bold = True
        tx_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    @staticmethod
    def get_wl_details(cl_list):

        wl_details_table = OrderedDict()
        row_index = 0

        # HEADER FORMATION
        wl_details_table['header'] = [_('Workload'), _('Type'), _('Details'), _('Cluster Parameters')]

        for cluster_data_index in cl_list:
            for cluster_data in cluster_data_index:

                # cluster_name = cluster_data['node_info'][0]['display_name']
                wl_list = cluster_data['wl_list']

                for wdata in wl_list:
                    wl_cluster_per_row = list()
                    wl_cluster_per_row.append(wdata['wl_name'])
                    if wdata['wl_type'] == "ROBO":
                        wl_cluster_per_row.append("Edge")
                    else:
                        wl_cluster_per_row.append(wdata['wl_type'])

                    details_dict = OrderedDict()

                    if wdata['wl_type'] == "VDI":
                        # skip VDI_HOME workload, info included in main workload
                        if 'primary_wl_name' in wdata:
                            continue

                        for key in range(0, len(VDI_KEY_LIST)):
                            if VDI_KEY_LIST[key] in wdata:
                                if VDI_KEY_LIST[key] == 'disk_per_desktop':
                                    key_holder = VDI_KEY_HOLDER[key] + '(' + wdata['disk_per_desktop_unit'] + ')'
                                    details_dict[key_holder] = wdata[VDI_KEY_LIST[key]]
                                elif VDI_KEY_LIST[key] == 'gold_image_size':
                                    key_holder = VDI_KEY_HOLDER[key] + '(' + wdata['gold_image_size_unit'] + ')'
                                    details_dict[key_holder] = wdata[VDI_KEY_LIST[key]]
                                elif VDI_KEY_LIST[key] == 'ram_per_desktop':
                                    key_holder = VDI_KEY_HOLDER[key] + '(' + wdata['ram_per_desktop_unit'] + ')'
                                    details_dict[key_holder] = wdata[VDI_KEY_LIST[key]]
                                else:
                                    # details_dict[VDI_KEY_HOLDER[key]] = wdata[VDI_KEY_LIST[key]]
                                    details_dict[VDI_KEY_HOLDER[key]] = "{:,}".format(wdata[VDI_KEY_LIST[key]]) \
                                        if type(wdata[VDI_KEY_LIST[key]]) is int else wdata[VDI_KEY_LIST[key]]
                        if wdata['provisioning_type'] == 'Pooled Desktops':
                            details_dict['Dedupe savings for inflight data(%)'] = 100 - wdata['inflight_dedupe_factor']
                        else:
                            details_dict['Dedupe savings for OS data(%)'] = 100 - wdata['image_dedupe_factor']

                    elif wdata['wl_type'] == "VSI":
                        for key in range(0, len(VSI_KEY_LIST)):
                            if wdata.get(VSI_KEY_LIST[key]):
                                if VSI_KEY_LIST[key] == 'disk_per_vm':
                                    key_holder = VSI_KEY_HOLDER[key] + '(' + wdata['disk_per_vm_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[VSI_KEY_LIST[key]])
                                elif VSI_KEY_LIST[key] == 'base_image_size':
                                    key_holder = VSI_KEY_HOLDER[key] + '(' + wdata['base_image_size_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[VSI_KEY_LIST[key]])
                                elif VSI_KEY_LIST[key] == 'ram_per_vm':
                                    key_holder = VSI_KEY_HOLDER[key] + '(' + wdata['ram_per_vm_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[VSI_KEY_LIST[key]])
                                else:
                                    # details_dict[VSI_KEY_HOLDER[key]] = wdata[VSI_KEY_LIST[key]]
                                    details_dict[VSI_KEY_HOLDER[key]] = "{:,}".format(wdata[VSI_KEY_LIST[key]]) \
                                        if type(wdata[VSI_KEY_LIST[key]]) is int else wdata[VSI_KEY_LIST[key]]

                    elif wdata['wl_type'] == "RAW" or wdata['wl_type'] == "RAW_FILE":
                        for key in range(0, len(RAW_KEY_LIST)):
                            if RAW_KEY_LIST[key] in wdata:
                                if RAW_KEY_LIST[key] == 'hdd_size':
                                    key_holder = RAW_KEY_HOLDER[key] + '(' + wdata['hdd_size_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[RAW_KEY_LIST[key]])
                                elif RAW_KEY_LIST[key] == 'ram_size':
                                    key_holder = RAW_KEY_HOLDER[key] + '(' + wdata['ram_size_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[RAW_KEY_LIST[key]])
                                elif RAW_KEY_LIST[key] == 'iops_value':
                                    details_dict[RAW_KEY_HOLDER[key]] = "{:,}".format(wdata[RAW_KEY_LIST[key]])

                                    iops_wl_type = wdata['io_block_size'] if 'io_block_size' in wdata else 'VSI'
                                    if iops_wl_type == 'VSI':
                                        details_dict['IO Block Size'] = "8K IO Block Size 70/30 Read/Write"
                                    else:
                                        details_dict['IO Block Size'] = "64K IO Block Size 100% Write"

                                elif RAW_KEY_LIST[key] == 'cpu_attribute':
                                    if wdata['cpu_attribute'] == 'cpu_clock':
                                        details_dict[RAW_KEY_HOLDER[key]] = 'Clock'
                                        details_dict['Clock (GHz)'] = "{:,}".format(wdata['cpu_clock'])
                                    else:
                                        details_dict[RAW_KEY_HOLDER[key]] = 'Cores'
                                        details_dict['vCPUs'] = "{:,}".format(wdata['vcpus'])
                                else:
                                    # details_dict[RAW_KEY_HOLDER[key]] = wdata[RAW_KEY_LIST[key]]
                                    details_dict[RAW_KEY_HOLDER[key]] = "{:,}".format(wdata[RAW_KEY_LIST[key]]) \
                                        if type(wdata[RAW_KEY_LIST[key]]) is int else wdata[RAW_KEY_LIST[key]]

                    elif wdata['wl_type'] == "DB":
                        for key in range(0, len(DB_KEY_LIST)):
                            if wdata.get(DB_KEY_LIST[key]):
                                if DB_KEY_LIST[key] == 'db_size':
                                    key_holder = DB_KEY_HOLDER[key] + '(' + wdata['db_size_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[DB_KEY_LIST[key]])
                                elif DB_KEY_LIST[key] == 'ram_per_db':
                                    key_holder = DB_KEY_HOLDER[key] + '(' + wdata['ram_per_db_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[DB_KEY_LIST[key]])
                                else:
                                    # details_dict[DB_KEY_HOLDER[key]] = wdata[DB_KEY_LIST[key]]
                                    details_dict[DB_KEY_HOLDER[key]] = "{:,}".format(wdata[DB_KEY_LIST[key]]) \
                                        if type(wdata[DB_KEY_LIST[key]]) is int else wdata[DB_KEY_LIST[key]]

                    elif wdata['wl_type'] == "ROBO" or wdata['wl_type'] == "ROBO_BACKUP":
                        for key in range(0, len(ROBO_KEY_LIST)):
                            if wdata.get(ROBO_KEY_LIST[key]):
                                if ROBO_KEY_LIST[key] == 'disk_per_vm':
                                    key_holder = ROBO_KEY_HOLDER[key] + '(' + wdata['disk_per_vm_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[ROBO_KEY_LIST[key]])
                                elif ROBO_KEY_LIST[key] == 'base_image_size':
                                    key_holder = ROBO_KEY_HOLDER[key] + '(' + wdata['base_image_size_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[ROBO_KEY_LIST[key]])
                                elif ROBO_KEY_LIST[key] == 'ram_per_vm':
                                    key_holder = ROBO_KEY_HOLDER[key] + '(' + wdata['ram_per_vm_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[ROBO_KEY_LIST[key]])
                                else:
                                    # details_dict[ROBO_KEY_HOLDER[key]] = wdata[ROBO_KEY_LIST[key]]
                                    details_dict[ROBO_KEY_HOLDER[key]] = "{:,}".format(wdata[ROBO_KEY_LIST[key]]) \
                                        if type(wdata[ROBO_KEY_LIST[key]]) is int else wdata[ROBO_KEY_LIST[key]]

                    elif wdata['wl_type'] == "ORACLE":
                        for key in range(0, len(ORACLE_KEY_LIST)):
                            if wdata.get(ORACLE_KEY_LIST[key]):
                                if ORACLE_KEY_LIST[key] == 'db_size':
                                    key_holder = ORACLE_KEY_HOLDER[key] + '(' + wdata['db_size_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[ORACLE_KEY_LIST[key]])
                                elif ORACLE_KEY_LIST[key] == 'ram_per_db':
                                    key_holder = ORACLE_KEY_HOLDER[key] + '(' + wdata['ram_per_db_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[ORACLE_KEY_LIST[key]])
                                else:
                                    # details_dict[ORACLE_KEY_HOLDER[key]] = wdata[ORACLE_KEY_LIST[key]]
                                    details_dict[ORACLE_KEY_HOLDER[key]] = "{:,}".format(wdata[ORACLE_KEY_LIST[key]]) \
                                        if type(wdata[ORACLE_KEY_LIST[key]]) is int else wdata[ORACLE_KEY_LIST[key]]

                    elif wdata['wl_type'] == "EXCHANGE":
                        for key in range(0, len(EXCHANGE_KEY_LIST)):
                            if EXCHANGE_KEY_LIST[key] in wdata:
                                # details_dict[EXCHANGE_KEY_HOLDER[key]] = wdata[EXCHANGE_KEY_LIST[key]]
                                details_dict[EXCHANGE_KEY_HOLDER[key]] = "{:,}".format(wdata[EXCHANGE_KEY_LIST[key]]) \
                                    if type(wdata[EXCHANGE_KEY_LIST[key]]) is int else wdata[EXCHANGE_KEY_LIST[key]]

                    elif wdata['wl_type'] == "VDI_INFRA":
                        for key in range(0, len(VDIINFRA_KEY_LIST)):
                            if VDIINFRA_KEY_LIST[key] in wdata:
                                details_dict[VDIINFRA_KEY_HOLDER[key]] = wdata[VDIINFRA_KEY_LIST[key]]

                        for vm, vm_detail in wdata['vm_details'].items():
                            for key in range(0, len(VMDETAILS_KEY_LIST)):
                                if VMDETAILS_KEY_LIST[key] in vm_detail:
                                    if VMDETAILS_KEY_LIST[key] == 'ram_per_vm':
                                        key_holder = VMDETAILS_KEY_HOLDER[key] + '(' + vm_detail['ram_per_vm_unit'] + ')'
                                        details_dict[vm + key_holder] = "{:,}".format(vm_detail[VMDETAILS_KEY_LIST[key]])
                                    else:
                                        details_dict[vm + VMDETAILS_KEY_HOLDER[key]] = vm_detail[VMDETAILS_KEY_LIST[key]]

                    elif wdata['wl_type'] == "EPIC":
                        for key in range(0, len(EPIC_KEY_LIST)):
                            if EPIC_KEY_LIST[key] in wdata:
                                details_dict[EPIC_KEY_HOLDER[key]] = wdata[EPIC_KEY_LIST[key]]

                    elif wdata['wl_type'] == "VEEAM":
                        for key in range(0, len(VEEAM_KEY_LIST)):
                            if VEEAM_KEY_LIST[key] in wdata:
                                if VEEAM_KEY_LIST[key] == 'hdd_size':
                                    key_holder = VEEAM_KEY_HOLDER[key] + '(' + wdata['hdd_size_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[VEEAM_KEY_LIST[key]])
                                else:
                                    details_dict[VEEAM_KEY_HOLDER[key]] = wdata[VEEAM_KEY_LIST[key]]

                    elif wdata['wl_type'] == "SPLUNK":
                        for key in range(0, len(SPLUNK_KEY_LIST)):
                            if SPLUNK_KEY_LIST[key] in wdata:
                                if SPLUNK_KEY_LIST[key] == 'daily_data_ingest':
                                    key_holder = SPLUNK_KEY_HOLDER[key] + '(' + wdata['daily_data_ingest_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[SPLUNK_KEY_LIST[key]])
                                elif SPLUNK_KEY_LIST[key] == 'max_vol_ind':
                                    key_holder = SPLUNK_KEY_HOLDER[key] + '(' + wdata['max_vol_ind_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[SPLUNK_KEY_LIST[key]])
                                else:
                                    details_dict[SPLUNK_KEY_HOLDER[key]] = wdata[SPLUNK_KEY_LIST[key]]

                    elif wdata['wl_type'] == "RDSH":
                        # skip RDSH_HOME workload, info included in main workload
                        if 'primary_wl_name' in wdata:
                            continue

                        for key in range(0, len(RDSH_KEY_LIST)):
                            if RDSH_KEY_LIST[key] in wdata:
                                if RDSH_KEY_LIST[key] == 'os_per_vm':
                                    key_holder = RDSH_KEY_HOLDER[key] + '(' + wdata['os_per_vm_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[RDSH_KEY_LIST[key]])
                                elif RDSH_KEY_LIST[key] == 'hdd_per_user':
                                    key_holder = RDSH_KEY_HOLDER[key] + '(' + wdata['hdd_per_user_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[RDSH_KEY_LIST[key]])
                                elif RDSH_KEY_LIST[key] == 'ram_per_vm':
                                    key_holder = RDSH_KEY_HOLDER[key] + '(' + wdata['ram_per_vm_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[RDSH_KEY_LIST[key]])
                                else:
                                    details_dict[RDSH_KEY_HOLDER[key]] = wdata[RDSH_KEY_LIST[key]]

                    elif wdata['wl_type'] == "CONTAINER":
                        for key in range(0, len(CONTAINER_KEY_LIST)):
                            if wdata.get(CONTAINER_KEY_LIST[key]):
                                if CONTAINER_KEY_LIST[key] == 'disk_per_container':
                                    key_holder = CONTAINER_KEY_HOLDER[key] + '(' + wdata[
                                        'disk_per_container_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[CONTAINER_KEY_LIST[key]])
                                elif CONTAINER_KEY_LIST[key] == 'base_image_size':
                                    key_holder = CONTAINER_KEY_HOLDER[key] + '(' + wdata['base_image_size_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[CONTAINER_KEY_LIST[key]])
                                elif CONTAINER_KEY_LIST[key] == 'ram_per_container':
                                    key_holder = CONTAINER_KEY_HOLDER[key] + '(' + wdata['ram_per_container_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[CONTAINER_KEY_LIST[key]])
                                else:
                                    details_dict[CONTAINER_KEY_HOLDER[key]] = "{:,}".format(
                                        wdata[CONTAINER_KEY_LIST[key]]) \
                                        if type(wdata[CONTAINER_KEY_LIST[key]]) is int else wdata[
                                        CONTAINER_KEY_LIST[key]]

                    elif wdata['wl_type'] == "AIML":
                        for key in range(0, len(AIML_KEY_LIST)):
                            if AIML_KEY_LIST[key] in wdata:
                                if AIML_KEY_LIST[key] == 'disk_per_ds':
                                    key_holder = AIML_KEY_HOLDER[key] + '(' + wdata['disk_per_ds_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[AIML_KEY_LIST[key]])
                                elif AIML_KEY_LIST[key] == 'enablestorage':
                                    key_holder = AIML_KEY_HOLDER[key]
                                    if wdata[AIML_KEY_LIST[key]]:
                                        details_dict[key_holder] = "True"
                                    else:
                                        details_dict[key_holder] = "False"
                                else:
                                    details_dict[AIML_KEY_HOLDER[key]] = "{:,}".format(wdata[AIML_KEY_LIST[key]]) \
                                        if type(wdata[AIML_KEY_LIST[key]]) is int else wdata[AIML_KEY_LIST[key]]

                    elif wdata['wl_type'] == "AWR_FILE":
                        for key in range(0, len(AWR_KEY_LIST)):
                            if AWR_KEY_LIST[key] in wdata:
                                if AWR_KEY_LIST[key] == 'db_size':
                                    key_holder = AWR_KEY_HOLDER[key] + '(' + wdata['db_size_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[AWR_KEY_LIST[key]])
                                elif AWR_KEY_LIST[key] == 'db_type':
                                    details_dict[AWR_KEY_HOLDER[key]] = "{:,}".format(wdata[AWR_KEY_LIST[key]]) \
                                        if type(wdata[AWR_KEY_LIST[key]]) is int else wdata[AWR_KEY_LIST[key]]

                                    if wdata[AWR_KEY_LIST[key]] == 'OLAP':
                                        details_dict['IOPS'] = "{:,}".format(wdata['avg_iops_per_db'])
                                    else:
                                        details_dict['IOPS'] = "{:,}".format(wdata['avg_mbps_per_db'])
                                elif AWR_KEY_LIST[key] == 'provisioned':
                                    key_holder = AWR_KEY_HOLDER[key]
                                    if wdata[AWR_KEY_LIST[key]]:
                                        details_dict[key_holder] = "Provisioned"
                                    else:
                                        details_dict[key_holder] = "Utilized"
                                else:
                                    details_dict[AWR_KEY_HOLDER[key]] = "{:,}".format(wdata[AWR_KEY_LIST[key]]) \
                                        if type(wdata[AWR_KEY_LIST[key]]) is int else wdata[AWR_KEY_LIST[key]]

                    elif wdata['wl_type'] == "ANTHOS":

                        for pod in wdata['pod_detail']:
                            for key in range(0, len(ANTHOS_POD_KEY_LIST)):
                                if ANTHOS_POD_KEY_LIST[key] in pod:
                                    if ANTHOS_POD_KEY_LIST[key] == 'pod_ram':
                                        key_holder = ANTHOS_POD_KEY_HOLDER[key] + '(' + pod['pod_ram_unit'] + ')'
                                        details_dict[pod['pod_name'] + '_' + key_holder] = \
                                            "{:,}".format(pod[ANTHOS_POD_KEY_LIST[key]])
                                    elif ANTHOS_POD_KEY_LIST[key] == 'pod_storage':
                                        key_holder = ANTHOS_POD_KEY_HOLDER[key] + '(' + pod['pod_storage_unit'] + ')'
                                        details_dict[pod['pod_name'] + '_' + key_holder] = \
                                            "{:,}".format(pod[ANTHOS_POD_KEY_LIST[key]])
                                    elif ANTHOS_POD_KEY_LIST[key] == 'worker_node_ram':
                                        key_holder = ANTHOS_POD_KEY_HOLDER[key] + '(' + pod['worker_node_ram_unit'] + ')'
                                        details_dict[pod['pod_name'] + '_' + key_holder] = \
                                            "{:,}".format(pod[ANTHOS_POD_KEY_LIST[key]])
                                    elif ANTHOS_POD_KEY_LIST[key] == 'worker_node_storage':
                                        key_holder = ANTHOS_POD_KEY_HOLDER[key] + '(' + pod['worker_node_storage_unit'] + ')'
                                        details_dict[pod['pod_name'] + '_' + key_holder] = \
                                            "{:,}".format(pod[ANTHOS_POD_KEY_LIST[key]])
                                    else:
                                        details_dict[pod['pod_name'] + '_' + ANTHOS_POD_KEY_HOLDER[key]] = \
                                            pod[ANTHOS_POD_KEY_LIST[key]]

                        for key in range(0, len(ANTHOS_KEY_LIST)):
                            if ANTHOS_KEY_LIST[key] in wdata:
                                if ANTHOS_KEY_LIST[key] == 'user_vm_ram':
                                    key_holder = ANTHOS_KEY_HOLDER[key] + '(' + wdata['user_vm_ram_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[ANTHOS_KEY_LIST[key]])
                                elif ANTHOS_KEY_LIST[key] == 'user_vm_storage':
                                    key_holder = ANTHOS_KEY_HOLDER[key] + '(' + wdata['user_vm_storage_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[ANTHOS_KEY_LIST[key]])
                                elif ANTHOS_KEY_LIST[key] == 'audit_log':
                                    key_holder = ANTHOS_KEY_HOLDER[key] + '(' + wdata['audit_log_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[ANTHOS_KEY_LIST[key]])
                                elif ANTHOS_KEY_LIST[key] == 'prometheous_storage':
                                    key_holder = ANTHOS_KEY_HOLDER[key] + '(' + wdata['prometheous_storage_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata[ANTHOS_KEY_LIST[key]])
                                else:
                                    details_dict[ANTHOS_KEY_HOLDER[key]] = wdata[ANTHOS_KEY_LIST[key]]

                            elif ANTHOS_KEY_LIST[key] in wdata['controller_panel']:
                                if ANTHOS_KEY_LIST[key] == 'controller_vm_ram':
                                    key_holder = ANTHOS_KEY_HOLDER[key] + '(' + \
                                                 wdata['controller_panel']['controller_vm_ram_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata['controller_panel'][ANTHOS_KEY_LIST[key]])
                                elif ANTHOS_KEY_LIST[key] == 'load_balancer_ram':
                                    key_holder = ANTHOS_KEY_HOLDER[key] + '(' + \
                                                 wdata['controller_panel']['load_balancer_ram_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata['controller_panel'][ANTHOS_KEY_LIST[key]])
                                else:
                                    details_dict[ANTHOS_KEY_HOLDER[key]] = wdata['controller_panel'][ANTHOS_KEY_LIST[key]]

                            elif ANTHOS_KEY_LIST[key] in wdata['anthos_master']:
                                if ANTHOS_KEY_LIST[key] == 'vm_ram':
                                    key_holder = ANTHOS_KEY_HOLDER[key] + '(' + \
                                                 wdata['anthos_master']['vm_ram_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata['anthos_master'][ANTHOS_KEY_LIST[key]])
                                elif ANTHOS_KEY_LIST[key] == 'vm_storage':
                                    key_holder = ANTHOS_KEY_HOLDER[key] + '(' + \
                                                 wdata['anthos_master']['vm_storage_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata['anthos_master'][ANTHOS_KEY_LIST[key]])
                                elif ANTHOS_KEY_LIST[key] == 'etcd_anthos_master':
                                    key_holder = ANTHOS_KEY_HOLDER[key] + '(' + \
                                                 wdata['anthos_master']['etcd_anthos_master_unit'] + ')'
                                    details_dict[key_holder] = "{:,}".format(wdata['anthos_master'][ANTHOS_KEY_LIST[key]])
                                else:
                                    details_dict[ANTHOS_KEY_HOLDER[key]] = wdata['anthos_master'][ANTHOS_KEY_LIST[key]]

                    else:
                        continue

                    req_data = str(json.dumps(details_dict))
                    output = req_data.replace(':', '-').replace(', ', ' | ').replace('{', '').replace('}', ''). \
                        replace("'", '').replace('"', '')

                    wl_cluster_per_row.append(output)

                    cluster_dict = OrderedDict()
                    # cluster_dict['Resilience Factor'] = wdata['replication_factor']
                    # cluster_dict['Fault Tolerance Nodes'] = wdata['fault_tolerance']

                    for key in range(0, len(CLUSTER_KEY_LIST)):
                        if wdata.get(CLUSTER_KEY_LIST[key]):
                            cluster_dict[CLUSTER_KEY_HOLDER[key]] = wdata[CLUSTER_KEY_LIST[key]]

                    cluster_placement = site_failure = ""
                    if len(cluster_data_index) == 1:
                        cluster_placement = ""

                    elif len(cluster_data_index) > 1 and 'remote_replication_enabled' in wdata and \
                            not wdata['remote_replication_enabled']:
                        cluster_placement = " Site A"

                    elif len(cluster_data_index) > 1 and 'remote_replication_enabled' in wdata and \
                            wdata['remote_replication_enabled'] and 'remote' in wdata:

                        if wdata['remote']:
                            cluster_placement = " Site B"
                        else:
                            cluster_placement = " Site A"

                        site_failure = wdata['replication_amt']

                    if cluster_placement != "":
                        cluster_dict['Placement'] = cluster_placement
                    if site_failure != "":
                        cluster_dict['Site failure protection(%)'] = str(site_failure)

                    req_data = str(json.dumps(cluster_dict))
                    output = req_data.replace(':', '-').replace(', ', '\n').replace('{', '').replace('}', ''). \
                        replace("'", '').replace('"', '')

                    wl_cluster_per_row.append(output)

                    row_name = "row" + str(row_index)
                    wl_details_table[row_name] = wl_cluster_per_row
                    row_index += 1
                break  # No need to add DR cluster workloads again

        return wl_details_table

    def create_wl_summary_page(self):

        cluster_list = self.scenario_data['workload_result'][0]['clusters']

        wl_details_table = self.get_wl_details(cluster_list)

        workload_list = self.scenario_data['workload_json']['wl_list']
        self.workload_note = False
        for wl in workload_list:
            if wl['wl_type'] == 'VSI' or wl['wl_type'] == 'DB':
                key = 'ram_per_vm' if wl['wl_type'] == 'VSI' else 'ram_per_db'
                if wl.get(key, 0) >= 128:
                    self.workload_note = True

        # Coordinates for table
        x = 0.32  # LEFT POSITION
        y = 0.4  # TOP POSITION
        wid = 10  # table Width
        hei = 0.9  # height
        slide_title = _("Workload Summary")
        col_width = [1.3, 0.9, 5.3, 1.8]

        rows_per_slide = 5
        anthoswl = any(wl_list_tmp['wl_type'] == 'ANTHOS' for wl_list_tmp in workload_list)
        if anthoswl:
            rows_per_slide = 3

        self.render_table_across_pages(slide_title, x, y, wid, hei, col_width, wl_details_table, rows_per_slide)

        self.workload_note = False

    def get_percluster_sizing_summary(self, site_info, cluster_data, cluster_index):

        percluster_sizing_summary_table = OrderedDict()
        row_index = 0

        no_of_workloads = 0
        wl_list = cluster_data['wl_list']
        for wl in wl_list:
            if wl['wl_type'] in ['VSI', 'DB', 'ORACLE']:
                if self.scenario_data['settings_json'][0]['dr_enabled']:
                    if (wl['remote'] and cluster_index == 1) or (not wl['remote'] and not cluster_index):
                        no_of_workloads += 1
                else:
                    no_of_workloads += 1
            else:
                no_of_workloads += 1

        # no_of_workloads = len(cluster_data['wl_list'])
        wls_details = ['No. of Workloads', str(no_of_workloads)]
        row_name = "row" + str(row_index)
        percluster_sizing_summary_table[row_name] = wls_details
        row_index += 1

        num_nodes = 0
        num_ft_nodes = 0
        power_consumption = 0
        node_info_list = cluster_data['node_info']
        for node_data in node_info_list:
            num_nodes += int(node_data['num_nodes'])
            if 'num_ft_nodes' in node_data:
                num_ft_nodes += int(node_data['num_ft_nodes'])

            power_consumption += int(node_data['power_consumption'])

        nodes = num_nodes - num_ft_nodes
        node_value = str(nodes) + " + " + str(num_ft_nodes) + " (FT)"
        nodes_details = ['Number of Nodes', node_value]
        row_name = "row" + str(row_index)
        percluster_sizing_summary_table[row_name] = nodes_details
        row_index += 1

        num_ru = cluster_data['rack_units']

        no_of_slots_per_chassis = 8
        total_slots_per_cluster = num_chassis_per_cluster * no_of_slots_per_chassis
        free_slots = total_slots_per_cluster - int(total_consumed_slots)

        if int(total_consumed_slots) and total_slots_per_cluster and free_slots:

            ru_details = ['Number of RU', str(num_ru) + "*"]
            slot_msg = '*includes ' + str(free_slots) + ' free chassis slots'

        else:

            ru_details = ['Number of RU', str(num_ru)]
            slot_msg = ''

        row_name = "row" + str(row_index)
        percluster_sizing_summary_table[row_name] = ru_details
        row_index += 1

        if site_info != "":
            placement_details = ['Cluster Placement', site_info]
            row_name = "row" + str(row_index)
            percluster_sizing_summary_table[row_name] = placement_details
            row_index += 1

        # capex = node_data['capex']
        # total_capex = "US $" + str(capex)
        # capex_details = ['Total Capex', total_capex]
        # row_name = "row" + str(row_index)
        # percluster_sizing_summary_table[row_name] = capex_details
        # row_index += 1

        '''
        power_consum = str(power_consumption/1000) + " KW" 
        power_details = ['Power Cooling', power_consum]
        row_name = "row" + str(row_index)
        percluster_sizing_summary_table[row_name] = power_details
        '''

        return percluster_sizing_summary_table, slot_msg

    @staticmethod
    def convert_raw_cpucorestocpuclock(util_list):

        base_cpu = SpecIntData.objects.get(is_base_model=True)

        for util_data in util_list:
            tag_name = util_data['tag_name']
            if tag_name == 'Physical Cores' or tag_name == 'CPU':
                util_data['node_val'] = util_data['node_val'] * float(base_cpu.speed)
                util_data['units'] = 'GHz'
                util_data['workload_val'] = util_data['workload_val'] * float(base_cpu.speed)
                util_data['best_practice'] = util_data['best_practice'] * float(base_cpu.speed)

    @staticmethod
    def epiccluster_addup(util_list, total_cluster_count):

        for util_data in util_list:
            tag_name = util_data['tag_name']
            if tag_name == 'Physical Cores' or tag_name == 'CPU' or tag_name == 'RAM':
                util_data['node_val'] = util_data['node_val'] * total_cluster_count
                util_data['workload_val'] = util_data['workload_val'] * total_cluster_count
                util_data['best_practice'] = util_data['best_practice'] * total_cluster_count
            if tag_name == 'Storage Capacity':
                util_data['node_val'] = util_data['node_val'] * total_cluster_count
                util_data['node_val_binarybyte'] = util_data['node_val_binarybyte'] * total_cluster_count
                util_data['usable'] = util_data['usable'] * total_cluster_count
                util_data['usable_binarybyte'] = util_data['usable_binarybyte'] * total_cluster_count
                util_data['workload_val'] = util_data['workload_val'] * total_cluster_count
                util_data['best_practice'] = util_data['best_practice'] * total_cluster_count
                util_data['best_practice_binarybyte'] = util_data['best_practice_binarybyte'] * total_cluster_count

    def create_per_cluster_sys_util_page(self, cluster_name, util_list, rep_factor, wl_list, ft, total_cluster_count=1):

        prs = self.prs
        slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BOX_LAYOUT_INDEX])
        slide_title = "%s %s" %(cluster_name, _('Overall Sizing Summary'))

        cpu = slide.placeholders[10]
        memory = slide.placeholders[11]
        disk = slide.placeholders[12]
        iops = slide.placeholders[13]

        cpu.text = "Physical Cores"
        memory.text = "RAM"
        disk.text = "Storage Capacity"
        iops.text = "Storage IOPS"

        # Added not to render Storage IOPS for RAW workloads.
        wl_type = ""
        wls_cores = 0
        wls_clock = 0
        for wldata in wl_list:
            wl_type = wldata['wl_type']
            if wl_type == 'RAW' or wl_type == 'RAW_FILE':
                if wldata['cpu_attribute'] == 'vcpus':
                    wls_cores = 1
                elif wldata['cpu_attribute'] == 'cpu_clock':
                    wls_clock = 1

        raw_wls_mixed = CORES  # Output will be in Cores
        if not wls_cores and wls_clock == 1:
            raw_wls_mixed = CLOCK  # Output will be in Clock

        if raw_wls_mixed == CLOCK:
            self.convert_raw_cpucorestocpuclock(util_list)

        if total_cluster_count > 1:
            self.epiccluster_addup(util_list, total_cluster_count)

        self.render_util(slide, slide_title, util_list, rep_factor, wl_type, ft, wl_list)

    # merge cells horizontally
    @staticmethod
    def merge_cells_horizontally(table, row_idx, start_col_idx, end_col_idx):

        col_count = end_col_idx - start_col_idx + 1
        row_cells = [c for c in table.rows[row_idx].cells][start_col_idx:end_col_idx]
        row_cells[0]._tc.set('gridSpan', str(col_count))

        for c in row_cells[1:]:
            c._tc.set('hMerge', '1')

    # merge cells vertically
    @staticmethod
    def merge_cells_vertically(table, start_row_idx, end_row_idx, col_idx):
        row_count = end_row_idx - start_row_idx + 1
        column_cells = [r.cells[col_idx] for r in table.rows][start_row_idx:]

        column_cells[0]._tc.set('rowSpan', str(row_count))
        for c in column_cells[1:]:
            c._tc.set('vMerge', '1')


    def get_replication_type(self, wldata, cluster_index):

        if wldata.get('replication_type', 'any') == 'any':
            return 'any'
        else:
            remote_flag = wldata['remote']
            if cluster_index and not remote_flag:
                return 'replicated'
            elif not cluster_index and remote_flag:
                return 'replicated'
            else:
                return 'normal'

    def get_workload_name(self, wldata, cluster_index):

        wl_name = wldata['wl_name']

        dr_workload_types = [HyperConstants.VSI, HyperConstants.DB, HyperConstants.OLTP, HyperConstants.OLAP,
                             HyperConstants.ORACLE, HyperConstants.OOLTP, HyperConstants.OOLAP, HyperConstants.AWR_FILE]
        wl_type = wldata['wl_type']
        replication_type = self.get_replication_type(wldata, cluster_index)

        if wl_type in dr_workload_types and replication_type == 'replicated':
            wl_name += ' [Replicated]'

        return wl_name


    def create_anthos_calc_table(self, workload):

        wl_name = workload['wl_name']

        total_user_cluster = 0
        cpu_req = 0
        ram_req = 0
        hdd_req = 0

        for pod in workload['pod_detail']:
            # CPU
            total_core_pod = float((pod['pod_cpu'] / 1000) * pod['pod_quantity'])
            cpu_overhead = 0.05 + (0.01 * pod['worker_node_cpu'])
            cpu_worker_node = float((total_core_pod + cpu_overhead) / pod['worker_node_cpu'])

            # RAM
            total_ram_pod = float(WL.unit_conversion(pod['pod_ram'], pod['pod_ram_unit'], 'GiB') * pod['pod_quantity'])
            ram_overhead = 0.25 + (0.05 * pod['worker_node_ram'])
            ram_worker_node = float((total_ram_pod + ram_overhead) / WL.unit_conversion(pod['worker_node_ram'],
                                                                                          pod['worker_node_ram_unit'],
                                                                                          'GiB'))

            max_worker_node_per_cluster = 100
            max_user_cluster = ceil(cpu_worker_node / max_worker_node_per_cluster)
            total_user_cluster += max_user_cluster

            pod_ha = 3 if pod['pod_ha'] else 1

            user_master_cpu = max_user_cluster * pod_ha * workload['user_vm_cpu']
            user_master_ram = max_user_cluster * pod_ha * WL.unit_conversion(workload['user_vm_ram'],
                                                                             workload['user_vm_ram_unit'], 'GiB')
            user_master_storage = max_user_cluster * pod_ha * WL.unit_conversion(workload['user_vm_storage'],
                                                                                 workload['user_vm_storage_unit'])
            audit_log_storage = max_user_cluster * pod_ha * WL.unit_conversion(workload['audit_log'],
                                                                               workload['audit_log_unit'])
            etcd_event_data = max_user_cluster * pod_ha * WL.unit_conversion(workload['etcd_event'],
                                                                             workload['etcd_event_unit'])

            max_worker_node = round(max(cpu_worker_node, ram_worker_node))

            worker_node_cpu_req = max_worker_node * pod['worker_node_cpu']
            worker_node_ram_req = max_worker_node * WL.unit_conversion(pod['worker_node_ram'],
                                                                         pod['worker_node_ram_unit'], 'GiB')

            #HDD
            total_storage_pod = float(
                WL.unit_conversion(pod['pod_storage'], pod['pod_storage_unit']) * pod['pod_quantity'])
            worker_node_storage_req = total_storage_pod + (
                        max_user_cluster * WL.unit_conversion(workload['gc_ops_overhead_user_vm'],
                                                              workload['gc_ops_overhead_user_vm_unit']))

            if pod['prometheous_on']:
                prometheous = max_user_cluster * WL.unit_conversion(workload['prometheous_storage'],
                                                                    workload['prometheous_storage_unit'])
            else:
                prometheous = 0


            cpu_req += user_master_cpu + worker_node_cpu_req
            ram_req += user_master_ram + worker_node_ram_req
            hdd_req += user_master_storage + worker_node_storage_req + audit_log_storage + etcd_event_data + prometheous


        # step-1
        # cpu_str = "Pod_CPU_Req += Master_CPU_Req + Worker_Node_CPU_Req"
        # ram_str = "Pod_RAM_Req += Master_RAM_Req + Worker_Node_RAM_Req "
        # hdd_str = "Pod_Storage_Req += Master_Storage_Req  + Worker_Storage + Audit_Log_Storage + ETCD_Event_Data + Prometheous"

        cpu_str = "*Pod_CPU_Req"
        ram_str = "*Pod_RAM_Req"
        hdd_str = "*Pod_Storage_Req"

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        max_cluster = 5
        load_balancer = 2
        controller_vm = 4

        total_admin_cluster = ceil(total_user_cluster / max_cluster)

        # load balancer - CPU and RAM
        load_balancer_cpu = load_balancer * workload['controller_panel']['load_balancer_cpu']
        load_balancer_ram = float(
            load_balancer * WL.unit_conversion(workload['controller_panel']['load_balancer_ram'],
                                               workload['controller_panel']['load_balancer_ram_unit'], 'GiB'))

        # controller vm - CPU and RAM
        controller_vm_cpu = controller_vm * workload["controller_panel"]['controller_vm_cpu']
        controller_vm_ram = float(
            controller_vm * WL.unit_conversion(workload["controller_panel"]['controller_vm_ram'],
                                               workload["controller_panel"]['controller_vm_ram_unit'], 'GiB'))

        # calculate control panel - CPU and RAM
        control_panel_cpu = load_balancer_cpu + controller_vm_cpu
        control_panel_ram = load_balancer_ram + controller_vm_ram

        #  calculate Anthos Master VMs - unit conversion is not taken here
        anthos_master_cpu = float(total_admin_cluster * 3 * workload["anthos_master"]['vm_cpu'])
        anthos_master_ram = float(total_admin_cluster * 3 * WL.unit_conversion(workload["anthos_master"]['vm_ram'],
                                                                               workload["anthos_master"][
                                                                                     'vm_ram_unit'], 'GiB'))
        anthos_master_storage = float(
            total_admin_cluster * 3 * WL.unit_conversion(workload["anthos_master"]['vm_storage'],
                                                         workload["anthos_master"]['vm_storage_unit']))
        anthos_master_gc_ops = float(
            total_admin_cluster * WL.unit_conversion(workload["anthos_master"]['gc_ops_overhead'],
                                                     workload["anthos_master"]['gc_ops_overhead_unit']))
        anthos_master_etcd = float(
            total_admin_cluster * WL.unit_conversion(workload["anthos_master"]['etcd_anthos_master'],
                                                     workload["anthos_master"]['etcd_anthos_master_unit']))

        # cpu_str = "Controller_Anthos_CPU = control_panel_cpu + anthos_master_cpu"
        # ram_str = "Controller_Anthos_RAM = control_panel_ram + anthos_master_ram"
        # hdd_str = "Controller_Anthos_HDD = anthos_master_storage + anthos_master_gc_ops +  anthos_master_etcd"

        cpu_str = "*Controller_Anthos_CPU"
        ram_str = "*Controller_Anthos_RAM"
        hdd_str = "*Controller_Anthos_HDD"

        controller_anthos_cpu = control_panel_cpu + anthos_master_cpu
        controller_anthos_ram = control_panel_ram + anthos_master_ram
        controller_anthos_hdd = anthos_master_storage + anthos_master_gc_ops + anthos_master_etcd

        self.create_workload_calc_table(wl_name, cpu_str, controller_anthos_cpu,
                                        ram_str, controller_anthos_ram, hdd_str, controller_anthos_hdd)

        # step-3
        cpu_str = "Pod_CPU_Req + Controller_Anthos_CPU"
        cpu_capsum = cpu_req + controller_anthos_cpu

        ram_str = "Pod_RAM_Req + Controller_Anthos_RAM"
        ram_capsum = ram_req + controller_anthos_ram

        hdd_str = "Pod_CPU_Storage + Controller_Anthos_HDD"
        hdd_capsum = hdd_req + controller_anthos_hdd

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_aiml_calc_table(self, workload):

        wl_name = workload['wl_name']

        hdd_size = WL.unit_conversion(workload[HyperConstants.HDD_PER_DS],
                                      workload[HyperConstants.HDD_PER_DS_UNIT])

        # step-1
        cpu_str = "vCPUs per Data Scientist / vCPUs per Core"
        cpu_req = (float(workload[HyperConstants.VCPUS_PER_DS]) /
                     float(workload[HyperConstants.VCPUS_PER_CORE])) * WL.normalise_cpu()

        ram_str = "RAM per Data Scientist"
        ram_size = WL.unit_conversion(workload[HyperConstants.RAM_PER_DS],
                                            workload[HyperConstants.RAM_PER_DS_UNIT], 'GiB')

        ram_req = ram_size

        hdd_str = "Storage Capacity per Data Scientist"
        hdd_req = hdd_size

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        req_text = "# of Data Scientist"
        number_of_ds = workload[HyperConstants.NUM_DATA_SCIENTISTS]

        self.create_workload_calc_table(wl_name, req_text, number_of_ds,
                                        req_text, number_of_ds, req_text, number_of_ds)

        # step-3
        cpu_str = "(vCPUs per Data Scientist / vCPUs per Core) * # of Data Scientist"
        cpu_capsum = cpu_req * number_of_ds

        ram_str = "RAM per Data Scientist * # of Data Scientist"
        ram_capsum = ram_req * number_of_ds

        hdd_str = "Storage Capacity per Data Scientist * # of Data Scientist"
        hdd_capsum = hdd_req * number_of_ds

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_container_calc_table(self, workload):

        wl_name = workload['wl_name']

        ram_size = WL.unit_conversion(workload[HyperConstants.RAM_PER_CONTAINER],
                                        workload[HyperConstants.RAM_PER_CONTAINER_UNIT], 'GiB')

        hdd_size = WL.unit_conversion(workload[HyperConstants.HDD_PER_CONTAINER],
                                        workload[HyperConstants.HDD_PER_CONTAINER_UNIT])

        image_size = WL.unit_conversion(workload[HyperConstants.BASE_IMG_SIZE],
                                          workload[HyperConstants.BASE_IMG_SIZE_UNIT])

        # step-1
        cpu_str = "vCPUs per Container / vCPUs per Core"
        cpu_req = (float(workload[HyperConstants.VCPUS_PER_CONTAINER]) /
                     float(workload[HyperConstants.VCPUS_PER_CORE])) * WL.normalise_cpu()

        ram_str = "RAM per Container"
        ram_req = ram_size

        hdd_str = "Storage Capacity per Container = User Data Size + OS Image Size"
        hdd_req = hdd_size + image_size

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        req_text = "# of Containers"
        number_of_container = workload[HyperConstants.NUM_CONTAINERS]

        self.create_workload_calc_table(wl_name, req_text, number_of_container,
                                        req_text, number_of_container, req_text, number_of_container)

        # step-3
        cpu_str = "(vCPUs per Container / vCPUs per Core) * # of Containers"
        cpu_capsum = cpu_req * number_of_container

        ram_str = "RAM per Container * # of Containers"
        ram_capsum = ram_req * number_of_container

        hdd_str = "Storage Capacity per Container * # of Containers"
        hdd_capsum = hdd_req * number_of_container

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_veeam_calc_table(self, workload):

        wl_name = workload['wl_name']

        hdd_size = WL.unit_conversion(workload['hdd_size'],
                                        workload['hdd_size_unit'])

        # step-1
        cpu_str = "--"
        cpu_req = "--"

        ram_str = "--"
        ram_req = "--"

        hdd_str = "Storage requirement"
        hdd_req = hdd_size

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        req_text = "--"
        self.create_workload_calc_table(wl_name, req_text, req_text,
                                        req_text, req_text, req_text, req_text)

        # step-3
        cpu_str = "Workload Requirement = N/A"
        cpu_capsum = 0

        ram_str = "Workload Requirement = N/A"
        ram_capsum = 0

        hdd_str = "Total Storage Requirement"
        hdd_capsum = hdd_req

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_raw_calc_table(self, workload):

        wl_name = workload['wl_name']

        safety_overhead = workload.get(HyperConstants.RAW_OVERHEAD_PERCENTAGE, 0) / 100.0

        ram_size = WL.unit_conversion(workload[BaseConstants.RAM_SIZE],
                                      workload[BaseConstants.RAM_SIZE_UNIT], 'GiB')

        hdd_size = WL.unit_conversion(workload[BaseConstants.HDD_SIZE],
                                      workload[BaseConstants.HDD_SIZE_UNIT])

        input_cpu = workload.get(HyperConstants.CPU_MODEL, None)

        try:
            if not input_cpu:
                cpu_model = SpecIntData.objects.get(is_base_model=True)
            else:
                cpu_model = SpecIntData.objects.get(model=input_cpu)
        except ObjectDoesNotExist:
            raise Exception("Input CPU model doesn't exist")

        # step-1
        if workload[HyperConstants.CPU_ATTRIBUTE] == HyperConstants.CPU_CLOCK:
            cpu_str = "(CPU Clock / CPU Speed) / vCPUs per Core * (1 + Future Growth/100.0)"
            vcpus = workload[HyperConstants.CPU_CLOCK] / float(cpu_model.speed)
            cpu_cores = vcpus / float(workload[HyperConstants.VCPUS_PER_CORE])
        else:
            cpu_str = "(vCPUs / vCPUs per Core) * (1 + Future Growth/100.0)"
            cpu_cores = workload[HyperConstants.VCPUS] / float(workload[HyperConstants.VCPUS_PER_CORE])

        normalized_cores = cpu_cores * WL.normalise_cpu(cpu_model)

        cpu_req = ceil(normalized_cores * (1 + safety_overhead))

        ram_str = "(RAM Size / RAM Overprovisioning ratio) * (1 + Future Growth/100.0)"
        ram_req = (ram_size / float(workload.get(HyperConstants.RAM_OPRATIO, 1))) * (1 + safety_overhead)

        hdd_str = "HDD Size * (1 + Future Growth/100.0)"
        hdd_req = hdd_size * (1 + safety_overhead)

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        req_text = "--"
        self.create_workload_calc_table(wl_name, req_text, req_text,
                                        req_text, req_text, req_text, req_text)

        # step-3
        cpu_str = "CPU requirement"
        cpu_capsum = cpu_req

        ram_str = "RAM requirement"
        ram_capsum = ram_req

        hdd_str = "Storage requirement"
        hdd_capsum = hdd_req

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_exchange_calc_table(self, workload):

        wl_name = workload['wl_name']

        safety_overhead = workload.get(HyperConstants.RAW_OVERHEAD_PERCENTAGE, 0) / 100.0

        # step-1
        cpu_str = "(vCPUs / vCPUs per Core) * (1 + Future Growth/100.0)"
        cpu_req = workload[HyperConstants.VCPUS] / float(workload[HyperConstants.VCPUS_PER_CORE]) * \
                    WL.normalise_cpu() * (1 + safety_overhead)

        ram_str = "RAM Size * (1 + Future Growth/100.0)"
        ram_size = WL.unit_conversion(workload[BaseConstants.RAM_SIZE],
                                      workload[BaseConstants.RAM_SIZE_UNIT], 'GiB')

        ram_req = ceil(ram_size * (1 + safety_overhead))

        hdd_str = "HDD Size * (1 + Future Growth/100.0)"
        hdd_req = workload[BaseConstants.HDD_SIZE] * (1 + safety_overhead)

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        req_text = "--"
        self.create_workload_calc_table(wl_name, req_text, req_text,
                                        req_text, req_text, req_text, req_text)

        # step-3
        cpu_str = "CPU requirement"
        cpu_capsum = cpu_req

        ram_str = "RAM requirement"
        ram_capsum = ram_req

        hdd_str = "Storage requirement"
        hdd_capsum = hdd_req

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_vsi_calc_table(self, workload, cluster_index):

        dr_enabled = self.scenario_data['settings_json'][0]['dr_enabled']

        wl_name = self.get_workload_name(workload, cluster_index)

        hdd_size = WL.unit_conversion(workload[HyperConstants.HDD_PER_VM],
                                      workload[HyperConstants.HDD_PER_VM_UNIT])

        image_size = WL.unit_conversion(workload[HyperConstants.VM_BASE_IMG_SIZE],
                                          workload[HyperConstants.VM_BASE_IMG_SIZE_UNIT])

        replication_type = 'any'
        if dr_enabled:
            replication_type = self.get_replication_type(workload, cluster_index)

        # step-1
        if replication_type == 'any' or replication_type == 'normal':

            cpu_str = "vCPUs per VM (in terms of Intel Platinum 8164) / vCPUs per Core"
            cpu_req = (float(workload[HyperConstants.VCPUS_PER_VM]) /
                     float(workload[HyperConstants.VCPUS_PER_CORE])) * WL.normalise_cpu()

            ram_str = "RAM per VM"
            ram_size = WL.unit_conversion(workload[HyperConstants.RAM_PER_VM],
                                          workload[HyperConstants.RAM_PER_VM_UNIT], 'GiB')
            ram_req = ram_size

        elif replication_type == 'replicated':

            cpu_str = "--"
            cpu_req = 0

            ram_str = "--"
            ram_req = 0

        hdd_str = "Storage Capacity per VM = User Data + OS Image + (0.02 * User Data * Number of Snapshots)"
        hdd_req = hdd_size + image_size + (0.02 * hdd_size * workload['snapshots'])

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        if replication_type == 'any' or replication_type == 'normal':

            req_text = "# of VMs"
            number_of_vms = workload[HyperConstants.NUM_VM]

        elif replication_type == 'replicated':

            req_text = "Number of Replicated VMs = Number of VMs * Replication Amount / 100"
            number_of_vms = int(ceil(workload['num_vms'] * workload['replication_amt'] / 100.0))

        self.create_workload_calc_table(wl_name, req_text, number_of_vms,
                                        req_text, number_of_vms, req_text, number_of_vms)

        # step-3
        if replication_type == 'any' or replication_type == 'normal':

            cpu_str = "vCPUs per VM * # of VMs"
            cpu_capsum = cpu_req * number_of_vms

            ram_str = "RAM per VM * # of VMs"
            ram_capsum = ram_req * number_of_vms

            hdd_str = "Storage Capacity per VM * # of VMs"
            hdd_capsum = hdd_req * number_of_vms

        elif replication_type == 'replicated':

            cpu_str = "Workload Requirement = N/A"
            cpu_capsum = 0

            ram_str = "Workload Requirement = N/A"
            ram_capsum = 0

            hdd_str = "Storage Capacity per VM * # of Replicated VMs"
            num_replica_vms = int(ceil(workload['num_vms'] * workload['replication_amt'] / 100.0))
            hdd_capsum = hdd_req * num_replica_vms

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_splunk_calc_table(self, workload):

        wl_name = workload['wl_name']

        cpu_req = list()
        cpu_capsum = 0
        ram_req = list()
        ram_capsum = 0

        # step-1
        for _, vm_detail in workload['vm_details'].items():
            cpu_req.append(
                float(vm_detail[HyperConstants.VCPUS_PER_VM]) / workload[HyperConstants.VCPUS_PER_CORE] *
                        vm_detail[HyperConstants.NUM_VM] )

            cpu_capsum += float(vm_detail[HyperConstants.VCPUS_PER_VM]) / workload[HyperConstants.VCPUS_PER_CORE] * \
                          vm_detail[HyperConstants.NUM_VM]

            ram_per_vm = WL.unit_conversion(vm_detail[HyperConstants.RAM_PER_VM],
                                                vm_detail[HyperConstants.RAM_PER_VM_UNIT], 'GiB')

            ram_req.append(round(ram_per_vm * vm_detail[HyperConstants.NUM_VM], 2))

            ram_capsum += ram_per_vm * vm_detail[HyperConstants.NUM_VM]


        cpu_str = "CPU per VM Type += (vCPU per VM / vCPU per Core) * # of VMs"
        ram_str = "RAM per VM Type += RAM per VM *  # of VMs"
        hdd_str = "HDD Size = Daily data ingest"
        hdd_req = WL.unit_conversion(workload['daily_data_ingest'],
                                     workload['daily_data_ingest_unit'])

        cpu_req_with_plus = ' + '.join([str(v) for v in cpu_req])
        ram_req_with_plus = ' + '.join([str(v) for v in ram_req])
        self.create_workload_calc_table(wl_name, cpu_str, cpu_req_with_plus,
                                        ram_str, ram_req_with_plus, hdd_str, hdd_req)

        # step - 2
        req_text = "--"
        if workload['acc_type'] == 'hx+splunk':

            hdd_text = "Total Days for HX+Splunk = Splunk Level Replication * (Hot Tier + Cold Tier + Frozen Tier)"
            total_days = workload.get('app_rf', 2) * (workload['storage_acc']['hot'] + workload['storage_acc']['cold'] +
                                                      workload['storage_acc']['frozen'])

        elif workload['acc_type'] == 'hx+splunk_smartstore':

            hdd_text = "Total Days for HX+Splunk Smartstore = (Hot Tier * 2) + Warm Tier"
            total_days = (workload['storage_acc']['hot'] * 2) + workload['storage_acc']['warm']


        self.create_workload_calc_table(wl_name, req_text, req_text,
                                        req_text, req_text, hdd_text, total_days)

        # step-3
        cpu_str = "Total CPU"
        ram_str = "Total RAM"
        hdd_str = "HDD Size * Total Days"
        hdd_capsum = hdd_req * total_days

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_db_oracle_calc_table(self, workload, cluster_index):

        dr_enabled = self.scenario_data['settings_json'][0]['dr_enabled']
        wl_name = self.get_workload_name(workload, cluster_index)

        db_size = WL.unit_conversion(workload[HyperConstants.DB_SIZE],
                                     workload[HyperConstants.DB_SIZE_UNIT])

        replication_type = 'any'
        if dr_enabled:
            replication_type = self.get_replication_type(workload, cluster_index)

        cpu_normalise = 1
        if workload[BaseConstants.WL_TYPE] == HyperConstants.AWR_FILE:

            input_cpu = workload.get(HyperConstants.CPU_MODEL, None)
            cpu_model = None

            try:
                if not input_cpu:
                    cpu_model = SpecIntData.objects.get(is_base_model=True)
                else:
                    cpu_model = SpecIntData.objects.get(model=input_cpu)
            except Exception as e:
                pass

            cpu_normalise = WL.normalise_cpu(cpu_model)

        # step-1
        if replication_type == 'any' or replication_type == 'normal':
            cpu_str = "vCPUs per DB / vCPUs per Core"
            cpu_req = (float(workload[HyperConstants.VCPUS_PER_DB]) / float(workload[HyperConstants.VCPUS_PER_CORE])) \
                      * cpu_normalise

            ram_str = "RAM per DB"
            ram_size = WL.unit_conversion(workload[HyperConstants.RAM_PER_DB],
                                     workload[HyperConstants.RAM_PER_DB_UNIT], 'GiB')
            ram_req = ram_size
        elif replication_type == 'replicated':
            cpu_str = "--"
            cpu_req = 0

            ram_str = "--"
            ram_req = 0

        hdd_str = "Total Database Size = DB Size + (1 + Meta Data / 100.0)"
        hdd_req = db_size * (1 + (workload[HyperConstants.META_DATA] / 100.0))

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        if replication_type == 'any' or replication_type == 'normal':
            req_text = "# of DBs"
            number_of_dbs = workload['num_db_instances']

        elif replication_type == 'replicated':

            req_text = "Number of Replicated DBs = Number of DBs * Replication Amount / 100"
            number_of_dbs = int(ceil(workload['num_db_instances'] * workload['replication_amt'] / 100.0))

        self.create_workload_calc_table(wl_name, req_text, number_of_dbs,
                                        req_text, number_of_dbs, req_text, number_of_dbs)

            # step-3
        if replication_type == 'any' or replication_type == 'normal':
            cpu_str = "CPU per DB * # of DBs"
            cpu_capsum = cpu_req * number_of_dbs

            ram_str = "RAM per DB * # of DBs"
            ram_capsum = ram_req * number_of_dbs

            hdd_str = "DB Size per DB * # of DBs"
            hdd_capsum = hdd_req * number_of_dbs
        elif replication_type == 'replicated':
            cpu_str = "Workload Requirement = N/A"
            cpu_capsum = 0

            ram_str = "Workload Requirement = N/A"
            ram_capsum = 0

            hdd_str = "DB Size per DB * # of Replicated DBs"
            hdd_capsum = hdd_req * number_of_dbs

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_rdsh_home_calc_table(self, workload):

        wl_name = workload['wl_name']

        # step-1
        profile = workload['profile']
        base_cpu = SpecIntData.objects.get(is_base_model=True)
        cpu_str = "CPU per VM based on Profile"
        cpu_req = float(HyperConstants.HOME_CONFIG[profile]['cpu']) * WL.normalise_cpu() * base_cpu.speed

        ram_str = "RAM per VM based on Profile"
        ram_req = HyperConstants.HOME_CONFIG[profile]['ram']

        hdd_str = "--"
        hdd_req = "--"

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        req_text = "# of VMs"
        number_of_vms = workload['number_of_vms']

        self.create_workload_calc_table(wl_name, req_text, number_of_vms,
                                        req_text, number_of_vms, hdd_str, hdd_req)

        # step-3
        cpu_str = "CPU per VM * # of VMs"
        cpu_capsum = cpu_req * number_of_vms

        ram_str = "RAM per VM * # of VMs"
        ram_capsum = ram_req * number_of_vms

        hdd_str = "Workload Requirement = N/A"
        hdd_capsum = 0

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_rdsh_calc_table(self, workload):

        wl_name = workload['wl_name']

        ram_size = WL.unit_conversion(workload[HyperConstants.RAM_PER_VM],
                                      workload[HyperConstants.RAM_PER_VM_UNIT], 'GiB')

        hdd_size = WL.unit_conversion(workload[HyperConstants.HDD_PER_USER],
                                      workload[HyperConstants.HDD_PER_USER_UNIT])

        image_size = WL.unit_conversion(workload[HyperConstants.OS_PER_VM],
                                        workload[HyperConstants.OS_PER_VM_UNIT])

        # step-1

        base_cpu = SpecIntData.objects.get(is_base_model=True)
        cpu_str = "Clock per Session / 1000"
        cpu_req = (workload[HyperConstants.CLOCK_PER_SESSION] / 1000.0) * WL.normalise_cpu()

        ram_str = "RAM per VM"
        ram_req = ram_size

        hdd_str = "User Data = HDD Size * Total Users"
        hdd_req = hdd_size * workload['total_users']

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        cpu_req_text = "Total Users"
        total_users = workload['total_users']

        ram_req_text = "# of VMs = Total Users / Sessions per VM"
        num_vms = int(ceil(workload['total_users'] / workload['sessions_per_vm']))

        hdd_req_text = "OS Data = Image Size * (Total Users / Sessions per VM)"
        os_data = image_size * num_vms

        self.create_workload_calc_table(wl_name, cpu_req_text, total_users,
                                        ram_req_text, num_vms, hdd_req_text, os_data)

        # step-3
        cpu_str = "(Clock per Session / 1000 ) * Total Users"
        cpu_capsum = cpu_req * total_users

        ram_str = "RAM per VM * # of VMs"
        ram_capsum = ram_req * num_vms

        hdd_str = "User Data + OS Data"
        hdd_capsum = hdd_req + os_data

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_vdi_infra_calc_table(self, workload):

        wl_name = workload['wl_name']

        cpu_req = list()
        cpu_capsum = 0
        ram_req = list()
        ram_capsum = 0
        hdd_req = list()
        hdd_capsum = 0

        base_cpu = SpecIntData.objects.get(is_base_model=True)
        # step-1

        for vm_detail in workload['vm_details'].values():
            cpu_req.append(
                float(vm_detail[HyperConstants.VCPUS_PER_VM]) / workload[HyperConstants.VCPUS_PER_CORE] * \
                WL.normalise_cpu() * vm_detail[HyperConstants.NUM_VM] * float(base_cpu.speed) )

            cpu_capsum += float(vm_detail[HyperConstants.VCPUS_PER_VM]) / workload[HyperConstants.VCPUS_PER_CORE] * \
                WL.normalise_cpu() * vm_detail[HyperConstants.NUM_VM] * float(base_cpu.speed)

            ram_per_vm = WL.unit_conversion(vm_detail[HyperConstants.RAM_PER_VM],
                                            vm_detail[HyperConstants.RAM_PER_VM_UNIT], 'GiB')

            ram_req.append(round(ram_per_vm / float(workload[HyperConstants.RAM_OPRATIO])
                           * vm_detail[HyperConstants.NUM_VM], 2))

            ram_capsum += (ram_per_vm / float(workload[HyperConstants.RAM_OPRATIO])
                           * vm_detail[HyperConstants.NUM_VM])

            hdd_size = WL.unit_conversion(vm_detail[HyperConstants.HDD_PER_VM],
                                          vm_detail[HyperConstants.HDD_PER_VM_UNIT])

            hdd_req.append(hdd_size * vm_detail[HyperConstants.NUM_VM])
            hdd_capsum += hdd_size * vm_detail[HyperConstants.NUM_VM]


        cpu_str = "Total CPU += (vCPU per VM / vCPU per Core) * # of VMs"
        ram_str = "Total RAM += (RAM per VM / RAM Overprovisioning Ratio) *  # of VMs"
        hdd_str = "Total HDD += HDD Size * # of VMs"

        cpu_req_with_plus = ' + '.join([str(v) for v in cpu_req])
        ram_req_with_plus = ' + '.join([str(v) for v in ram_req])
        hdd_req_with_plus = ' + '.join([str(v) for v in hdd_req])
        self.create_workload_calc_table(wl_name, cpu_str, cpu_req_with_plus,
                                        ram_str, ram_req_with_plus, hdd_str, hdd_req_with_plus)

        # step - 2
        req_text = "--"
        self.create_workload_calc_table(wl_name, req_text, req_text,
                                        req_text, req_text, req_text, req_text)

        # step-3
        cpu_str = "Total CPU for Infra type"
        ram_str = "Total RAM for Infra type"
        hdd_str = "Total HDD for Infra type"

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_epic_calc_table(self, workload, total_cluster_count=1):

        wl_name = workload['wl_name']

        user_concurrency = workload['concurrent_user_pcnt'] / 100.0
        num_clusters = workload['num_clusters']

        # step-1
        input_cpu = SpecIntData.objects.get(model=workload['cpu'])
        total_clock_per_host = 2 * input_cpu.speed * input_cpu.cores
        cpu_str = "Total Clock per Host / Users per Host"
        cpu_req = total_clock_per_host * WL.normalise_cpu(input_cpu) / workload['users_per_host']

        ram_str = "RAM per Guest * Guest per Host"
        ram_per_guest = WL.unit_conversion(workload['ram_per_guest'],
                                             workload['ram_per_guest_unit'], 'GiB')
        ram_req = ram_per_guest * workload['guests_per_host']

        hdd_per_guest = 70
        hdd_str = "HDD Size = Guest per Host * HDD per Guest\nNote: HDD per Guest is fixed at 70 GB per host"
        hdd_req = workload['guests_per_host'] * hdd_per_guest

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        req_text = "Users = Concurrency * Total Users / # of Clusters"
        num_users = round(user_concurrency * workload['total_users'] / float(num_clusters))
        num_users *= total_cluster_count

        req_text_2 = "Expected hosts per Cluster = Concurrency * Expected hosts / # of Clusters"
        ex_hosts_per_cluster = round(user_concurrency * workload['expected_hosts'] / float(num_clusters))
        ex_hosts_per_cluster *= total_cluster_count

        self.create_workload_calc_table(wl_name, req_text, num_users,
                                        req_text_2, ex_hosts_per_cluster, req_text_2, ex_hosts_per_cluster)

        # step-3
        cpu_str = "Total Clock per Host / Users per Host * # of Users"
        cpu_capsum = cpu_req * num_users

        ram_str = "RAM per Guest * Guest per Host * Expected hosts per Cluster"
        ram_capsum = ram_req * ex_hosts_per_cluster

        hdd_str = "HDD Size * # Expected hosts per Cluster"
        hdd_capsum = hdd_req * ex_hosts_per_cluster

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_vdi_home_calc_table(self, workload):

        wl_name = workload['wl_name']

        # step-1
        profile = workload['profile']
        base_cpu = SpecIntData.objects.get(is_base_model=True)
        cpu_str = "CPU per Desktop based on Profile * (Concurrency Ratio / 100)"
        cpu_req = float(HyperConstants.HOME_CONFIG[profile]['cpu']) * float(base_cpu.speed)* \
                  WL.normalise_cpu() * (workload['concurrency']/100.0)

        ram_str = "RAM per Desktop based on Profile* (Concurrency Ratio / 100)"
        ram_req = HyperConstants.HOME_CONFIG[profile]['ram'] * (workload['concurrency'] / 100.0)

        hdd_str = "--"
        hdd_req = 0

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        req_text = "# of VMs"
        num_desktops = workload['number_of_vms']

        self.create_workload_calc_table(wl_name, req_text, num_desktops,
                                        req_text, num_desktops, req_text, num_desktops)

        # step-3
        cpu_str = "CPU per Desktop * # of VMs"
        cpu_capsum = cpu_req * num_desktops

        ram_str = "RAM per Desktop * # of VMs"
        ram_capsum = ram_req * num_desktops

        hdd_str = "Workload Requirement = N/A"
        hdd_capsum = 0

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_workload_calc_table(self, wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req):

        wl_calc_per_row = list()
        wl_calc_per_row.append(wl_name)
        wl_calc_per_row.append(_(cpu_str))
        cpu_req_value = cpu_req if type(cpu_req) != float else round(cpu_req, 2)
        wl_calc_per_row.append(str(cpu_req_value))

        wl_calc_per_row.append(_(ram_str))
        ram_req_value = ram_req if type(ram_req) != float else round(ram_req, 2)
        wl_calc_per_row.append(str(ram_req_value))

        wl_calc_per_row.append(_(hdd_str))
        hdd_req_value = hdd_req if type(hdd_req) != float else round(hdd_req, 2)
        wl_calc_per_row.append(str(hdd_req_value))

        row_name = "row" + str(self.row_index)
        self.wl_calc_table[row_name] = wl_calc_per_row
        self.row_index += 1

    def create_vdi_calc_table(self, workload):

        wl_name = workload['wl_name']

        ram_size = WL.unit_conversion(workload[HyperConstants.RAM_PER_DT],
                                        workload[HyperConstants.RAM_PER_DT_UNIT], 'GiB')

        hdd_size = WL.unit_conversion(workload[HyperConstants.HDD_PER_DT],
                                        workload[HyperConstants.HDD_PER_DT_UNIT])
        image_size = WL.unit_conversion(workload[HyperConstants.GOLD_IMG_SIZE],
                                          workload[HyperConstants.GOLD_IMG_SIZE_UNIT])
        snapshot_factor = workload[HyperConstants.DT_SNAPSHOTS] * 0.02 + 1
        inflight_size = workload[HyperConstants.INFLIGHT_DATA]

        # step-1
        cpu_str = "(Total Clock per Desktop / 1000) * (Concurrency Ratio / 100)"
        cpu_req = (workload['clock_per_desktop'] / 1000.0) * (workload['concurrency']/100.0)

        ram_str = "RAM per Desktop * (Concurrency Ratio / 100)"
        ram_req = ram_size * (workload['concurrency'] / 100.0)

        if workload[HyperConstants.DT_PROV_TYPE] in [HyperConstants.VIEW_FULL]:
            hdd_str = "Storage Capacity per Desktop = ( HDD Size + OS Image Size ) * # of Snapshots"
            hdd_req = (hdd_size + image_size) * snapshot_factor
        else:
            hdd_str = "Storage Capacity per Desktop = ( HDD Size + Inflight Size ) * # of Snapshots"
            hdd_req = (hdd_size + inflight_size) * snapshot_factor

        self.create_workload_calc_table(wl_name, cpu_str, cpu_req,
                                        ram_str, ram_req, hdd_str, hdd_req)

        # step - 2
        req_text = "# of Desktops"
        num_desktops = workload['num_desktops']

        self.create_workload_calc_table(wl_name, req_text, num_desktops,
                                        req_text, num_desktops, req_text, num_desktops)

        # step-3
        cpu_str = "CPU per Desktop * # of Desktops"
        cpu_capsum = cpu_req * num_desktops

        ram_str = "RAM per Desktop * # of Desktops"
        ram_capsum = ram_req * num_desktops

        if workload[HyperConstants.DT_PROV_TYPE] in [HyperConstants.VIEW_FULL]:
            hdd_str = "Storage Capacity per Desktop * # of Desktops"
            hdd_capsum = hdd_req * num_desktops
        else:
            hdd_str = "Storage Capacity per Desktop * # of Desktops + OS Image Size"
            hdd_capsum = hdd_req * num_desktops + image_size

        self.create_workload_calc_table(wl_name, cpu_str, cpu_capsum,
                                        ram_str, ram_capsum, hdd_str, hdd_capsum)

        self.wl_capsum['CPU'] += cpu_capsum
        self.wl_capsum['RAM'] += ram_capsum
        self.wl_capsum['HDD'] += hdd_capsum

        return

    def create_wl_calc_page(self, cluster_name, wl_list, cluster_index=1):

        if "Workload Calculation" not in self.slides:
            return

        workload_list = copy.deepcopy(wl_list)

        self.wl_calc_table = OrderedDict()
        self.row_index = 0

        for cap in ['CPU', 'RAM', 'HDD']:
            self.wl_capsum[cap] = 0

        cluster_wl_type = wl_list[0]['wl_type']

        # HEADER FORMATION - The Compute, Memory and Disk are merged to one column.
        if cluster_wl_type in [HyperConstants.VDI, HyperConstants.VDI_HOME, HyperConstants.VDI_INFRA,
                       HyperConstants.RDSH, HyperConstants.RDSH_HOME]:
            self.wl_calc_table['header'] = [_('Workload Name'), _('Compute (GHz)'), _('Compute (GHz)'),
                                        _('Memory (GiB)'), _('Memory (GiB)'), _('Disk (GB)'), _('Disk (GB)')]
        else:
            self.wl_calc_table['header'] = [_('Workload Name'), _('Compute (Cores)'), _('Compute (Cores)'),
                                            _('Memory (GiB)'), _('Memory (GiB)'), _('Disk (GB)'), _('Disk (GB)')]


        for wl in workload_list:
            wl_type = wl['wl_type']

            if wl_type == HyperConstants.VDI:
                if 'primary_wl_name' in wl:
                    self.create_vdi_home_calc_table(wl)
                else:
                    self.create_vdi_calc_table(wl)

            elif wl_type == HyperConstants.EPIC:
                total_cluster_count = cluster_index
                self.create_epic_calc_table(wl, total_cluster_count)

            elif wl_type == HyperConstants.VDI_INFRA:
                self.create_vdi_infra_calc_table(wl)

            elif wl_type == HyperConstants.RDSH:
                if 'primary_wl_name' in wl:
                    self.create_rdsh_home_calc_table(wl)
                else:
                    self.create_rdsh_calc_table(wl)

            elif wl_type == HyperConstants.DB or wl_type == HyperConstants.ORACLE or wl_type == HyperConstants.AWR_FILE:
                self.create_db_oracle_calc_table(wl, cluster_index)

            elif wl_type == HyperConstants.SPLUNK:
                self.create_splunk_calc_table(wl)

            elif wl_type == HyperConstants.VSI or wl_type == HyperConstants.ROBO:
                self.create_vsi_calc_table(wl, cluster_index)

            elif wl_type == HyperConstants.EXCHANGE:
                self.create_exchange_calc_table(wl)

            elif wl_type == HyperConstants.RAW or wl_type == HyperConstants.RAW_FILE:
                self.create_raw_calc_table(wl)

            elif wl_type == HyperConstants.VEEAM:
                self.create_veeam_calc_table(wl)

            elif wl_type == HyperConstants.CONTAINER:
                self.create_container_calc_table(wl)

            elif wl_type == HyperConstants.AIML:
                self.create_aiml_calc_table(wl)

            elif wl_type == HyperConstants.ANTHOS:
                self.create_anthos_calc_table(wl)

        if cluster_wl_type in [HyperConstants.VDI, HyperConstants.VDI_INFRA, HyperConstants.VDI_HOME, HyperConstants.RDSH,
                           HyperConstants.RDSH_HOME, HyperConstants.EPIC, HyperConstants.ANTHOS]:
            op_ratio = 1
            op_ratio_note = ""
            if cluster_wl_type in [HyperConstants.ANTHOS]:
                op_ratio_note = "*Workload Calculation formula/details are provided in next slide"

        else:
            ops = Counter(wl[HyperConstants.VCPUS_PER_CORE] for wl in wl_list)
            op_ratio = ops.most_common(1)[0][0]

            total_cpu_req = self.wl_capsum['CPU'] * op_ratio
            op_ratio_note = _("Note: The Total CPU workload requirement is multiple of most common over provisioning ratio " \
                            "of all the workloads in this cluster. This case the common OP ratio is") + " " + str(op_ratio) + \
                            "\n    Ex: " + str(self.wl_capsum['CPU']) + " * " + str(op_ratio) + " = " + str(round(total_cpu_req, 2))

        if len(self.wl_calc_table) > 1:
            # Coordinates for table
            x = 0.32    # LEFT POSITION
            y = 0.4            # TOP POSITION
            wid = 10    # table Width
            hei = 0.9       # height
            slide_title = "%s %s" %(cluster_name, _('Total Resources Required By Workloads'))
            col_width = [1.3, 1.3, 1.3, 1.3, 1.3, 1.3, 1.3]

            rows_per_slide = 9
            self.wl_capsum['CPU'] *= op_ratio
            self.render_calc_table_across_pages(slide_title, x, y, wid, hei, col_width,
                                                self.wl_calc_table, self.wl_capsum, rows_per_slide, op_ratio_note)

        if cluster_wl_type in [HyperConstants.ANTHOS]:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[PPTReport.SLIDEMASTER_ANTHOS_LAYOUT_INDEX])
            shapes = slide.shapes
            shapes.title.text = 'Anthos Sizing Workload Calculations'

    def get_node_calc_details(self, cluster_name, cluster_data, total_cluster_count=1):

        if "Node Calculation" not in self.slides:
            return

        total_raw_capacity_desc = defaultdict(str)
        total_raw_capacity_pre_spec = defaultdict(float)
        total_raw_capacity = defaultdict(float)

        node_overhead_desc = defaultdict(str)
        node_overhead = defaultdict(float)
        op_ratio = defaultdict(float)
        op_ratio_cap = defaultdict(float)

        fault_tolerance_capacity = defaultdict(float)
        usable_capacity = defaultdict(float)
        node_reserve = defaultdict(str)
        reserve_capacity = defaultdict(float)

        cpu_normalization = "NA"
        fault_tolerance = 1
        CAPLIST = ['CPU', 'RAM', 'HDD']

        node_info_list = cluster_data['node_info']
        threshold_factor = self.scenario_data['settings_json'][0]['threshold']
        replication_factor = float(cluster_data['settings']['replication_factor'])
        wl_type = cluster_data['wl_list'][0]['internal_type']

        wl_list = cluster_data['wl_list']
        wls_cores = 0
        wls_clock = 0

        for wldata in wl_list:
            wl_type = wldata['internal_type']
            if wl_type == 'RAW' or wl_type == 'RAW_FILE':
                if wldata['cpu_attribute'] == 'vcpus':
                    wls_cores = 1
                elif wldata['cpu_attribute'] == 'cpu_clock':
                    wls_clock = 1

        raw_wls_mixed = CORES  # Output will be in Cores
        if not wls_cores and wls_clock == 1:
            raw_wls_mixed = CLOCK  # Output will be in Clock

        # effective_storage_savings = round(cluster_data['Utilization'][2]['ratio'], 1)
        scaling_factor = 1

        # There is no performance data for CPU and RAM for these workloads
        if wl_type == "VEEAM":
            CAPLIST = ['HDD']
            total_raw_capacity_desc['CPU'] = "--"
            total_raw_capacity_desc['RAM'] = "--"

        for node_data in node_info_list:

            model_details = node_data['model_details']

            overprovisioning_ratio = node_data['model_details']['sizing_calculation']
            node_subtype = model_details['subtype']
            node_disk_cage = model_details['disk_cage']

            for cap in CAPLIST:
                overhead = node_data['model_details']['sizing_calculation']['node_overhead']

                threshold_key = cap
                if cap == 'HDD':
                    if node_subtype in [HyperConstants.ALL_FLASH, HyperConstants.AF_ROBO_NODE,
                                                      HyperConstants.ALLNVME_NODE, HyperConstants.ALLNVME_NODE_8TB,
                                                      HyperConstants.ALL_FLASH_7_6TB, HyperConstants.AF_ROBO_TWO_NODE,
                                                      HyperConstants.ROBO_AF_240]:
                        threshold_key = HyperConstants.ALL_FLASH_HDD

                    if node_disk_cage == 'LFF':
                        threshold_key = HyperConstants.LFF_HDD

                reservation = self.get_threshold_value(wl_type, threshold_factor, threshold_key)

                if cap == 'CPU':

                    if wl_type != 'VDI' and raw_wls_mixed == CORES and wl_type != 'RDSH' and wl_type != 'VDI_INFRA':
                        total_raw_capacity_desc[cap] = "# (HX + Compute) Nodes * # Sockets * # of Cores"
                        total_raw_capacity_pre_spec[cap] += \
                            model_details['cpu_socket_count'] * model_details['cores_per_cpu_prespeclnt'] * \
                            node_data['num_nodes']

                        raw_capacity_per_node = \
                            model_details['cpu_socket_count'] * model_details['cores_per_cpu_prespeclnt'] * \
                            model_details['speclnt']

                        raw_capacity = raw_capacity_per_node * node_data['num_nodes']

                        total_raw_capacity[cap] += raw_capacity

                        if model_details['subtype'] != 'compute' and model_details['subtype'] != 'aiml':
                            cpu_normalization = model_details['cpu_part'] + ': ' + \
                                                str(round(model_details['speclnt'], 2))
                            node_overhead_desc[cap] = str(overhead[cap]) + " pCPUs / node"
                            op_ratio[cap] = overprovisioning_ratio.get('cpu_opratio', 1)
                            node_reserve[cap] = reservation

                            fault_tolerance_nodes = node_data['num_ft_nodes']
                            fault_tolerance = fault_tolerance_nodes

                        else:
                            fault_tolerance_nodes = 0

                        overhead_cap = raw_capacity_per_node - overhead[cap]
                        node_overhead[cap] += overhead_cap * node_data['num_nodes']

                        op_ratio_overhead = overhead_cap * op_ratio[cap]
                        op_ratio_cap[cap] += op_ratio_overhead * node_data['num_nodes']

                        fault_tolerance_capacity[cap] += op_ratio_overhead * fault_tolerance_nodes

                        usable_cap = op_ratio_overhead * (node_data['num_nodes'] - fault_tolerance_nodes)

                        usable_capacity[cap] += usable_cap

                        reserve_capacity[cap] += usable_cap * (100 - reservation) / 100.0

                    elif raw_wls_mixed == CLOCK:
                        total_raw_capacity_desc[cap] = "# (HX + Compute) Nodes * # Sockets * # of Cores * Frequency"
                        total_raw_capacity_pre_spec[cap] += \
                            model_details['cpu_socket_count'] * model_details['cores_per_cpu_prespeclnt'] * \
                            node_data['num_nodes'] * float(model_details['base_frequency'])

                        raw_capacity_per_node = \
                            model_details['cpu_socket_count'] * model_details['cores_per_cpu_prespeclnt'] * \
                            model_details['speclnt'] * float(model_details['base_frequency'])

                        raw_capacity = raw_capacity_per_node * node_data['num_nodes']

                        total_raw_capacity[cap] += raw_capacity

                        if model_details['subtype'] != 'compute':
                            cpu_normalization = model_details['cpu_part'] + ': ' + \
                                                str(round(model_details['speclnt'], 2))
                            node_overhead_desc[cap] = str(overhead[cap]) + " pCPUs / node"
                            op_ratio[cap] = overprovisioning_ratio['cpu_opratio']
                            node_reserve[cap] = reservation
                            fault_tolerance_nodes = node_data['num_ft_nodes']
                            fault_tolerance = fault_tolerance_nodes
                        else:
                            fault_tolerance_nodes = 0

                        overhead_cap = \
                            raw_capacity_per_node - (overhead[cap] * float(model_details['base_frequency']))
                        node_overhead[cap] += overhead_cap * node_data['num_nodes']

                        op_ratio_overhead = overhead_cap * op_ratio[cap]
                        op_ratio_cap[cap] += op_ratio_overhead * node_data['num_nodes']

                        fault_tolerance_capacity[cap] += op_ratio_overhead * fault_tolerance_nodes

                        usable_cap = op_ratio_overhead * (node_data['num_nodes'] - fault_tolerance_nodes)

                        usable_capacity[cap] += usable_cap

                        reserve_capacity[cap] += usable_cap * (100 - reservation) / 100.0

                    else:
                        total_raw_capacity_desc[cap] = "# (HX + Compute) Nodes * # Sockets * # Cores"
                        total_raw_capacity_pre_spec[cap] += \
                            model_details['cpu_socket_count'] * model_details['cores_per_cpu_prespeclnt'] * \
                            float(model_details['frequency']) * node_data['num_nodes']

                        raw_capacity_per_node = \
                            model_details['cpu_socket_count'] * model_details['cores_per_cpu_prespeclnt'] * \
                            float(model_details['base_frequency']) * model_details['speclnt']

                        raw_capacity = raw_capacity_per_node * node_data['num_nodes']

                        total_raw_capacity[cap] += raw_capacity

                        if model_details['subtype'] != 'compute':
                            cpu_normalization = \
                                model_details['cpu_part'] + ': ' + \
                                str(round(raw_capacity / float(total_raw_capacity_pre_spec[cap]), 2))

                            node_overhead_desc[cap] = str(overhead[cap]) + " pCPUs / node"
                            op_ratio[cap] = overprovisioning_ratio['cpu_opratio']
                            node_reserve[cap] = reservation
                            fault_tolerance_nodes = node_data['num_ft_nodes']
                            fault_tolerance = fault_tolerance_nodes
                        else:
                            fault_tolerance_nodes = 0

                        overhead_cap = (raw_capacity_per_node -
                                        (overhead[cap] * float(model_details['base_frequency'])))

                        node_overhead[cap] += overhead_cap * node_data['num_nodes']

                        op_ratio_overhead = overhead_cap * op_ratio[cap]
                        op_ratio_cap[cap] += op_ratio_overhead * node_data['num_nodes']

                        fault_tolerance_capacity[cap] += op_ratio_overhead * fault_tolerance_nodes

                        usable_cap = op_ratio_overhead * (node_data['num_nodes'] - fault_tolerance_nodes)

                        usable_capacity[cap] += usable_cap

                        reserve_capacity[cap] += usable_cap * (100 - reservation) / 100.0

                elif cap == 'RAM':
                    total_raw_capacity_desc[cap] = "# (HX + Compute) Nodes * # DIMMs * DIMM Size"
                    raw_capacity_per_node = model_details['ram_slots'] * model_details['ram_size']
                    raw_capacity = raw_capacity_per_node * node_data['num_nodes']

                    total_raw_capacity[cap] += raw_capacity

                    if model_details['subtype'] != 'compute' and model_details['subtype'] != 'aiml':
                        node_overhead_desc[cap] = str(overhead[cap])
                        op_ratio[cap] = overprovisioning_ratio['ram_opratio']
                        node_reserve[cap] = reservation
                        fault_tolerance_nodes = node_data['num_ft_nodes']
                        fault_tolerance = fault_tolerance_nodes
                    else:
                        fault_tolerance_nodes = 0

                    overhead_cap = raw_capacity_per_node - overhead[cap]
                    node_overhead[cap] += overhead_cap * node_data['num_nodes']

                    op_ratio_overhead = overhead_cap * op_ratio[cap]
                    op_ratio_cap[cap] += op_ratio_overhead * node_data['num_nodes']

                    fault_tolerance_capacity[cap] += op_ratio_overhead * fault_tolerance_nodes

                    usable_cap = op_ratio_overhead * (node_data['num_nodes'] - fault_tolerance_nodes)

                    usable_capacity[cap] += usable_cap

                    reserve_capacity[cap] += usable_cap * (100 - reservation) / 100.0

                elif cap == 'HDD' and model_details['subtype'] != 'compute' and model_details['subtype'] != 'aiml':

                    total_raw_capacity_desc[cap] = "# HX Nodes * # Disk Slots * Disk Size"

                    hdd_size = model_details['hdd_size'] / (1.0 if model_details['HDD_GB_TB'] == 'TB' else 1000.0)

                    raw_capacity_per_node = model_details['hdd_slots'] * hdd_size
                    raw_capacity = raw_capacity_per_node * node_data['num_nodes']

                    total_raw_capacity[cap] += raw_capacity

                    if model_details['subtype'] != 'compute' and model_details['subtype'] != 'aiml':
                        node_overhead_desc[cap] = "8%"
                        op_ratio[cap] = overprovisioning_ratio['hdd_opratio']
                        node_reserve[cap] = reservation

                    hdd_overhead = 8 * raw_capacity_per_node / 100.0
                    overhead_cap = raw_capacity_per_node - hdd_overhead
                    node_overhead[cap] += overhead_cap * node_data['num_nodes']

                    hdd_capacity_after_rf = (overhead_cap / replication_factor)
                    op_ratio_cap[cap] += (hdd_capacity_after_rf - overhead_cap) * node_data['num_nodes']

                    usable_cap = \
                        (raw_capacity_per_node - hdd_overhead) / replication_factor * node_data['num_nodes']
                    usable_capacity[cap] += usable_cap

                    reserve_capacity[cap] += usable_cap * (100 - reservation) / 100.0
                    scaling_factor = model_details['sizing_calculation']['scaling_factor']

        # Since this is lot customized table need to add here only, rather than calling general function.
        prs = self.prs
        slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
        shapes = slide.shapes
        shapes.title.text = cluster_name + ' ' + _('Total Available Cluster Resources')

        rows = 11  # rows = 14
        cols = 8
        left = Inches(0.05)
        top = Inches(0.39)
        width = Inches(9.5)
        height = Inches(0.5)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        for colindex in range(0, cols):
            table.columns[colindex].width = Inches(0.9)
            if colindex == 0:  # last column
                table.columns[colindex].width = Inches(1.0)
            if colindex == 6:  # last column
                table.columns[colindex].width = Inches(1.5)
            if colindex == 7:  # last column
                table.columns[colindex].width = Inches(2.9)  # 2.4

        if wl_type == 'VDI' or raw_wls_mixed == CLOCK or \
                wl_type == 'RDSH' or wl_type == 'VDI_INFRA':
            cpu_text = 'Compute (GHz)'
        else:
            cpu_text = 'Compute (Cores)'

        # HEADER FORMATION
        row_index = 0
        table.cell(row_index, 0).text = ''
        table.cell(row_index, 1).text = _(cpu_text)
        table.cell(row_index, 3).text = _('Memory (GiB)')
        table.cell(row_index, 5).text = _('Disk (TB/TiB)')
        table.cell(row_index, 7).text = _('Notes')

        for colindex in range(0, cols):
            table.cell(row_index, colindex).text_frame.paragraphs[0].font.size = Pt(9)
            table.cell(row_index, colindex).text_frame.paragraphs[0].font.bold = True
            table.cell(row_index, colindex).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            table.cell(row_index, colindex).fill.solid()
            table.cell(row_index, colindex).fill.fore_color.rgb = RGBColor(0xB2, 0xB2, 0xB2)

        # merge cells horizontally
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=1, end_col_idx=2)
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=3, end_col_idx=4)
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=5, end_col_idx=6)

        for rowindex in range(1, rows):
            for colindex in range(0, cols):
                table.cell(rowindex, colindex).text_frame.paragraphs[0].font.size = Pt(7)
                table.cell(rowindex, colindex).text_frame.paragraphs[0].font.italic = True
                table.cell(rowindex, colindex).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                table.cell(rowindex, colindex).vertical_anchor = MSO_ANCHOR.MIDDLE

                # if colindex == 7:  # last column
                #     table.cell(rowindex, colindex).text_frame.paragraphs[0].font.size = Pt(6)

        if total_cluster_count > 1:
            total_raw_capacity_pre_spec['CPU'] = total_raw_capacity_pre_spec['CPU'] * total_cluster_count

            for cap in CAPLIST:
                total_raw_capacity[cap] = total_raw_capacity[cap] * total_cluster_count
                node_overhead[cap] = node_overhead[cap] * total_cluster_count
                op_ratio_cap[cap] = op_ratio_cap[cap] * total_cluster_count
                fault_tolerance_capacity[cap] = fault_tolerance_capacity[cap] * total_cluster_count
                usable_capacity[cap] = usable_capacity[cap] * total_cluster_count
                reserve_capacity[cap] = reserve_capacity[cap] * total_cluster_count

        # STARTING EACH ROWS ONE BY ONE..
        row_index += 1
        table.cell(row_index, 0).text = ''
        table.cell(row_index, 0).text_frame.paragraphs[0].font.bold = True
        table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        table.cell(row_index, 1).text = _(total_raw_capacity_desc['CPU'])
        table.cell(row_index, 2).text = "{:,}".format(int(ceil(total_raw_capacity_pre_spec['CPU'])))
        table.cell(row_index, 3).text = _(total_raw_capacity_desc['RAM'])
        table.cell(row_index, 4).text = "{:,}".format(round(total_raw_capacity['RAM'], 1))
        table.cell(row_index, 5).text = _(total_raw_capacity_desc['HDD'])
        binary_unit = "{:,}".format(round(total_raw_capacity['HDD'] * HyperConstants.TB_TO_TIB_CONVERSION, 1))
        table.cell(row_index, 6).text = \
            "{:,}".format(round(total_raw_capacity['HDD'], 1)) + ' TB (' + binary_unit + ' TiB)'

        row_index += 1
        table.cell(row_index, 0).text = _('CPU Normalization')
        table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        if wl_type == "VEEAM":
            table.cell(row_index, 1).text = 'Not Applicable'
        else:
            table.cell(row_index, 1).text = _('Factor for ') + cpu_normalization

        table.cell(row_index, 2).text = "{:,}".format(int(ceil(total_raw_capacity['CPU'])))
        tooltip = \
            "Hx Sizer uses an Intel Platinum 8164 CPU as a baseline. Selected CPUs are normalized using the CPU SPEC " \
            "ratings against the baseline CPU"
        table.cell(row_index, 7).text = _(tooltip)
        table.cell(row_index, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        table.cell(row_index, 7).text_frame.paragraphs[0].font.italic = False

        row_index += 1
        table.cell(row_index, 0).text = _('Total Physical Resources')
        table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        table.cell(row_index, 0).text_frame.paragraphs[0].font.bold = True
        for colindex in range(0, cols):
            table.cell(row_index, colindex).fill.solid()
            table.cell(row_index, colindex).fill.fore_color.rgb = RGBColor(0xB0, 0xE0, 0xE6)
            table.cell(row_index, colindex).text_frame.paragraphs[0].font.italic = False
        table.cell(row_index, 1).text = "{:,}".format(int(ceil(total_raw_capacity['CPU'])))
        table.cell(row_index, 3).text = "{:,}".format(round(total_raw_capacity['RAM'], 1))
        binary_unit = "{:,}".format(round(total_raw_capacity['HDD'] * HyperConstants.TB_TO_TIB_CONVERSION, 1))
        table.cell(row_index, 5).text = \
            "{:,}".format(round(total_raw_capacity['HDD'], 1)) + ' TB (' + binary_unit + ' TiB)'
        tooltip = "Total physical resources available for sizing purposes"
        table.cell(row_index, 7).text = _(tooltip)
        table.cell(row_index, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        # merge cells horizontally
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=1, end_col_idx=2)
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=3, end_col_idx=4)
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=5, end_col_idx=6)

        # row_index += 1  # Empty Row Needed??
        row_index += 1
        table.cell(row_index, 0).text = _("Controller VM & Metadata Reserves")
        table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        table.cell(row_index, 1).text = node_overhead_desc['CPU']
        table.cell(row_index, 2).text = "{:,}".format(round(node_overhead['CPU'] - total_raw_capacity['CPU'], 1))
        table.cell(row_index, 3).text = node_overhead_desc['RAM']
        table.cell(row_index, 4).text = "{:,}".format(round(node_overhead['RAM'] - total_raw_capacity['RAM'], 1))
        table.cell(row_index, 5).text = node_overhead_desc['HDD']
        hdd_cap = node_overhead['HDD'] - total_raw_capacity['HDD']
        binary_unit = "{:,}".format(round(hdd_cap * HyperConstants.TB_TO_TIB_CONVERSION, 1))
        table.cell(row_index, 6).text = "{:,}".format(round(hdd_cap, 1)) + ' TB (' + binary_unit + ' TiB)'
        tooltip = "These are resourced used to run the HXDP stack that enables data services such as dedupe, " \
                  "compression, instantaneous snapshots, clones, etc."
        table.cell(row_index, 7).text = _(tooltip)
        table.cell(row_index, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        table.cell(row_index, 7).text_frame.paragraphs[0].font.italic = False

        row_index += 1
        table.cell(row_index, 0).text = _("Overprovisioning/Replication Factor")
        table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        table.cell(row_index, 1).text = str(op_ratio['CPU'])
        table.cell(row_index, 2).text = "{:,}".format(round((op_ratio_cap['CPU']), 1))
        table.cell(row_index, 3).text = str(op_ratio['RAM'])
        table.cell(row_index, 4).text = "{:,}".format(round((op_ratio_cap['RAM']), 1))
        table.cell(row_index, 5).text = "RF" + str(replication_factor)
        binary_unit = "{:,}".format(round((op_ratio_cap['HDD']) * HyperConstants.TB_TO_TIB_CONVERSION, 1))
        table.cell(row_index, 6).text = "{:,}".format(round(op_ratio_cap['HDD'], 1)) + \
                                        ' TB (' + binary_unit + ' TiB)'
        tooltip = _("Overprovisioning includes CPU and RAM Overprovisioning Ratios.")
        tooltip += "\nRF" + str(replication_factor) + _(" ensures that data is available even after the cluster has " + \
                   str(int(replication_factor - 1)) + " uncorrelated and simultaneous failures")
        table.cell(row_index, 7).text = _(tooltip)
        table.cell(row_index, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        table.cell(row_index, 7).text_frame.paragraphs[0].font.italic = False

        row_index += 1
        table.cell(row_index, 0).text = _("Performance headroom (N+X)")
        table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        col = 1
        for cap in ['CPU', 'RAM']:
            table.cell(row_index, col).text = 'N+' + str(fault_tolerance)
            col += 1
            table.cell(row_index, col).text = "{:,}".format(round(-1 * fault_tolerance_capacity[cap], 1))
            col += 1

        tooltip = "N+X: This ensures there is enough performance capability to sustain applications after X failures"
        table.cell(row_index, 7).text = _(tooltip)
        table.cell(row_index, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        table.cell(row_index, 7).text_frame.paragraphs[0].font.italic = False

        row_index += 1
        table.cell(row_index, 0).text = _('Total Usable')
        table.cell(row_index, 0).text_frame.paragraphs[0].font.bold = True
        table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        for colindex in range(0, cols):
            table.cell(row_index, colindex).fill.solid()
            table.cell(row_index, colindex).fill.fore_color.rgb = RGBColor(0xB0, 0xE0, 0xE6)
            table.cell(row_index, colindex).text_frame.paragraphs[0].font.italic = False
        table.cell(row_index, 1).text = "{:,}".format(int(ceil(usable_capacity['CPU'])))
        table.cell(row_index, 3).text = "{:,}".format(round(usable_capacity['RAM'], 1))
        binary_unit = "{:,}".format(round(usable_capacity['HDD'] * HyperConstants.TB_TO_TIB_CONVERSION, 1))
        table.cell(row_index, 5).text = \
            "{:,}".format(round(usable_capacity['HDD'], 1)) + ' TB (' + binary_unit + ' TiB)'

        # merge cells horizontally
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=1, end_col_idx=2)
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=3, end_col_idx=4)
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=5, end_col_idx=6)
        table.cell(row_index, 7).text = _('Total resources usable after accounting for controller VM overhead, ' \
                                        'and performance headroom')
        table.cell(row_index, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        row_index += 1
        table.cell(row_index, 0).text = _("Best Practice Threshold")
        table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        table.cell(row_index, 1).text = str(node_reserve['CPU']) + "%"
        reserve_diff = reserve_capacity['CPU'] - usable_capacity['CPU']
        table.cell(row_index, 2).text = "{:,}".format(round(reserve_diff, 1))
        table.cell(row_index, 3).text = str(node_reserve['RAM']) + "%"
        reserve_diff = reserve_capacity['RAM'] - usable_capacity['RAM']
        table.cell(row_index, 4).text = "{:,}".format(round(reserve_diff, 1))
        table.cell(row_index, 5).text = str(node_reserve['HDD']) + "%"
        reserve_diff = reserve_capacity['HDD'] - usable_capacity['HDD']
        binary_unit = "{:,}".format(round(reserve_diff * HyperConstants.TB_TO_TIB_CONVERSION, 1))
        table.cell(row_index, 6).text = "{:,}".format(round(reserve_diff, 1)) + ' TB (' + binary_unit + ' TiB)'

        tooltip = _("CPU & Memory: This ensures predictable behavior for user VMs")

        if self.scenario_data['settings_json'][0]['threshold'] == FULL_CAPACITY:
            tooltip += _(" and does not include reservations for the hypervisor.")

        tooltip += "\n" + _("Storage: Allows for better resilience and performance.")
        if self.scenario_data['settings_json'][0]['threshold'] == FULL_CAPACITY:
            tooltip += "\n" + _("Hypervisor overheads must be accounted for separately.")

        table.cell(row_index, 7).text = _(tooltip)
        table.cell(row_index, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        table.cell(row_index, 7).text_frame.paragraphs[0].font.italic = False

        row_index += 1
        if self.hercules_conf:
            table.cell(row_index, 0).text = _("Storage Efficiency Savings Multiplier") + "*"
        else:
            table.cell(row_index, 0).text = _("Storage Efficiency Savings Multiplier")
        table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        effective_storage_savings_mul = round((100.0 + op_ratio['HDD']) / 100.0, 2)
        table.cell(row_index, 5).text = str(effective_storage_savings_mul)
        total_effective_resource = round((reserve_capacity['HDD'] / scaling_factor), 1)
        binary_unit = "{:,}".format(round(total_effective_resource * HyperConstants.TB_TO_TIB_CONVERSION, 1))
        table.cell(row_index, 6).text = "{:,}".format(total_effective_resource) + ' TB (' + binary_unit + ' TiB)'
        tooltip = _("This multiplier shows how much storage has increased " \
                  "due to Dedupe and Compression Savings")
        if self.hercules_conf:
            tooltip += "\n" + _("*Savings has been increased due to Hardware Acceleration")
        table.cell(row_index, 7).text = tooltip
        table.cell(row_index, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        table.cell(row_index, 7).text_frame.paragraphs[0].font.italic = False

        row_index += 1
        table.cell(row_index, 0).text = _('Best Practice Resources')
        table.cell(row_index, 0).text_frame.paragraphs[0].font.bold = True
        table.cell(row_index, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        for colindex in range(0, cols):
            table.cell(row_index, colindex).fill.solid()
            table.cell(row_index, colindex).fill.fore_color.rgb = RGBColor(0xB0, 0xE0, 0xE6)
            table.cell(row_index, colindex).text_frame.paragraphs[0].font.italic = False
        table.cell(row_index, 1).text = "{:,}".format(int(ceil(reserve_capacity['CPU'])))
        table.cell(row_index, 3).text = "{:,}".format(round(reserve_capacity['RAM'], 1))
        binary_unit = "{:,}".format(round(total_effective_resource * HyperConstants.TB_TO_TIB_CONVERSION, 1))
        table.cell(row_index, 5).text = "{:,}".format(total_effective_resource) + ' TB (' + binary_unit + ' TiB)'
        table.cell(row_index, 7).text = _('Includes Best Practice Not-to-Exceed Threshold')
        table.cell(row_index, 7).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        # merge cells horizontally
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=1, end_col_idx=2)
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=3, end_col_idx=4)
        self.merge_cells_horizontally(table=table, row_idx=row_index, start_col_idx=5, end_col_idx=6)

        return usable_capacity, usable_capacity

    def create_epiccluster_page(self, cluster_data_index, cluster_no, dc_count_map):

        if len(dc_count_map) == 2:
            dc1_end = dc_count_map['DC1'] + cluster_no - 1
            if dc_count_map['DC1'] == 1:
                cluster_name = "DC1 - Cluster " + str(cluster_no)
            else:
                cluster_name = "DC1 - Cluster [" + str(cluster_no) + " - " + str(dc1_end) + "]"

            cluster_data = cluster_data_index[0]
            util_list = cluster_data['Utilization']
            wl_list = cluster_data['wl_list']
            rep_factor = 3
            if 'replication_factor' in cluster_data['settings']:
                rep_factor = cluster_data['settings']['replication_factor']
            elif 'Replication_Factor' in cluster_data['settings']:
                rep_factor = cluster_data['settings']['Replication_Factor']
            ft = cluster_data['settings']['fault_tolerance']

            total_cluster_count = dc_count_map['DC1']
            self.create_per_cluster_sys_util_page(cluster_name, util_list, rep_factor, wl_list, ft, total_cluster_count)
            # This function is for workload calculation
            self.create_wl_calc_page(cluster_name, wl_list, total_cluster_count)
            self.get_node_calc_details(cluster_name, cluster_data, total_cluster_count)

            # self.create_UCSLayout_slide(cluster_name, cluster_data, self.fi_option)

            cluster_data_dc2 = cluster_data_index[dc_count_map['DC1']]
            util_list = cluster_data_dc2['Utilization']
            wl_list = cluster_data_dc2['wl_list']
            dc2_begin = dc_count_map['DC1'] + cluster_no
            dc2_end = dc2_begin + dc_count_map['DC2'] - 1

            if dc_count_map['DC2'] == 1:
                cluster_name = "DC2 - Cluster " + str(dc2_begin)
            else:
                cluster_name = "DC2 - Cluster [" + str(dc2_begin) + " - " + str(dc2_end) + "]"

            total_cluster_count = dc_count_map['DC2']
            self.create_per_cluster_sys_util_page(cluster_name, util_list, rep_factor, wl_list, ft, total_cluster_count)
            self.create_wl_calc_page(cluster_name, wl_list, total_cluster_count)
            self.get_node_calc_details(cluster_name, cluster_data_dc2, total_cluster_count)

            # self.create_UCSLayout_slide(cluster_name, cluster_data_dc2, self.fi_option)

        elif len(dc_count_map) == 1 and 'DC1' in dc_count_map:
            dc1_end = dc_count_map['DC1'] + cluster_no - 1
            if dc_count_map['DC1'] == 1:
                cluster_name = "DC1 - Cluster " + str(cluster_no)
            else:
                cluster_name = "DC1 - Cluster [" + str(cluster_no) + " - " + str(dc1_end) + "]"

            cluster_data = cluster_data_index[0]
            util_list = cluster_data['Utilization']
            wl_list = cluster_data['wl_list']
            rep_factor = 3
            if 'replication_factor' in cluster_data['settings']:
                rep_factor = cluster_data['settings']['replication_factor']
            elif 'Replication_Factor' in cluster_data['settings']:
                rep_factor = cluster_data['settings']['Replication_Factor']
            ft = cluster_data['settings']['fault_tolerance']

            total_cluster_count = dc_count_map['DC1']
            self.create_per_cluster_sys_util_page(cluster_name, util_list, rep_factor, wl_list, ft, total_cluster_count)
            self.create_wl_calc_page(cluster_name, wl_list, total_cluster_count)
            self.get_node_calc_details(cluster_name, cluster_data, total_cluster_count)

            # self.create_UCSLayout_slide(cluster_name, cluster_data, self.fi_option)

        elif len(dc_count_map) == 1 and 'DC2' in dc_count_map:
            dc2_end = dc_count_map['DC2'] + cluster_no - 1
            if dc_count_map['DC2'] == 1:
                cluster_name = "DC2 - Cluster " + str(cluster_no)
            else:
                cluster_name = "DC2 - Cluster [" + str(cluster_no) + " - " + str(dc2_end) + "]"

            cluster_data = cluster_data_index[0]
            util_list = cluster_data['Utilization']
            wl_list = cluster_data['wl_list']
            rep_factor = 3
            if 'replication_factor' in cluster_data['settings']:
                rep_factor = cluster_data['settings']['replication_factor']
            elif 'Replication_Factor' in cluster_data['settings']:
                rep_factor = cluster_data['settings']['Replication_Factor']
            ft = cluster_data['settings']['fault_tolerance']

            total_cluster_count = dc_count_map['DC2']
            self.create_per_cluster_sys_util_page(cluster_name, util_list, rep_factor, wl_list, ft, total_cluster_count)
            self.create_wl_calc_page(cluster_name, wl_list, total_cluster_count)
            self.get_node_calc_details(cluster_name, cluster_data, total_cluster_count)

            # self.create_UCSLayout_slide(cluster_name, cluster_data, self.fi_option)

    def create_UCSLayout_slide(self, cluster_name, cluster_data, fi_option):

        if "UCS Layout" not in self.slides:
            return

        hyper_node = cluster_data['node_info'][0]['model_details']
        hyper_count = cluster_data['node_info'][0]['num_nodes']

        comp_node = None
        comp_count = 0
        if self.fi_count_cluster > 0:
            fi_opt = fi_option
        else:
            fi_opt = None

        if len(cluster_data['node_info']) > 1:
            comp_node = cluster_data['node_info'][1]['model_details']
            comp_count = cluster_data['node_info'][1]['num_nodes']

        images = list()
        try:
            images = fetch_layout(hyper_node, comp_node, hyper_count, comp_count, fi_opt)

            if images:
                prs = self.prs
                slide = prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_BLANK_LAYOUT_INDEX])
                slide.shapes.title.text = cluster_name + " " + _("UCS Layout")
                shapes = slide.shapes

                urllib.urlretrieve(images[0], self.util_image_name)
                if len(images) == 1:
                    slide.shapes.add_picture(self.util_image_name, Inches(1.5), Inches(0.5), height=Inches(4.7),
                                             width=Inches(7.5))
                else:
                    slide.shapes.add_picture(self.util_image_name, Inches(0.5), Inches(0.5),
                                             width=Inches(2.5), height=Inches(4.7))

                # Ignoring UCS wiring image if count excceds 25 as its too cluttered
                total_nodes = hyper_count + comp_count
                if len(images) > 1 and total_nodes <= 25:
                    urllib.urlretrieve(images[1], self.util_image_name)
                    slide.shapes.add_picture(self.util_image_name, Inches(3.0), Inches(0.5),
                                             width=Inches(7.0), height=Inches(4.7))

                self.fi_count_cluster -= 1

        except Exception as e:
            pass

    def create_per_cluster_pages(self, wl_result_data):

        cluster_list = wl_result_data['clusters']

        cluster_no = 0
        self.fi_count_cluster = self.fi_count

        for cluster_data_index in cluster_list:
            for cluster_data in cluster_data_index:

                cluster_index = cluster_data_index.index(cluster_data)

                cluster_no += 1
                cluster_name = self.fixed_cluster_name if self.fixed_cluster_name else "Cluster " + str(cluster_no)

                stretchcluster = any(wl_list['cluster_type'] == 'stretch' for wl_list in cluster_data['wl_list'])
                if stretchcluster:
                    cluster_name = "Stretch Cluster " + str(cluster_no)

                epiccluster = any(wl_list_tmp['wl_type'] == 'EPIC' for wl_list_tmp in cluster_data['wl_list'])
                dcname = list()
                dc_count_map = defaultdict()
                if epiccluster:
                    for cluster_tmp_data in cluster_data_index:
                        epicdc = [tmp_wl_list['dc_name'] for tmp_wl_list in cluster_tmp_data['wl_list']]
                        dcname.extend(epicdc)

                    dc_count_map = dict(Counter(dcname))

                # Add Per Cluster System Utilization
                # cluster_name = cluster_data['node_info'][0]['display_name']
                util_list = cluster_data['Utilization']
                wl_list = cluster_data['wl_list']
                rep_factor = 3
                if 'replication_factor' in cluster_data['settings']:
                    rep_factor = cluster_data['settings']['replication_factor']
                elif 'Replication_Factor' in cluster_data['settings']:
                    rep_factor = cluster_data['settings']['Replication_Factor']
                ft = cluster_data['settings']['fault_tolerance']

                self.num_nodes = cluster_data['node_info'][0]['num_nodes']

                # is cluster is Hercules node ?
                self.hercules_conf = any(
                    True if node_data['hercules_conf'] else False for node_data in cluster_data['node_info'])

                if epiccluster:
                    self.create_epiccluster_page(cluster_data_index, cluster_no, dc_count_map)
                else:
                    self.create_per_cluster_sys_util_page(cluster_name, util_list, rep_factor, wl_list, ft)

                    # This function is for workload calculation
                    self.create_wl_calc_page(cluster_name, wl_list, cluster_index)

                    # Below functions will get calculation and create excel sheet
                    self.get_node_calc_details(cluster_name, cluster_data)

                    # UCS Layout code is commented as product management decision of not using at the moment
                    # self.create_UCSLayout_slide(cluster_name, cluster_data, self.fi_option)

                # If Cluster_type is STRETCH, no need to provide cluster information for cluster pair
                if stretchcluster:
                    cluster_no += 1
                    break

                if epiccluster:
                    cluster_no += len(cluster_data_index) - 1
                    break

    def create_glossary_page(self):

        if self.language == JAPANESE:
            prs = self.prs
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_JA_GLOSSARY_INDEX])
            return

        glossarydict = OrderedDict()
        glossarydict['header'] = [_('Term'), _('Definition')]
        row_index = 0

        glossary_path = os.path.join(BASE_DIR, "sizer/sizerglossary.txt")
        with open(glossary_path, 'r', encoding='utf-8') as f:
            for line in f:
                if ':' in line:
                    (key, val) = line.split(':')
                    glossary_per_row = list()
                    glossary_per_row.append(key)
                    output = val.replace('For example', '\n For example')
                    glossary_per_row.append(output)

                    row_name = "row" + str(row_index)
                    glossarydict[row_name] = glossary_per_row
                    row_index += 1

        # Coordinates for table
        x = 0.32  # LEFT POSITION
        y = 0.6  # TOP POSITION
        wid = 10  # table Width
        hei = 0.9  # height
        slide_title = _("Glossary")
        # col_width = [1.3, 0.8, 5.3, 1.9]
        col_width = [1.5, 7.5]

        rows_per_slide = 6
        self.render_table_across_pages(slide_title, x, y, wid, hei, col_width, glossarydict, rows_per_slide)

    def close(self):

        prs = self.prs

        # LAST SLIDE
        if self.language == JAPANESE:
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_LAST_LAYOUT_INDEX_JA])
        else:
            prs.slides.add_slide(prs.slide_layouts[PPTReport.SLIDEMASTER_LAST_LAYOUT_INDEX])

        scenario_name = self.scenario_data['name']
        currentdatetime = datetime.datetime.now().strftime('%Y-%m-%d_%H-%M')
        report_name = scenario_name + "_" + str(currentdatetime) + ".pptx"

        report_name = os.path.join(BASE_DIR, report_name)
        prs.save(report_name)

        # Delete the ucsimage created.
        try:
            os.system("rm " + str(self.util_image_name))
        except Exception as e:
            pass

        return report_name


def generate_report(scenario_data, request_data):
    global _

    sdata = json.loads(scenario_data)

    if 'dr_enabled' not in sdata['settings_json'][0]:
        sdata['settings_json'][0]['dr_enabled'] = False

    scenario_data = json.dumps(sdata)
    req_data = json.loads(request_data)

    report = PPTReport(scenario_data, req_data)

    report.create_header_page()

    report.create_disclaimer_page()

    if "Agenda" in req_data['slides']:
        report.create_agenda_page()

    report.create_overview_page()

    # This is Optimal sizing report to include only Lowest_Cost, All-Flash and All-NVMe
    sizing_results = len(sdata['workload_result'])
    if sizing_results:

        issameresult = False
        if "Sizing Report" in req_data['slides']:
            report.create_page_break(_("Sizing Report"))
            # sizing_type = sdata['workload_json'].get('sizing_type', 'optimal')

            # if sizing_type == 'optimal':
            #     report.create_user_in_page()
            # else:
            #     report.create_fixed_user_in_page()
            report.create_user_in_page()

        if len(sdata['workload_result']) > 2 and "Lowest_Cost/All-Flash/All NVMe Comparison" in req_data['slides']:
            report.create_comparison_page()

            wl_result_cmp = copy.deepcopy(sdata['workload_result'])
            wl_result_cmp[0]['result_name'] = ""
            wl_result_cmp[1]['result_name'] = ""

            if wl_result_cmp[0] == wl_result_cmp[1]:
                issameresult = True

        # Sizing Report - For Optimal
        wl_result_list = sdata['workload_result']
        for wl_result_data in wl_result_list:
            title = wl_result_data['result_name'] + " Sizing Report"
            lang_title = _(wl_result_data['result_name'] + " " + "Sizing Report")

            if title in req_data['slides']:
                report.create_page_break(lang_title)
                report.create_wl_cluster_map_page(wl_result_data)
                report.create_sizing_results_page(wl_result_data)
                report.create_per_cluster_pages(wl_result_data)

            # Commented to get all Sizing results even if lowestcost and All-Flash have same result
            # if issameresult:
            #     break

    #Sizing Report - For Fixed
    if "Fixed Config Sizing Report" in req_data['slides'] and len(sdata['fixed_workload_result']):

        for fixed_setting_json in sdata['fixed_setting_json']:
            calc_result = dict()
            # calc_result['node_properties'] = sdata['settings_json'][0]['node_properties']
            # calc_result['scenario_settings'] = sdata['settings_json'][0]['cluster_properties']
            # calc_result['scenario_settings']['threshold'] = sdata['settings_json'][0]['threshold']
            # calc_result['scenario_settings']['hypervisor'] = sdata['settings_json'][0]['hypervisor']

            # calc_result['results'] = sdata['workload_result'][0]['sizing_calculator']

            calc_result['node_properties'] = fixed_setting_json['node_properties']
            calc_result['scenario_settings'] = fixed_setting_json['cluster_properties']
            calc_result['scenario_settings']['threshold'] = fixed_setting_json['threshold']
            calc_result['scenario_settings']['hypervisor'] = fixed_setting_json['hypervisor']
            cluster_name = fixed_setting_json['cluster_name']

            fixed_wl_result = None
            if (not cluster_name in sdata['fixed_workload_result']) and len(sdata['fixed_setting_json']) == 1:
                if 'No_cluster' not in sdata['fixed_workload_result']:
                    continue

                fixed_wl_result = sdata['fixed_workload_result']['No_cluster']
                report.create_page_break(_("Fixed Config Sizing Report"))

            else:
                fixed_wl_result = sdata['fixed_workload_result'][cluster_name]
                lang_title = _("Fixed Config Sizing Report") + "-" + cluster_name
                report.create_page_break(lang_title)

            calc_result['results'] = fixed_wl_result[0]['sizing_calculator']

            calc_scenario_data = json.dumps(calc_result)
            SizingCalculatorReport(report, calc_scenario_data)

            if cluster_name in sdata['fixed_workload_result']:
                for wl_result_data in fixed_wl_result:

                    if 'fixed_error_data' in wl_result_data or 'clusters' not in wl_result_data:
                        break

                    report.fixed_cluster_name = cluster_name
                    report.create_wl_cluster_map_page(wl_result_data)
                    report.create_sizing_results_page(wl_result_data)
                    report.create_per_cluster_pages(wl_result_data)


    if "Workload Summary" in req_data['slides']:
        wls = sdata['workload_json']['wl_list'] if 'wl_list' in sdata['workload_json'] else list()
        if wls and sizing_results:
            report.create_page_break(_("Backup"))
            report.create_wl_summary_page()

    if "Glossary" in req_data['slides']:
        report.create_glossary_page()

    report_path = report.close()
    return report_path

